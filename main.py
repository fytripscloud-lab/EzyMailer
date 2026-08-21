import sys
import html
import csv
import hashlib
import json
import os
import re
import platform
import random
import time
import secrets
import subprocess
import shutil
import ssl
import sqlite3
import socket
import string
import threading
import uuid
import certifi
import xml.etree.ElementTree
import xml.parsers.expat
from datetime import datetime
from io import BytesIO
import urllib.error
import urllib.request
from urllib.parse import quote_plus
from math import ceil, sqrt
from dataclasses import dataclass, field
from pathlib import Path
import tempfile
from typing import Callable

from backend.local_api import (
    API_BASE_URL,
    ensure_api_server,
    get_content as api_get_content,
    get_customer_variables as api_get_customer_variables,
    get_settings as api_get_settings,
    get_tags as api_get_tags,
    login as api_login,
    record_activity,
    delete_content as api_delete_content,
    delete_customer_variables as api_delete_customer_variables,
    delete_tags as api_delete_tags,
    save_content as api_save_content,
    save_customer_variables as api_save_customer_variables,
    save_tags as api_save_tags,
    update_content as api_update_content,
    upsert_setting,
)
from PySide6.QtCore import (
    QDateTime,
    QEasingCurve,
    QObject,
    Property,
    QPropertyAnimation,
    QThread,
    QTimer,
    Qt,
    QEvent,
    QPoint,
    QSize,
    Signal,
)
from PySide6.QtGui import QColor, QFont, QPainter, QPen, QKeySequence, QGuiApplication, QClipboard
from PySide6.QtWidgets import (
    QApplication,
    QAbstractItemView,
    QButtonGroup,
    QCheckBox,
    QDialog,
    QDialogButtonBox,
    QComboBox,
    QFrame,
    QFormLayout,
    QGridLayout,
    QHBoxLayout,
    QLabel,
    QLineEdit,
    QListWidget,
    QListWidgetItem,
    QGraphicsOpacityEffect,
    QMainWindow,
    QMessageBox,
    QPushButton,
    QProgressBar,
    QStyle,
    QSpinBox,
    QStackedWidget,
    QTabWidget,
    QTableWidget,
    QTableWidgetItem,
    QFileDialog,
    QTextEdit,
    QDoubleSpinBox,
    QRadioButton,
    QScrollArea,
    QSizePolicy,
    QSplitter,
    QVBoxLayout,
    QTextBrowser,
    QWidget,
    QHeaderView,
)
APP_TITLE = "EzyMailer"
DEFAULT_USERNAME = "admin"
DEFAULT_PASSWORD = "admin"
IS_MAC = sys.platform == "darwin"
IS_WINDOWS = sys.platform.startswith("win")
MAX_BODY_TABS = 50
MAX_ATTACHMENT_TABS = 50
MAX_SUBJECTS = 100
LOCAL_CACHE_DIR = (
    Path.home() / "Library" / "Application Support" / APP_TITLE
    if IS_MAC
    else Path.home() / ".ezymailer"
)
LEGACY_LOCAL_CACHE_DIR = Path.home() / ".ezymailer"
LEGACY_LOCAL_CACHE_DB = LEGACY_LOCAL_CACHE_DIR / "local_drafts.sqlite3"
LOCAL_CACHE_DB = LOCAL_CACHE_DIR / "local_drafts.sqlite3"
LOCAL_ATTACHMENT_STATE_KEY = "attachment_content_state"
LOCAL_SUBJECT_STATE_KEY = "subject_content_state"
LOCAL_BODY_STATE_KEY = "body_content_state"
LOCAL_PENDING_EMAILS_STATE_KEY = "pending_emails_state"
LOCAL_TAG_STATE_KEY = "tag_state"
LOCAL_CUSTOMER_VARIABLES_TABLE = "customer_variables"
LOCAL_BROWSER_STATE_KEY = "browser_controls_state"
LOCAL_SETTINGS_STATE_KEY = "sending_settings_state"
ROLE_LOCAL_ONLY = Qt.UserRole + 10
ROLE_LOCAL_DRAFT_ID = Qt.UserRole + 11
BUILTIN_BROWSER_DIR_NAME = "playwright-browsers"
DEFAULT_BROWSER_DOWNLOAD_HOST = "https://cdn.playwright.dev"
DEPENDENCY_RELEASE_TAG = "dependencies-v4"
DEPENDENCY_RELEASE_BASE = (
    "https://github.com/fytripscloud-lab/EzyMailer/releases/download/"
    f"{DEPENDENCY_RELEASE_TAG}"
)


def _runtime_root() -> Path:
    if getattr(sys, "frozen", False):
        return Path(sys.executable).resolve().parent
    return Path(__file__).resolve().parent


def _external_dependency_dir() -> Path:
    if getattr(sys, "frozen", False) and IS_MAC:
        root = _runtime_root().parents[2]
    elif getattr(sys, "frozen", False):
        root = _runtime_root()
    else:
        root = _runtime_root()
    return root / ".ezymailer" / "dependencies" / DEPENDENCY_RELEASE_TAG


def _external_dependency_asset() -> str:
    if IS_MAC and platform.machine().lower() in {"arm64", "aarch64"}:
        return "ezymailer-dependencies-macos-arm64.zip"
    if IS_WINDOWS:
        return "ezymailer-dependencies-windows-x64.zip"
    raise RuntimeError(f"No external dependency archive is available for {platform.system()} {platform.machine()}")


def _restore_runtime_permissions(root: Path) -> None:
    """Restore executable bits lost when the dependency ZIP is extracted."""
    if not root.exists():
        return
    executable_paths = [
        root / "playwright" / "driver" / "node",
        root / "playwright" / "driver" / "node.exe",
        root / "playwright" / "driver.sh",
    ]
    for path in executable_paths:
        if not path.is_file() or IS_WINDOWS:
            continue
        try:
            path.chmod(path.stat().st_mode | 0o111)
        except OSError:
            pass


_external_dll_directory_handles: list[object] = []
_external_dll_directories: set[str] = set()


def _configure_external_dll_search(root: Path) -> None:
    """Make extracted native dependency DLLs discoverable on Windows."""
    if not IS_WINDOWS:
        return
    for directory in (root, root / "PIL"):
        if not directory.is_dir():
            continue
        try:
            directory_text = str(directory.resolve())
            if directory_text in _external_dll_directories:
                continue
            handle = os.add_dll_directory(directory_text)
            _external_dll_directory_handles.append(handle)
            _external_dll_directories.add(directory_text)
        except (AttributeError, OSError):
            pass


def ensure_external_dependencies(progress: Callable[[str, str, int, int], None] | None = None) -> Path:
    """Download and extract the versioned GitHub dependency pack once."""
    target = _external_dependency_dir()
    marker = target / ".ready"
    target.mkdir(parents=True, exist_ok=True)
    if marker.exists():
        _restore_runtime_permissions(target)
        _configure_external_dll_search(target)
        if str(target) not in sys.path:
            sys.path.append(str(target))
        return target

    asset_name = _external_dependency_asset()
    archive_path = target.with_suffix(".zip.part")
    url = f"{DEPENDENCY_RELEASE_BASE}/{asset_name}"
    if progress:
        progress("Downloading app dependencies", "Connecting to GitHub CDN • ETA calculating", 1, 100)
    request = urllib.request.Request(url, headers={"User-Agent": "EazyMailer"})
    try:
        tls_context = ssl.create_default_context(cafile=certifi.where())
        with urllib.request.urlopen(request, timeout=30, context=tls_context) as response, archive_path.open("wb") as output:
            total = int(response.headers.get("Content-Length") or 0)
            downloaded = 0
            started_at = time.monotonic()
            while True:
                chunk = response.read(1024 * 1024)
                if not chunk:
                    break
                output.write(chunk)
                downloaded += len(chunk)
                elapsed = max(time.monotonic() - started_at, 0.1)
                speed = downloaded / elapsed
                eta = (total - downloaded) / speed if total and speed > 0 else None
                if eta is None:
                    eta_text = "ETA calculating"
                else:
                    eta_text = f"ETA {int(eta // 60)}m {int(eta % 60):02d}s"
                percent = int(downloaded * 85 / total) if total else 5
                if progress:
                    progress(
                        "Downloading app dependencies",
                        f"{downloaded / 1048576:.0f} MB downloaded • {eta_text}",
                        max(1, min(85, percent)),
                        100,
                    )
    except Exception:
        archive_path.unlink(missing_ok=True)
        raise

    if progress:
        progress("Extracting app dependencies", "Automatically extracting the downloaded runtime", 90, 100)
    import zipfile

    extraction_dir = target.with_name(target.name + ".extracting")
    shutil.rmtree(extraction_dir, ignore_errors=True)
    extraction_dir.mkdir(parents=True, exist_ok=True)
    with zipfile.ZipFile(archive_path) as archive:
        members = archive.infolist()
        total_members = max(1, len(members))
        for index, member in enumerate(members, start=1):
            archive.extract(member, extraction_dir)
            if progress and (index == 1 or index == total_members or index % 25 == 0):
                progress(
                    "Extracting app dependencies",
                    f"Extracted {index:,} of {total_members:,} files",
                    min(99, 90 + int(index * 9 / total_members)),
                    100,
                )
    if progress:
        progress("Finalizing app dependencies", "Completing the local runtime setup", 99, 100)
    shutil.rmtree(target, ignore_errors=True)
    extraction_dir.rename(target)
    _restore_runtime_permissions(target)
    _configure_external_dll_search(target)
    archive_path.unlink(missing_ok=True)
    packaged_browser = target / BUILTIN_BROWSER_DIR_NAME
    browser_cache = _browser_cache_dir()
    if packaged_browser.exists():
        if browser_cache.exists():
            shutil.rmtree(packaged_browser, ignore_errors=True)
        else:
            browser_cache.parent.mkdir(parents=True, exist_ok=True)
            packaged_browser.rename(browser_cache)
    marker = target / ".ready"
    marker.write_text(DEPENDENCY_RELEASE_TAG, encoding="utf-8")
    if str(target) not in sys.path:
        sys.path.append(str(target))
    if progress:
        progress("Dependencies ready", "The downloaded runtime is ready to use", 100, 100)
    return target


def _find_first_existing(paths: list[Path]) -> Path | None:
    for path in paths:
        if path.exists():
            return path
    return None


def _find_executable(root: Path, names: list[str]) -> Path | None:
    for name in names:
        direct = root / name
        if direct.exists():
            return direct
    if not root.exists():
        return None

    # Playwright's browser cache has a stable two-level layout. Avoid a full
    # recursive scan here because a completed Chromium cache is hundreds of
    # megabytes and can otherwise delay the post-login loader for minutes.
    version_roots = [root] if root.name.startswith("chromium-") else list(root.glob("chromium-*"))
    for version_root in version_roots:
        platform_roots = list(version_root.glob("chrome-*"))
        if not platform_roots:
            platform_roots = list(version_root.glob("chromium-*"))
        for platform_root in platform_roots:
            mac_app = platform_root / "Google Chrome for Testing.app" / "Contents" / "MacOS" / "Google Chrome for Testing"
            if mac_app.exists():
                return mac_app
            for name in names:
                candidate = platform_root / name
                if candidate.exists():
                    return candidate
    return None


def _browser_cache_dir() -> Path:
    override = os.getenv("EZYM_MAILER_BROWSER_CACHE_DIR", "").strip()
    if override:
        return Path(override).expanduser()

    # Use one deterministic cache beside the portable build. Do not inspect
    # legacy or system caches; this makes first-run behavior predictable.
    if getattr(sys, "frozen", False):
        if IS_MAC:
            adjacent_root = _runtime_root().parents[2]
        else:
            adjacent_root = _runtime_root()
        return adjacent_root / ".ezymailer" / BUILTIN_BROWSER_DIR_NAME
    return _runtime_root() / ".ezymailer" / BUILTIN_BROWSER_DIR_NAME


def _cached_browser_binary(root: Path) -> Path | None:
    """Resolve only the expected Playwright Chromium layout."""
    for version_root in sorted(root.glob("chromium-*")):
        direct_paths = [
            version_root / "chrome-mac-arm64" / "Google Chrome for Testing.app" / "Contents" / "MacOS" / "Google Chrome for Testing",
            version_root / "chrome-win64" / "chrome.exe",
            version_root / "chrome-win" / "chrome.exe",
            version_root / "chrome-linux64" / "chrome",
            version_root / "chrome-linux" / "chrome",
        ]
        for candidate in direct_paths:
            if candidate.exists():
                return candidate
    return None


def _installed_browser_binary() -> Path | None:
    """Prefer an installed Edge or Chrome before downloading Chromium."""
    candidates: list[Path] = []
    if IS_MAC:
        for app_root in (Path("/Applications"), Path.home() / "Applications"):
            candidates.extend(
                [
                    app_root / "Microsoft Edge.app/Contents/MacOS/Microsoft Edge",
                    app_root / "Google Chrome.app/Contents/MacOS/Google Chrome",
                    app_root / "Chromium.app/Contents/MacOS/Chromium",
                ]
            )
    elif IS_WINDOWS:
        local_app_data = Path(os.getenv("LOCALAPPDATA", ""))
        program_files = Path(os.getenv("ProgramFiles", "C:/Program Files"))
        program_files_x86 = Path(os.getenv("ProgramFiles(x86)", "C:/Program Files (x86)"))
        candidates.extend(
            [
                program_files / "Microsoft/Edge/Application/msedge.exe",
                program_files_x86 / "Microsoft/Edge/Application/msedge.exe",
                local_app_data / "Microsoft/Edge/Application/msedge.exe",
                program_files / "Google/Chrome/Application/chrome.exe",
                program_files_x86 / "Google/Chrome/Application/chrome.exe",
                local_app_data / "Google/Chrome/Application/chrome.exe",
            ]
        )
    else:
        candidates.extend(
            [
                Path("/usr/bin/microsoft-edge"),
                Path("/usr/bin/microsoft-edge-stable"),
                Path("/usr/bin/google-chrome"),
                Path("/usr/bin/google-chrome-stable"),
                Path("/usr/bin/chromium"),
            ]
        )
    return _find_first_existing(candidates)


def _browser_product_name(binary: Path) -> str:
    """Return the user-facing browser name for a selected executable."""
    normalized = str(binary).lower().replace("\\", "/")
    if "microsoft edge" in normalized or normalized.endswith("/msedge.exe") or "microsoft-edge" in normalized:
        return "Microsoft Edge"
    if "google chrome" in normalized or normalized.endswith("/chrome.exe") or "google-chrome" in normalized:
        return "Google Chrome"
    return "Chromium"


def _browser_private_flag(browser_name: str) -> str:
    """Use the private-window switch supported by the selected browser."""
    return "--inprivate" if browser_name == "Microsoft Edge" else "--incognito"


class BrowserBootstrapWorker(QObject):
    progress = Signal(str, str, int, int)
    finished = Signal(bool, str)

    def __init__(self, browser_cache_dir: Path, parent=None):
        super().__init__(parent)
        self._browser_cache_dir = browser_cache_dir
        self.latest_progress: tuple[str, str, int, int] | None = None
        self.result: tuple[bool, str] | None = None
        self._state_lock = threading.Lock()

    def _report(self, title: str, subtitle: str, value: int, total: int) -> None:
        with self._state_lock:
            self.latest_progress = (title, subtitle, value, total)

    def _finish(self, success: bool, message: str) -> None:
        with self._state_lock:
            self.result = (success, message)

    def run(self) -> None:
        try:
            self._report(
                "The app is building",
                "Checking the local browser cache.",
                10,
                100,
            )
            self._browser_cache_dir.mkdir(parents=True, exist_ok=True)
            ensure_external_dependencies(self._report)
            installed_browser = _installed_browser_binary()
            existing_browser = _cached_browser_binary(self._browser_cache_dir)
            selected_browser = installed_browser or existing_browser
            if selected_browser is not None:
                browser_name = "installed browser" if installed_browser and not existing_browser else "browser runtime"
                if existing_browser is not None:
                    self._report(
                    "Browser ready",
                    f"Using {_browser_product_name(selected_browser)} for browser sessions.",
                    100,
                    100,
                )
                else:
                    self._report(
                        "Browser ready",
                        f"Using the detected {browser_name}.",
                        100,
                        100,
                    )
                self._finish(True, "Browser runtime already configured.")
                return

            # Move visibly into the download phase before Playwright resolves
            # its driver, so a slow CDN/driver startup cannot look hung.
            self._report(
                "Downloading Chromium",
                "Connecting to the Chromium CDN • ETA calculating",
                15,
                100,
            )
            from playwright._impl._driver import compute_driver_executable, get_driver_env

            driver_executable, driver_cli = compute_driver_executable()
            self._report(
                "The app is building",
                "Additional files are downloading from the CDN.",
                2,
                100,
            )
            env = os.environ.copy()
            env.update(get_driver_env())
            env["PLAYWRIGHT_BROWSERS_PATH"] = str(self._browser_cache_dir)
            # Keep the first-run browser download on the official Playwright
            # Chromium CDN instead of relying on an installed system browser.
            env["PLAYWRIGHT_DOWNLOAD_HOST"] = os.getenv(
                "EZYM_MAILER_BROWSER_DOWNLOAD_HOST",
                DEFAULT_BROWSER_DOWNLOAD_HOST,
            )
            process = subprocess.Popen(
                [driver_executable, driver_cli, "install", "chromium"],
                env=env,
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL,
            )
            started_at = time.monotonic()
            expected_bytes = 360 * 1024 * 1024
            while process.poll() is None:
                downloaded_bytes = 0
                try:
                    for file_path in self._browser_cache_dir.rglob("*"):
                        if file_path.is_file():
                            downloaded_bytes += file_path.stat().st_size
                except OSError:
                    pass

                elapsed = max(time.monotonic() - started_at, 0.1)
                speed = downloaded_bytes / elapsed
                remaining = max(expected_bytes - downloaded_bytes, 0)
                eta = remaining / speed if speed > 0 else None
                if eta is None:
                    eta_text = "calculating ETA"
                else:
                    eta_text = f"ETA {int(eta // 60)}m {int(eta % 60):02d}s"
                # Progress reflects the measured bytes, not a synthetic 20%
                # offset. Keep a small visible value while the first chunk is
                # being established.
                progress = min(98, max(2, int((downloaded_bytes / expected_bytes) * 98)))
                downloaded_mb = downloaded_bytes / (1024 * 1024)
                self._report(
                    "Downloading Chromium",
                    f"{downloaded_mb:.0f} MB downloaded • {eta_text}",
                    progress,
                    100,
                )
                time.sleep(0.5)

            if process.returncode != 0:
                raise RuntimeError("Playwright could not download Chromium.")
            self._report(
                "The app is building",
                "Automatically configuring the downloaded files.",
                90,
                100,
            )
            self._report(
                "Browser ready",
                "The workspace is ready to use.",
                100,
                100,
            )
            self._finish(True, "Browser runtime configured.")
        except Exception as exc:
            self._finish(False, str(exc))


def _scaled_int(value: float, scale: float, minimum: int = 1) -> int:
    return max(minimum, int(round(value * scale)))


def _compute_layout_scale(screen) -> float:
    if screen is None:
        return 1.00

    geometry = screen.availableGeometry()
    width_boost = max(0.0, (geometry.width() - 1280.0) / 8000.0)
    height_boost = max(0.0, (geometry.height() - 768.0) / 8000.0)
    scale = 1.0 + min(0.06, width_boost + height_boost)
    return max(1.0, min(1.06, scale))


def _compute_text_scale(screen) -> float:
    if screen is None:
        return 0.92 if IS_WINDOWS else 1.28

    geometry = screen.availableGeometry()
    width_boost = max(0.0, (geometry.width() - 1280.0) / 2200.0)
    height_boost = max(0.0, (geometry.height() - 768.0) / 2600.0)
    # Keep Windows compact, but readable on standard laptop displays.
    base = 1.04 if IS_WINDOWS else 1.28
    ceiling = 1.12 if IS_WINDOWS else 1.40
    scale = base + min(0.12, (width_boost + height_boost) * 0.8)
    return max(base, min(ceiling, scale))


def _device_fingerprint() -> str:
    node = uuid.getnode()
    hostname = socket.gethostname()
    machine = platform.machine()
    system = platform.system()
    raw = f"{node:x}|{hostname}|{machine}|{system}"
    return hashlib.sha256(raw.encode("utf-8")).hexdigest()


def _device_name() -> str:
    return f"{platform.system()} {platform.release()} ({socket.gethostname()})"


def _is_subject_content_row(row: dict[str, object]) -> bool:
    if not isinstance(row, dict):
        return False
    if bool(row.get("local_only")) and str(row.get("content_type") or "") == "subject":
        return True
    content_type = str(row.get("content_type") or "")
    details = row.get("details_json")
    kind = ""
    if isinstance(details, str):
        try:
            parsed = json.loads(details)
            if isinstance(parsed, dict):
                kind = str(parsed.get("kind") or "")
        except Exception:
            kind = ""
    elif isinstance(details, dict):
        kind = str(details.get("kind") or "")
    return content_type in {"subject", "subject-draft"} or kind in {"subject", "subject-draft"}


def _subject_rows_from_content(rows: list[dict[str, object]]) -> list[dict[str, object]]:
    subject_rows: list[dict[str, object]] = []
    for row in rows:
        if not isinstance(row, dict):
            continue
        if not _is_subject_content_row(row):
            continue
        subject = str(row.get("subject") or "").strip()
        if not subject:
            continue
        subject_rows.append(row)
    return subject_rows


def _subject_count_text(count: int) -> str:
    return f"{count}/{MAX_SUBJECTS} subjects"


def _ensure_local_cache_db() -> None:
    LOCAL_CACHE_DIR.mkdir(parents=True, exist_ok=True)
    try:
        os.chmod(LOCAL_CACHE_DIR, 0o700)
    except Exception:
        pass
    if not LOCAL_CACHE_DB.exists() and LEGACY_LOCAL_CACHE_DB.exists():
        try:
            shutil.copy2(LEGACY_LOCAL_CACHE_DB, LOCAL_CACHE_DB)
        except Exception:
            pass
    connection = sqlite3.connect(LOCAL_CACHE_DB)
    try:
        connection.execute(
            """
            CREATE TABLE IF NOT EXISTS ui_state (
                state_key TEXT PRIMARY KEY,
                payload_json TEXT NOT NULL DEFAULT '{}',
                updated_at TEXT NOT NULL DEFAULT CURRENT_TIMESTAMP
            )
            """
        )
        connection.execute(
            """
            CREATE TABLE IF NOT EXISTS tag_state (
                state_key TEXT PRIMARY KEY,
                payload_json TEXT NOT NULL DEFAULT '{}',
                updated_at TEXT NOT NULL DEFAULT CURRENT_TIMESTAMP
            )
            """
        )
        connection.execute(
            """
            CREATE TABLE IF NOT EXISTS customer_variables (
                email TEXT NOT NULL,
                variables_json TEXT NOT NULL DEFAULT '{}',
                updated_at TEXT NOT NULL DEFAULT CURRENT_TIMESTAMP,
                PRIMARY KEY (email)
            )
            """
        )
        connection.execute(
            """
            CREATE TABLE IF NOT EXISTS browser_sessions (
                session_id TEXT PRIMARY KEY,
                username TEXT NOT NULL DEFAULT '',
                title TEXT NOT NULL,
                browser_name TEXT NOT NULL,
                browser_mode TEXT NOT NULL,
                status TEXT NOT NULL,
                browser_pid INTEGER NULL,
                launch_preset TEXT NOT NULL DEFAULT 'Default',
                details_json TEXT NOT NULL DEFAULT '{}',
                started_at TEXT NOT NULL DEFAULT CURRENT_TIMESTAMP,
                updated_at TEXT NOT NULL DEFAULT CURRENT_TIMESTAMP,
                closed_at TEXT NULL
            )
            """
        )
        connection.execute(
            """
            CREATE TABLE IF NOT EXISTS draft_content (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                content_type TEXT NOT NULL,
                title TEXT NOT NULL,
                subject TEXT NOT NULL DEFAULT '',
                body_text TEXT NOT NULL DEFAULT '',
                body_html TEXT NOT NULL DEFAULT '',
                details_json TEXT NOT NULL DEFAULT '{}',
                created_at TEXT NOT NULL DEFAULT CURRENT_TIMESTAMP,
                updated_at TEXT NOT NULL DEFAULT CURRENT_TIMESTAMP
            )
            """
        )
        connection.execute(
            """
            CREATE TABLE IF NOT EXISTS attachment_state (
                state_key TEXT PRIMARY KEY,
                payload_json TEXT NOT NULL DEFAULT '{}',
                updated_at TEXT NOT NULL DEFAULT CURRENT_TIMESTAMP
            )
            """
        )
        connection.commit()
    finally:
        connection.close()
    try:
        os.chmod(LOCAL_CACHE_DB, 0o600)
    except Exception:
        pass


def _save_local_draft(
    content_type: str,
    title: str,
    subject: str = "",
    body_text: str = "",
    body_html: str = "",
    details: dict[str, object] | None = None,
) -> int:
    _ensure_local_cache_db()
    connection = sqlite3.connect(LOCAL_CACHE_DB)
    try:
        cursor = connection.cursor()
        cursor.execute(
            """
            INSERT INTO draft_content (content_type, title, subject, body_text, body_html, details_json, created_at, updated_at)
            VALUES (?, ?, ?, ?, ?, ?, CURRENT_TIMESTAMP, CURRENT_TIMESTAMP)
            """,
            (
                content_type,
                title,
                subject,
                body_text,
                body_html,
                json.dumps(details or {}, ensure_ascii=False),
            ),
        )
        connection.commit()
        return int(cursor.lastrowid or 0)
    finally:
        connection.close()


def _update_local_draft(
    draft_id: int,
    content_type: str,
    title: str,
    subject: str = "",
    body_text: str = "",
    body_html: str = "",
    details: dict[str, object] | None = None,
) -> None:
    _ensure_local_cache_db()
    connection = sqlite3.connect(LOCAL_CACHE_DB)
    try:
        connection.execute(
            """
            UPDATE draft_content
            SET content_type = ?,
                title = ?,
                subject = ?,
                body_text = ?,
                body_html = ?,
                details_json = ?,
                updated_at = CURRENT_TIMESTAMP
            WHERE id = ?
            """,
            (
                content_type,
                title,
                subject,
                body_text,
                body_html,
                json.dumps(details or {}, ensure_ascii=False),
                draft_id,
            ),
        )
        connection.commit()
    finally:
        connection.close()


def _delete_local_draft(draft_id: int) -> None:
    _ensure_local_cache_db()
    connection = sqlite3.connect(LOCAL_CACHE_DB)
    try:
        connection.execute("DELETE FROM draft_content WHERE id = ?", (draft_id,))
        connection.commit()
    finally:
        connection.close()


def _delete_local_drafts(content_type: str | None = None) -> int:
    _ensure_local_cache_db()
    connection = sqlite3.connect(LOCAL_CACHE_DB)
    try:
        cursor = connection.cursor()
        if content_type:
            cursor.execute("DELETE FROM draft_content WHERE content_type = ?", (content_type,))
        else:
            cursor.execute("DELETE FROM draft_content")
        connection.commit()
        return int(cursor.rowcount or 0)
    finally:
        connection.close()


def _list_local_drafts(content_type: str | None = None) -> list[dict[str, object]]:
    _ensure_local_cache_db()
    connection = sqlite3.connect(LOCAL_CACHE_DB)
    connection.row_factory = sqlite3.Row
    try:
        cursor = connection.cursor()
        if content_type:
            cursor.execute(
                """
                SELECT id, content_type, title, subject, body_text, body_html, details_json, created_at, updated_at
                FROM draft_content
                WHERE content_type = ?
                ORDER BY id ASC
                """,
                (content_type,),
            )
        else:
            cursor.execute(
                """
                SELECT id, content_type, title, subject, body_text, body_html, details_json, created_at, updated_at
                FROM draft_content
                ORDER BY id ASC
                """
            )
        rows = cursor.fetchall() or []
        payload: list[dict[str, object]] = []
        for row in rows:
            payload.append(
                {
                    "id": int(row["id"]),
                    "content_type": str(row["content_type"]),
                    "title": str(row["title"]),
                    "subject": str(row["subject"]),
                    "body_text": str(row["body_text"]),
                    "body_html": str(row["body_html"]),
                    "details_json": row["details_json"],
                    "created_at": str(row["created_at"]),
                    "updated_at": str(row["updated_at"]),
                    "local_only": True,
                }
            )
        return payload
    finally:
        connection.close()


def _upsert_attachment_state(payload: dict[str, object]) -> None:
    _ensure_local_cache_db()
    connection = sqlite3.connect(LOCAL_CACHE_DB)
    try:
        connection.execute(
            """
            INSERT INTO attachment_state (state_key, payload_json, updated_at)
            VALUES (?, ?, CURRENT_TIMESTAMP)
            ON CONFLICT(state_key) DO UPDATE SET
                payload_json = excluded.payload_json,
                updated_at = CURRENT_TIMESTAMP
            """,
            (LOCAL_ATTACHMENT_STATE_KEY, json.dumps(payload, ensure_ascii=False)),
        )
        connection.commit()
    finally:
        connection.close()


def _load_attachment_state() -> dict[str, object]:
    _ensure_local_cache_db()
    connection = sqlite3.connect(LOCAL_CACHE_DB)
    connection.row_factory = sqlite3.Row
    try:
        cursor = connection.cursor()
        cursor.execute(
            """
            SELECT payload_json
            FROM attachment_state
            WHERE state_key = ?
            """,
            (LOCAL_ATTACHMENT_STATE_KEY,),
        )
        row = cursor.fetchone()
        if not row:
            return {}
        raw_payload = row["payload_json"]
        if not raw_payload:
            return {}
        try:
            payload = json.loads(str(raw_payload))
            return payload if isinstance(payload, dict) else {}
        except Exception:
            return {}
    finally:
        connection.close()


def _delete_attachment_state() -> None:
    _ensure_local_cache_db()
    connection = sqlite3.connect(LOCAL_CACHE_DB)
    try:
        connection.execute(
            "DELETE FROM attachment_state WHERE state_key = ?",
            (LOCAL_ATTACHMENT_STATE_KEY,),
        )
        connection.commit()
    finally:
        connection.close()


def _upsert_ui_state(state_key: str, payload: dict[str, object]) -> None:
    _ensure_local_cache_db()
    connection = sqlite3.connect(LOCAL_CACHE_DB)
    try:
        connection.execute(
            """
            CREATE TABLE IF NOT EXISTS ui_state (
                state_key TEXT PRIMARY KEY,
                payload_json TEXT NOT NULL DEFAULT '{}',
                updated_at TEXT NOT NULL DEFAULT CURRENT_TIMESTAMP
            )
            """
        )
        connection.execute(
            """
            INSERT INTO ui_state (state_key, payload_json, updated_at)
            VALUES (?, ?, CURRENT_TIMESTAMP)
            ON CONFLICT(state_key) DO UPDATE SET
                payload_json = excluded.payload_json,
                updated_at = CURRENT_TIMESTAMP
            """,
            (state_key, json.dumps(payload, ensure_ascii=False)),
        )
        connection.commit()
    finally:
        connection.close()


def _load_ui_state(state_key: str) -> dict[str, object]:
    _ensure_local_cache_db()
    connection = sqlite3.connect(LOCAL_CACHE_DB)
    connection.row_factory = sqlite3.Row
    try:
        connection.execute(
            """
            CREATE TABLE IF NOT EXISTS ui_state (
                state_key TEXT PRIMARY KEY,
                payload_json TEXT NOT NULL DEFAULT '{}',
                updated_at TEXT NOT NULL DEFAULT CURRENT_TIMESTAMP
            )
            """
        )
        cursor = connection.cursor()
        cursor.execute(
            """
            SELECT payload_json
            FROM ui_state
            WHERE state_key = ?
            """,
            (state_key,),
        )
        row = cursor.fetchone()
        if not row:
            return {}
        raw_payload = str(row["payload_json"] or "")
        if not raw_payload:
            return {}
        try:
            payload = json.loads(raw_payload)
            return payload if isinstance(payload, dict) else {}
        except Exception:
            return {}
    finally:
        connection.close()


def _delete_ui_state(state_key: str) -> None:
    _ensure_local_cache_db()
    connection = sqlite3.connect(LOCAL_CACHE_DB)
    try:
        connection.execute(
            """
            CREATE TABLE IF NOT EXISTS ui_state (
                state_key TEXT PRIMARY KEY,
                payload_json TEXT NOT NULL DEFAULT '{}',
                updated_at TEXT NOT NULL DEFAULT CURRENT_TIMESTAMP
            )
            """
        )
        connection.execute("DELETE FROM ui_state WHERE state_key = ?", (state_key,))
        connection.commit()
    finally:
        connection.close()


def _upsert_tag_state(payload: dict[str, object], state_key: str = LOCAL_TAG_STATE_KEY) -> None:
    _ensure_local_cache_db()
    connection = sqlite3.connect(LOCAL_CACHE_DB)
    try:
        connection.execute(
            """
            INSERT INTO tag_state (state_key, payload_json, updated_at)
            VALUES (?, ?, CURRENT_TIMESTAMP)
            ON CONFLICT(state_key) DO UPDATE SET
                payload_json = excluded.payload_json,
                updated_at = CURRENT_TIMESTAMP
            """,
            (state_key, json.dumps(payload, ensure_ascii=False)),
        )
        connection.commit()
    finally:
        connection.close()


def _load_tag_state(state_key: str = LOCAL_TAG_STATE_KEY) -> dict[str, object]:
    _ensure_local_cache_db()
    connection = sqlite3.connect(LOCAL_CACHE_DB)
    connection.row_factory = sqlite3.Row
    try:
        cursor = connection.cursor()
        cursor.execute(
            """
            SELECT payload_json
            FROM tag_state
            WHERE state_key = ?
            """,
            (state_key,),
        )
        row = cursor.fetchone()
        if not row:
            return {}
        raw_payload = str(row["payload_json"] or "")
        if not raw_payload:
            return {}
        try:
            payload = json.loads(raw_payload)
            return payload if isinstance(payload, dict) else {}
        except Exception:
            return {}
    finally:
        connection.close()


def _delete_tag_state(state_key: str = LOCAL_TAG_STATE_KEY) -> None:
    _ensure_local_cache_db()
    connection = sqlite3.connect(LOCAL_CACHE_DB)
    try:
        connection.execute("DELETE FROM tag_state WHERE state_key = ?", (state_key,))
        connection.commit()
    finally:
        connection.close()


def _upsert_customer_variables(email: str, variables: dict[str, object]) -> None:
    _ensure_local_cache_db()
    email_key = (email or "").strip().lower()
    if not email_key:
        return
    connection = sqlite3.connect(LOCAL_CACHE_DB)
    try:
        connection.execute(
            """
            INSERT INTO customer_variables (email, variables_json, updated_at)
            VALUES (?, ?, CURRENT_TIMESTAMP)
            ON CONFLICT(email) DO UPDATE SET
                variables_json = excluded.variables_json,
                updated_at = CURRENT_TIMESTAMP
            """,
            (email_key, json.dumps(variables or {}, ensure_ascii=False)),
        )
        connection.commit()
    finally:
        connection.close()


def _load_customer_variables(email: str) -> dict[str, object]:
    _ensure_local_cache_db()
    email_key = (email or "").strip().lower()
    if not email_key:
        return {}
    connection = sqlite3.connect(LOCAL_CACHE_DB)
    connection.row_factory = sqlite3.Row
    try:
        cursor = connection.cursor()
        cursor.execute(
            """
            SELECT variables_json
            FROM customer_variables
            WHERE email = ?
            """,
            (email_key,),
        )
        row = cursor.fetchone()
        if not row:
            return {}
        raw_payload = str(row["variables_json"] or "")
        if not raw_payload:
            return {}
        try:
            payload = json.loads(raw_payload)
            return payload if isinstance(payload, dict) else {}
        except Exception:
            return {}
    finally:
        connection.close()


def _load_all_customer_variables() -> dict[str, dict[str, object]]:
    _ensure_local_cache_db()
    connection = sqlite3.connect(LOCAL_CACHE_DB)
    connection.row_factory = sqlite3.Row
    try:
        cursor = connection.cursor()
        cursor.execute(
            """
            SELECT email, variables_json
            FROM customer_variables
            ORDER BY email ASC
            """
        )
        rows = cursor.fetchall() or []
        result: dict[str, dict[str, object]] = {}
        for row in rows:
            email = str(row["email"] or "").strip().lower()
            if not email:
                continue
            raw_payload = str(row["variables_json"] or "")
            payload: dict[str, object] = {}
            if raw_payload:
                try:
                    decoded = json.loads(raw_payload)
                    if isinstance(decoded, dict):
                        payload = decoded
                except Exception:
                    payload = {}
            result[email] = payload
        return result
    finally:
        connection.close()


def _delete_customer_variables(email: str | None = None) -> None:
    _ensure_local_cache_db()
    connection = sqlite3.connect(LOCAL_CACHE_DB)
    try:
        if email:
            connection.execute("DELETE FROM customer_variables WHERE email = ?", ((email or "").strip().lower(),))
        else:
            connection.execute("DELETE FROM customer_variables")
        connection.commit()
    finally:
        connection.close()


def _upsert_local_browser_session(
    username: str,
    session_id: str,
    title: str,
    browser_name: str,
    browser_mode: str,
    status: str,
    browser_pid: int | None = None,
    launch_preset: str = "Default",
    details: dict[str, object] | None = None,
) -> None:
    _ensure_local_cache_db()
    connection = sqlite3.connect(LOCAL_CACHE_DB)
    try:
        connection.execute(
            """
            INSERT INTO browser_sessions (
                session_id, username, title, browser_name, browser_mode, status,
                browser_pid, launch_preset, details_json, started_at, updated_at, closed_at
            )
            VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, CURRENT_TIMESTAMP, CURRENT_TIMESTAMP,
                    CASE WHEN ? IN ('Closed', 'Stopped') THEN CURRENT_TIMESTAMP ELSE NULL END)
            ON CONFLICT(session_id) DO UPDATE SET
                username = excluded.username,
                title = excluded.title,
                browser_name = excluded.browser_name,
                browser_mode = excluded.browser_mode,
                status = excluded.status,
                browser_pid = excluded.browser_pid,
                launch_preset = excluded.launch_preset,
                details_json = excluded.details_json,
                updated_at = CURRENT_TIMESTAMP,
                closed_at = CASE
                    WHEN excluded.status IN ('Closed', 'Stopped') THEN CURRENT_TIMESTAMP
                    ELSE browser_sessions.closed_at
                END
            """,
            (
                username,
                session_id,
                title,
                browser_name,
                browser_mode,
                status,
                browser_pid,
                launch_preset,
                json.dumps(details or {}, ensure_ascii=False),
                status,
            ),
        )
        connection.commit()
    finally:
        connection.close()


def _is_local_draft_item(item: QTableWidgetItem | QListWidgetItem | None) -> bool:
    if item is None:
        return False
    return bool(item.data(ROLE_LOCAL_ONLY))


def _local_draft_id_from_item(item: QTableWidgetItem | QListWidgetItem | None) -> int | None:
    if item is None:
        return None
    value = item.data(ROLE_LOCAL_DRAFT_ID)
    return int(value) if value else None


@dataclass
class AppState:
    username: str = ""
    role: str = ""
    logged_in: bool = False
    auth_token: str = ""
    browser_mode: str = "Incognito"
    window_count: int = 1
    launch_preset: str = "Default"
    active_sessions: list[str] = field(default_factory=list)
    activity_log: list[str] = field(default_factory=list)
    pending_recipients: list[str] = field(default_factory=list)
    pending_emails_validated: bool = False
    custom_tag_1: str = ""
    custom_tag_2: str = ""
    tag_samples: dict[str, str] = field(default_factory=dict)
    body_mode: str = "Normal Message"
    subject_text: str = ""
    plain_body_text: str = ""
    html_message_text: str = ""
    html_template_text: str = ""
    ai_provider: str = "ChatGPT"
    ai_api_key: str = ""
    ai_model: str = ""
    ai_connected: bool = False
    sender_limit: int = 300
    delay_from: float = 0.5
    delay_to: float = 1.0
    retry_count: int = 3
    retry_enabled: bool = True
    delay_type: str = "Random range"
    email_send_order: str = "Sequential"
    window_send_mode: str = "Parallel"
    ai_available_models: list[str] = field(default_factory=list)


@dataclass
class BrowserSessionHandle:
    session_id: str
    title: str
    mode: str
    browser_name: str = "Google Chrome"
    process: subprocess.Popen[str] | None = None
    status: str = "Starting"
    profile_dir: Path | None = None
    debug_port: int | None = None
    send_completed: int = 0
    send_total: int = 0

    def is_alive(self) -> bool:
        return self.process is not None and self.process.poll() is None


class AIValidationWorker(QObject):
    finished = Signal(str, list)
    failed = Signal(str)

    def __init__(self, provider: str, api_key: str):
        super().__init__()
        self.provider = provider
        self.api_key = api_key

    def run(self) -> None:
        try:
            models = self._fetch_models()
        except Exception as exc:
            self.failed.emit(str(exc))
            return
        self.finished.emit(self.provider, models)

    def _fetch_models(self) -> list[str]:
        provider = self.provider
        context = self._ssl_context()
        if provider == "Claude":
            request = urllib.request.Request(
                "https://api.anthropic.com/v1/models",
                headers={
                    "x-api-key": self.api_key,
                    "anthropic-version": "2023-06-01",
                    "accept": "application/json",
                },
                method="GET",
            )
        elif provider == "DeepSeek":
            request = urllib.request.Request(
                "https://api.deepseek.com/models",
                headers={
                    "Authorization": f"Bearer {self.api_key}",
                    "accept": "application/json",
                },
                method="GET",
            )
        else:
            request = urllib.request.Request(
                "https://api.openai.com/v1/models",
                headers={
                    "Authorization": f"Bearer {self.api_key}",
                    "accept": "application/json",
                },
                method="GET",
            )

        with urllib.request.urlopen(request, timeout=20, context=context) as response:
            payload = json.loads(response.read().decode("utf-8"))

        return self._extract_model_ids(payload)

    def _ssl_context(self) -> ssl.SSLContext:
        try:
            import certifi

            return ssl.create_default_context(cafile=certifi.where())
        except Exception:
            return ssl.create_default_context()

    def _extract_model_ids(self, payload) -> list[str]:
        models: list[str] = []
        if isinstance(payload, dict):
            entries = payload.get("data") or payload.get("models") or []
        elif isinstance(payload, list):
            entries = payload
        else:
            entries = []

        for item in entries:
            if isinstance(item, dict):
                model_id = item.get("id") or item.get("name") or item.get("model")
                if model_id:
                    models.append(str(model_id))

        unique_models = list(dict.fromkeys(models))
        if not unique_models:
            raise ValueError("No models were returned by the provider.")
        return unique_models


class CampaignSendWorker(QObject):
    log = Signal(str)
    progress = Signal(str, str, int, int)
    finished = Signal(str, str, bool, int, int, str)

    def __init__(
        self,
        session: BrowserSessionHandle,
        tasks: list[dict[str, str]],
        pause_event: threading.Event,
        cancel_event: threading.Event,
        send_callback: Callable[..., None],
        *,
        window_label: str,
        delay_mode: str,
        delay_from: float,
        delay_to: float,
        retry_count: int,
        retry_enabled: bool,
        convert_enabled: bool,
        attachment_formats: list[str],
        file_name_mode: str,
        window_send_mode: str,
    ):
        super().__init__()
        self.session = session
        self.tasks = tasks
        self.pause_event = pause_event
        self.cancel_event = cancel_event
        self.send_callback = send_callback
        self.window_label = window_label
        self.delay_mode = delay_mode
        self.delay_from = max(0.0, float(delay_from))
        self.delay_to = max(self.delay_from, float(delay_to))
        self.retry_count = max(0, int(retry_count))
        self.retry_enabled = bool(retry_enabled)
        self.convert_enabled = bool(convert_enabled)
        self.attachment_formats = list(attachment_formats)
        self.file_name_mode = file_name_mode
        self.window_send_mode = window_send_mode

    def _timestamp(self) -> str:
        return QDateTime.currentDateTime().toString("hh:mm:ss")

    def _wait_for_resume(self) -> bool:
        while not self.cancel_event.is_set():
            if self.pause_event.is_set():
                return True
            time.sleep(0.1)
        return False

    def _sleep_with_controls(self, seconds: float) -> bool:
        deadline = time.time() + max(0.0, float(seconds))
        while time.time() < deadline:
            if self.cancel_event.is_set():
                return False
            if not self._wait_for_resume():
                return False
            time.sleep(0.1)
        return not self.cancel_event.is_set()

    def _delay_seconds(self) -> float:
        if self.delay_mode == "Fixed":
            return self.delay_from
        if self.delay_mode == "Human-like pattern":
            base = random.uniform(self.delay_from, self.delay_to)
            return max(self.delay_from, min(self.delay_to, base * random.uniform(0.75, 1.15)))
        return random.uniform(self.delay_from, self.delay_to)

    def run(self) -> None:
        completed = 0
        total = len(self.tasks)
        error_message = ""
        for task in self.tasks:
            if self.cancel_event.is_set():
                break
            if not self._wait_for_resume():
                break

            recipient = task["recipient"]
            subject = task["subject"]
            body_text = task["body_text"]
            attachment_html = task["attachment_html"]
            attachment_formats = self.attachment_formats or [task["attachment_format"]]
            file_name_value = task["file_name_value"]
            attempts = self.retry_count + 1 if self.retry_enabled else 1
            last_error: Exception | None = None
            sent_ok = False

            for attempt in range(1, attempts + 1):
                if self.cancel_event.is_set() or not self._wait_for_resume():
                    break
                try:
                    self.send_callback(
                        self.session,
                        recipient,
                        subject,
                        body_text,
                        attachment_html,
                        attachment_formats,
                        self.file_name_mode,
                        file_name_value,
                        self.convert_enabled,
                        True,
                        False,
                    )
                    sent_ok = True
                    break
                except Exception as exc:
                    last_error = exc
                    if attempt < attempts and not self.cancel_event.is_set():
                        if not self._sleep_with_controls(self._delay_seconds()):
                            break

            completed += 1
            if sent_ok:
                self.log.emit(f"[{self._timestamp()}] {self.window_label} sent {recipient}")
            elif last_error is not None:
                error_message = str(last_error)
                self.log.emit(f"[{self._timestamp()}] {self.window_label} failed {recipient}: {error_message}")
            else:
                self.log.emit(f"[{self._timestamp()}] {self.window_label} cancelled {recipient}")
            self.progress.emit(self.session.session_id, self.window_label, completed, total)

            if completed < total and not self.cancel_event.is_set():
                if not self._sleep_with_controls(self._delay_seconds()):
                    break

        self.finished.emit(
            self.session.session_id,
            self.window_label,
            self.cancel_event.is_set(),
            completed,
            total,
            error_message,
        )

class AnimatedLogoBadge(QWidget):
    def __init__(self, parent=None, scale: float = 1.0):
        super().__init__(parent)
        self._scale = scale
        self._pulse = 0.0
        self.setFixedSize(_scaled_int(40, self._scale), _scaled_int(40, self._scale))
        self._anim = QPropertyAnimation(self, b"pulse", self)
        self._anim.setStartValue(0.0)
        self._anim.setEndValue(1.0)
        self._anim.setDuration(1800)
        self._anim.setLoopCount(-1)
        self._anim.setEasingCurve(QEasingCurve.InOutSine)
        self._anim.start()

    def getPulse(self) -> float:
        return self._pulse

    def setPulse(self, value: float) -> None:
        self._pulse = value
        self.update()

    pulse = Property(float, getPulse, setPulse)

    def paintEvent(self, _event) -> None:
        painter = QPainter(self)
        painter.setRenderHint(QPainter.Antialiasing)
        pulse = 0.35 + (self._pulse * 0.65)

        painter.setPen(Qt.NoPen)
        painter.setBrush(QColor(37, 99, 235, int(42 + pulse * 80)))
        painter.drawEllipse(self.rect().adjusted(1, 1, -1, -1))

        painter.setBrush(QColor("#2563eb"))
        painter.drawEllipse(self.rect().adjusted(5, 5, -5, -5))

        pen_width = max(1.0, 1.4 * self._scale)
        painter.setPen(QPen(QColor("#67e8f9"), pen_width))
        painter.drawLine(_scaled_int(18, self._scale), _scaled_int(10, self._scale), _scaled_int(15, self._scale), _scaled_int(19, self._scale))
        painter.drawLine(_scaled_int(15, self._scale), _scaled_int(19, self._scale), _scaled_int(21, self._scale), _scaled_int(19, self._scale))
        painter.drawLine(_scaled_int(21, self._scale), _scaled_int(19, self._scale), _scaled_int(18, self._scale), _scaled_int(29, self._scale))

        painter.setPen(QColor("#ffffff"))
        painter.setFont(QFont("Segoe UI", _scaled_int(10, self._scale), QFont.Bold))
        painter.drawText(self.rect(), Qt.AlignCenter, "EZ")


class MacTrafficLightButton(QPushButton):
    def __init__(self, kind: str, scale: float = 1.0, parent=None):
        super().__init__(parent)
        self._kind = kind
        self._scale = scale
        self._hover = False
        self.setObjectName("macTrafficLightButton")
        self.setCursor(Qt.PointingHandCursor)
        self.setFlat(True)
        self.setText("")
        size = _scaled_int(14, self._scale)
        self.setFixedSize(size, size)
        self.setToolTip({
            "close": "Close the window",
            "minimize": "Minimize the window",
            "maximize": "Maximize or restore the window",
        }.get(kind, "Window control"))

    def enterEvent(self, event) -> None:
        self._hover = True
        self.update()
        super().enterEvent(event)

    def leaveEvent(self, event) -> None:
        self._hover = False
        self.update()
        super().leaveEvent(event)

    def paintEvent(self, _event) -> None:
        painter = QPainter(self)
        painter.setRenderHint(QPainter.Antialiasing)

        colors = {
            "close": QColor("#ff5f57"),
            "minimize": QColor("#febc2e"),
            "maximize": QColor("#28c840"),
        }
        base = colors.get(self._kind, QColor("#7a7a7a"))
        if self._hover:
            base = base.lighter(110)

        painter.setPen(Qt.NoPen)
        painter.setBrush(base)
        painter.drawEllipse(self.rect().adjusted(0, 0, -1, -1))

        if self._hover:
            painter.setPen(QPen(QColor(0, 0, 0, 180), max(1, _scaled_int(1, self._scale))))
            w = self.width()
            h = self.height()
            if self._kind == "close":
                painter.drawLine(_scaled_int(4, self._scale), _scaled_int(4, self._scale), w - _scaled_int(4, self._scale), h - _scaled_int(4, self._scale))
                painter.drawLine(w - _scaled_int(4, self._scale), _scaled_int(4, self._scale), _scaled_int(4, self._scale), h - _scaled_int(4, self._scale))
            elif self._kind == "minimize":
                painter.drawLine(_scaled_int(4, self._scale), h // 2, w - _scaled_int(4, self._scale), h // 2)
            else:
                painter.drawLine(_scaled_int(4, self._scale), h // 2, w - _scaled_int(4, self._scale), h // 2)
                painter.drawLine(w // 2, _scaled_int(4, self._scale), w // 2, h - _scaled_int(4, self._scale))


class TitleBar(QWidget):
    def __init__(self, window, on_close, scale: float = 1.0):
        super().__init__()
        self._window = window
        self._on_close = on_close
        self._drag_pos = None
        self._is_maximized = False
        self._scale = scale
        self.setObjectName("topBar")
        self._build_ui()

    def _build_ui(self) -> None:
        layout = QHBoxLayout(self)
        layout.setContentsMargins(
            _scaled_int(8, self._scale),
            _scaled_int(1, self._scale),
            _scaled_int(8, self._scale),
            _scaled_int(1, self._scale),
        )
        layout.setSpacing(_scaled_int(4, self._scale))

        controls = QHBoxLayout()
        controls.setSpacing(_scaled_int(5, self._scale))
        self.close_button = MacTrafficLightButton("close", self._scale)
        self.close_button.clicked.connect(self._on_close)
        self.minimize_button = MacTrafficLightButton("minimize", self._scale)
        self.minimize_button.clicked.connect(self._window.showMinimized)
        self.maximize_button = MacTrafficLightButton("maximize", self._scale)
        self.maximize_button.clicked.connect(self._toggle_maximize)
        controls.addWidget(self.close_button)
        controls.addWidget(self.minimize_button)
        controls.addWidget(self.maximize_button)
        controls.addSpacing(_scaled_int(6, self._scale))

        brand = QHBoxLayout()
        brand.setSpacing(_scaled_int(6, self._scale))
        title = QLabel("EazyMailer")
        title.setObjectName("brandTitle")
        title.setStyleSheet("font-weight: 800;")
        brand.addWidget(title)

        spacer = QWidget()
        spacer.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Preferred)

        version_badge = QLabel("v2.0")
        version_badge.setObjectName("versionBadge")
        self.status_badge = QLabel("LOCKED")
        self.status_badge.setObjectName("statusBadge")
        self.theme_badge = QLabel("◐")
        self.theme_badge.setObjectName("statusBadge")
        self.theme_badge.setToolTip("System appearance")
        self.user_id_label = QLabel("")
        self.user_id_label.setObjectName("statusBadge")
        for badge in (version_badge, self.status_badge):
            badge.setFixedHeight(_scaled_int(24 if IS_MAC else 28, self._scale))
            badge.setMinimumWidth(_scaled_int(58 if IS_MAC else 64, self._scale))
            badge.setAlignment(Qt.AlignCenter)
        version_badge.setToolTip("Application version")
        self.status_badge.setToolTip("Current login status")
        version_badge.hide()
        self.status_badge.hide()
        self.theme_badge.setFixedHeight(_scaled_int(24 if IS_MAC else 28, self._scale))
        self.user_id_label.setFixedHeight(_scaled_int(24 if IS_MAC else 28, self._scale))

        self.logout_button = QPushButton("Logout")
        self.logout_button.setObjectName("secondaryButton")
        self.logout_button.setFixedHeight(_scaled_int(24, self._scale))
        self.logout_button.setMinimumWidth(_scaled_int(58, self._scale))
        self.logout_button.setToolTip("Sign out and return to login")

        layout.addLayout(controls)
        layout.addLayout(brand)
        layout.addWidget(spacer)
        layout.addWidget(version_badge)
        layout.addWidget(self.status_badge)
        layout.addWidget(self.theme_badge)
        layout.addWidget(self.user_id_label)
        layout.addWidget(self.logout_button)

    def set_state(self, username: str, logged_in: bool) -> None:
        self.status_badge.setText("READY" if logged_in else "LOCKED")
        self.user_id_label.setText(username if logged_in else "")
        self.theme_badge.setText("☾" if self._window._system_is_dark() else "☀")
        self.theme_badge.setVisible(logged_in)
        self.user_id_label.setVisible(logged_in)
        self.logout_button.setVisible(logged_in)

    def sync_window_state(self) -> None:
        self._is_maximized = self._window.isMaximized()
        self.maximize_button.setToolTip("Restore window" if self._is_maximized else "Maximize window")

    def _toggle_maximize(self) -> None:
        if self._window.isMaximized():
            self._window.showNormal()
        else:
            self._window.showMaximized()
        self.sync_window_state()

    def set_logout_handler(self, handler) -> None:
        self.logout_button.clicked.connect(handler)

    def mousePressEvent(self, event) -> None:
        if event.button() == Qt.LeftButton:
            self._drag_pos = event.globalPosition().toPoint() - self._window.frameGeometry().topLeft()
            event.accept()

    def mouseMoveEvent(self, event) -> None:
        if self._drag_pos is not None and event.buttons() & Qt.LeftButton:
            self._window.move(event.globalPosition().toPoint() - self._drag_pos)
            event.accept()

    def mouseReleaseEvent(self, event) -> None:
        self._drag_pos = None
        event.accept()

    def mouseDoubleClickEvent(self, event) -> None:
        if event.button() == Qt.LeftButton:
            self._toggle_maximize()
            event.accept()


class Toast(QFrame):
    def __init__(self, parent: QWidget, message: str, kind: str = "info", scale: float = 1.0):
        super().__init__(parent)
        self._scale = scale
        self.setObjectName("toast")
        self.setProperty("kind", kind)
        self.setAttribute(Qt.WA_StyledBackground, True)
        self.setWindowFlags(Qt.SubWindow | Qt.FramelessWindowHint)
        self.setFixedWidth(_scaled_int(360, self._scale))
        self._effect = QGraphicsOpacityEffect(self)
        self.setGraphicsEffect(self._effect)
        self._effect.setOpacity(0.0)
        self._build_ui(message, kind)

        self._fade_in = QPropertyAnimation(self._effect, b"opacity", self)
        self._fade_in.setDuration(180)
        self._fade_in.setStartValue(0.0)
        self._fade_in.setEndValue(1.0)

        self._fade_out = QPropertyAnimation(self._effect, b"opacity", self)
        self._fade_out.setDuration(220)
        self._fade_out.setStartValue(1.0)
        self._fade_out.setEndValue(0.0)
        self._fade_out.finished.connect(self.deleteLater)

        self._timer = QTimer(self)
        self._timer.setSingleShot(True)
        self._timer.timeout.connect(self.close_out)

    def _build_ui(self, message: str, kind: str) -> None:
        layout = QHBoxLayout(self)
        layout.setContentsMargins(_scaled_int(12, self._scale), _scaled_int(8, self._scale), _scaled_int(12, self._scale), _scaled_int(8, self._scale))
        layout.setSpacing(_scaled_int(8, self._scale))

        icon = QLabel("●")
        icon.setObjectName("toastIcon")
        icon.setFixedSize(_scaled_int(14, self._scale), _scaled_int(14, self._scale))
        icon.setAlignment(Qt.AlignCenter)
        if kind == "warning":
            icon.setText("!")
        elif kind == "success":
            icon.setText("✓")
        elif kind == "error":
            icon.setText("×")

        text = QLabel(message)
        text.setWordWrap(False)
        text.setObjectName("toastText")
        text.setAlignment(Qt.AlignVCenter | Qt.AlignLeft)
        text.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Preferred)

        layout.addWidget(icon)
        layout.addWidget(text, 1)

    def showEvent(self, event) -> None:
        super().showEvent(event)
        self._fade_in.start()
        self._timer.start(2200)

    def close_out(self) -> None:
        self._fade_out.start()


class RobotLoaderBadge(QWidget):
    def __init__(self, parent=None, scale: float = 1.0):
        super().__init__(parent)
        self._scale = scale
        self._pulse = 0.0
        self._blink = False
        self.setFixedSize(_scaled_int(104, self._scale), _scaled_int(104, self._scale))
        self._pulse_anim = QPropertyAnimation(self, b"pulse", self)
        self._pulse_anim.setStartValue(0.0)
        self._pulse_anim.setEndValue(1.0)
        self._pulse_anim.setDuration(1500)
        self._pulse_anim.setLoopCount(-1)
        self._pulse_anim.setEasingCurve(QEasingCurve.InOutSine)
        self._pulse_anim.start()

        self._blink_timer = QTimer(self)
        self._blink_timer.setInterval(520)
        self._blink_timer.timeout.connect(self._toggle_blink)
        self._blink_timer.start()

    def _toggle_blink(self) -> None:
        self._blink = not self._blink
        self.update()

    def getPulse(self) -> float:
        return self._pulse

    def setPulse(self, value: float) -> None:
        self._pulse = value
        self.update()

    pulse = Property(float, getPulse, setPulse)

    def paintEvent(self, _event) -> None:
        painter = QPainter(self)
        painter.setRenderHint(QPainter.Antialiasing)
        rect = self.rect()
        pulse = 0.35 + (self._pulse * 0.65)

        painter.setPen(Qt.NoPen)
        painter.setBrush(QColor(14, 99, 156, int(34 + pulse * 60)))
        inset = _scaled_int(4, self._scale)
        painter.drawEllipse(rect.adjusted(inset, inset, -inset, -inset))

        # antenna
        painter.setPen(QPen(QColor("#8bd5ff"), 2))
        painter.drawLine(rect.center().x(), _scaled_int(10, self._scale), rect.center().x(), _scaled_int(24, self._scale))
        painter.setBrush(QColor("#f9fafb"))
        painter.drawEllipse(rect.center().x() - _scaled_int(4, self._scale), _scaled_int(6, self._scale), _scaled_int(8, self._scale), _scaled_int(8, self._scale))

        # head
        head = rect.adjusted(_scaled_int(18, self._scale), _scaled_int(24, self._scale), -_scaled_int(18, self._scale), -_scaled_int(24, self._scale))
        painter.setBrush(QColor("#2d2d30"))
        painter.setPen(QPen(QColor("#4b4b4b"), 1))
        painter.drawRoundedRect(head, 18, 18)

        # eyes
        eye_y = head.center().y() - _scaled_int(10, self._scale)
        eye_color = QColor("#9cdcfe") if not self._blink else QColor("#3a3d41")
        painter.setBrush(eye_color)
        painter.setPen(Qt.NoPen)
        painter.drawRoundedRect(head.center().x() - _scaled_int(19, self._scale), eye_y, _scaled_int(12, self._scale), _scaled_int(12, self._scale), _scaled_int(4, self._scale), _scaled_int(4, self._scale))
        painter.drawRoundedRect(head.center().x() + _scaled_int(7, self._scale), eye_y, _scaled_int(12, self._scale), _scaled_int(12, self._scale), _scaled_int(4, self._scale), _scaled_int(4, self._scale))

        # mouth / badge line
        painter.setBrush(QColor("#0e639c"))
        painter.drawRoundedRect(head.center().x() - _scaled_int(18, self._scale), head.center().y() + _scaled_int(10, self._scale), _scaled_int(36, self._scale), _scaled_int(8, self._scale), _scaled_int(4, self._scale), _scaled_int(4, self._scale))

        # body
        body = rect.adjusted(_scaled_int(30, self._scale), _scaled_int(56, self._scale), -_scaled_int(30, self._scale), -_scaled_int(16, self._scale))
        painter.setBrush(QColor("#1e1e1e"))
        painter.setPen(QPen(QColor("#3c3c3c"), 1))
        painter.drawRoundedRect(body, 10, 10)
        painter.setPen(QPen(QColor("#6b7280"), 2))
        painter.drawLine(body.left() + 10, body.bottom() - 8, body.right() - 10, body.bottom() - 8)
        painter.drawLine(body.left() + 14, body.bottom() - 4, body.left() + 14, body.bottom() + 4)
        painter.drawLine(body.right() - 14, body.bottom() - 4, body.right() - 14, body.bottom() + 4)


class LaunchLoaderDialog(QDialog):
    def __init__(self, parent=None, scale: float = 1.0):
        super().__init__(parent)
        self._scale = scale
        self._status_prefix = "Preparing"
        self.setModal(False)
        self.setWindowFlags(Qt.Dialog | Qt.FramelessWindowHint | Qt.Tool)
        self.setAttribute(Qt.WA_TranslucentBackground, True)
        self.setObjectName("launchLoader")
        self._dots = 0
        self._build_ui()

        self._dot_timer = QTimer(self)
        self._dot_timer.setInterval(_scaled_int(260, self._scale))
        self._dot_timer.timeout.connect(self._animate_dots)

    def _build_ui(self) -> None:
        root = QVBoxLayout(self)
        root.setContentsMargins(_scaled_int(24, self._scale), _scaled_int(24, self._scale), _scaled_int(24, self._scale), _scaled_int(24, self._scale))
        root.setSpacing(0)

        root.addStretch()

        card = QFrame()
        card.setObjectName("loaderCard")
        card.setFixedWidth(_scaled_int(360, self._scale))
        card_layout = QVBoxLayout(card)
        card_layout.setContentsMargins(_scaled_int(18, self._scale), _scaled_int(18, self._scale), _scaled_int(18, self._scale), _scaled_int(18, self._scale))
        card_layout.setSpacing(_scaled_int(10, self._scale))
        card_layout.setAlignment(Qt.AlignCenter)

        self.robot = RobotLoaderBadge(scale=self._scale)
        self.loader_title = QLabel("Launching browser windows")
        self.loader_title.setObjectName("loaderTitle")
        self.loader_title.setAlignment(Qt.AlignCenter)
        self.loader_subtitle = QLabel("Applying browser mode and launch preset.")
        self.loader_subtitle.setObjectName("loaderSubtitle")
        self.loader_subtitle.setAlignment(Qt.AlignCenter)
        self.loader_subtitle.setWordWrap(True)
        self.loader_status = QLabel("Preparing")
        self.loader_status.setObjectName("loaderStatus")
        self.loader_status.setAlignment(Qt.AlignCenter)
        self.loader_progress = QProgressBar()
        self.loader_progress.setObjectName("loaderProgress")
        self.loader_progress.setRange(0, 0)
        self.loader_progress.setValue(0)
        self.loader_progress.setTextVisible(True)

        card_layout.addWidget(self.robot, alignment=Qt.AlignCenter)
        card_layout.addWidget(self.loader_title)
        card_layout.addWidget(self.loader_subtitle)
        card_layout.addWidget(self.loader_progress)
        card_layout.addWidget(self.loader_status)

        root.addWidget(card, alignment=Qt.AlignCenter)
        root.addStretch()

    def set_message(self, title: str, subtitle: str | None = None, status_prefix: str | None = None) -> None:
        self.loader_title.setText(title)
        if subtitle:
            self.loader_subtitle.setText(subtitle)
        if status_prefix:
            self._status_prefix = status_prefix
            self.loader_status.setText(status_prefix)
        self.loader_progress.setRange(0, 0)
        self.loader_progress.setValue(0)

    def set_busy(self, title: str | None = None, subtitle: str | None = None, status_prefix: str = "Preparing") -> None:
        if title:
            self.loader_title.setText(title)
        if subtitle:
            self.loader_subtitle.setText(subtitle)
        self._status_prefix = status_prefix
        self.loader_status.setText(status_prefix)
        self.loader_progress.setRange(0, 0)
        self.loader_progress.setValue(0)

    def set_progress(self, value: int, total: int = 100, status_prefix: str | None = None) -> None:
        total = max(total, 1)
        value = max(0, min(value, total))
        self.loader_progress.setRange(0, total)
        self.loader_progress.setValue(value)
        if status_prefix:
            self._status_prefix = status_prefix
        self.loader_status.setText(self._status_prefix)

    def set_status_prefix(self, prefix: str) -> None:
        self._status_prefix = prefix
        self.loader_status.setText(prefix)

    def _animate_dots(self) -> None:
        self._dots = (self._dots + 1) % 4
        self.loader_status.setText(self._status_prefix + ("." * self._dots))

    def showEvent(self, event) -> None:
        super().showEvent(event)
        parent = self.parentWidget()
        if parent is not None:
            self.setGeometry(parent.frameGeometry())
        self._dot_timer.start()

    def closeEvent(self, event) -> None:
        self._dot_timer.stop()
        super().closeEvent(event)


class ConfirmDialog(QDialog):
    def __init__(self, parent: QWidget, title: str, message: str, scale: float = 1.0):
        super().__init__(parent)
        self._scale = scale
        self.setModal(True)
        self.setWindowFlags(Qt.Dialog | Qt.FramelessWindowHint)
        self.setObjectName("confirmDialog")
        self._build_ui(title, message)

    def _build_ui(self, title: str, message: str) -> None:
        layout = QVBoxLayout(self)
        layout.setContentsMargins(_scaled_int(16, self._scale), _scaled_int(16, self._scale), _scaled_int(16, self._scale), _scaled_int(16, self._scale))
        layout.setSpacing(_scaled_int(10, self._scale))

        card = QFrame()
        card.setObjectName("confirmCard")
        card_layout = QVBoxLayout(card)
        card_layout.setContentsMargins(_scaled_int(16, self._scale), _scaled_int(16, self._scale), _scaled_int(16, self._scale), _scaled_int(16, self._scale))
        card_layout.setSpacing(_scaled_int(10, self._scale))

        title_label = QLabel(title)
        title_label.setObjectName("confirmTitle")
        title_label.setAlignment(Qt.AlignCenter)
        message_label = QLabel(message)
        message_label.setObjectName("confirmText")
        message_label.setWordWrap(True)
        message_label.setAlignment(Qt.AlignCenter)

        buttons = QDialogButtonBox(QDialogButtonBox.Yes | QDialogButtonBox.No)
        buttons.button(QDialogButtonBox.Yes).setText("Confirm")
        buttons.button(QDialogButtonBox.No).setText("Cancel")
        buttons.accepted.connect(self.accept)
        buttons.rejected.connect(self.reject)

        card_layout.addWidget(title_label)
        card_layout.addWidget(message_label)
        card_layout.addWidget(buttons)
        layout.addWidget(card)


class BodyDraftEditor(QWidget):
    contentChanged = Signal()
    titleChanged = Signal(str)
    modeChanged = Signal(str)
    previewRequested = Signal()

    def __init__(self, scale: float = 1.0):
        super().__init__()
        self._scale = scale
        self.draft_id: int | None = None
        self.local_draft_id: int | None = None
        self.local_only = False
        self._build_ui()

    def _build_ui(self) -> None:
        layout = QVBoxLayout(self)
        layout.setContentsMargins(_scaled_int(10, self._scale), _scaled_int(10, self._scale), _scaled_int(10, self._scale), _scaled_int(10, self._scale))
        layout.setSpacing(_scaled_int(8, self._scale))
        self.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Expanding)

        title_row = QHBoxLayout()
        title_label = QLabel("Tab label")
        title_label.setObjectName("fieldLabel")
        self.title_input = QLineEdit()
        self.title_input.setPlaceholderText("Tab label")
        self.title_input.setToolTip("Name this body tab")
        title_row.addWidget(title_label)
        title_row.addWidget(self.title_input, 1)
        layout.addLayout(title_row)

        mode_row = QHBoxLayout()
        mode_label = QLabel("Mode")
        mode_label.setObjectName("fieldLabel")
        self.plain_button = QPushButton("Plain Text")
        self.html_button = QPushButton("HTML Body")
        self.preview_button = QPushButton("Preview")
        self.plain_button.setCheckable(True)
        self.html_button.setCheckable(True)
        self.plain_button.setChecked(True)
        self._plain_group = QButtonGroup(self)
        self._plain_group.setExclusive(True)
        self._plain_group.addButton(self.plain_button)
        self._plain_group.addButton(self.html_button)
        self.plain_button.clicked.connect(lambda: self.set_mode("Normal Message"))
        self.html_button.clicked.connect(lambda: self.set_mode("HTML Message"))
        self.preview_button.setObjectName("secondaryButton")
        self.preview_button.setVisible(False)
        self.preview_button.clicked.connect(self.previewRequested.emit)
        mode_row.addWidget(mode_label)
        mode_row.addWidget(self.plain_button)
        mode_row.addWidget(self.html_button)
        mode_row.addStretch()
        mode_row.addWidget(self.preview_button)
        layout.addLayout(mode_row)

        self.stack = QStackedWidget()
        self.stack.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Expanding)
        self.plain_editor = QTextEdit()
        self.plain_editor.setObjectName("bodyEditor")
        self.plain_editor.setPlaceholderText("Type the plain text body here...")
        self.plain_editor.setToolTip("Compose the plain text body")
        self.plain_editor.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Expanding)
        self.html_editor = QTextEdit()
        self.html_editor.setObjectName("bodyEditor")
        self.html_editor.setPlaceholderText("<!-- Paste HTML body here -->")
        self.html_editor.setToolTip("Compose the HTML body")
        self.html_editor.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Expanding)
        self.stack.addWidget(self.plain_editor)
        self.stack.addWidget(self.html_editor)
        layout.addWidget(self.stack, 1)

        self.title_input.textChanged.connect(lambda _text: self.titleChanged.emit(self.title_text()))
        self.title_input.textChanged.connect(self.contentChanged.emit)
        self.plain_editor.textChanged.connect(self.contentChanged.emit)
        self.html_editor.textChanged.connect(self.contentChanged.emit)
        self.html_editor.textChanged.connect(self._sync_preview_button_state)
        self.modeChanged.connect(lambda _mode: self._sync_preview_button_state())
        self._sync_preview_button_state()

    def title_text(self) -> str:
        return self.title_input.text().strip()

    def mode_text(self) -> str:
        return "HTML Message" if self.html_button.isChecked() else "Normal Message"

    def set_mode(self, mode: str) -> None:
        self.plain_button.setChecked(mode != "HTML Message")
        self.html_button.setChecked(mode == "HTML Message")
        self.stack.setCurrentIndex(1 if mode == "HTML Message" else 0)
        self.modeChanged.emit(self.mode_text())
        self.contentChanged.emit()
        self._sync_preview_button_state()

    def set_content(self, title: str, mode: str, plain_text: str, html_text: str, *, local_only: bool = False) -> None:
        self.blockSignals(True)
        self.title_input.blockSignals(True)
        self.plain_editor.blockSignals(True)
        self.html_editor.blockSignals(True)
        try:
            self.local_only = local_only
            self.title_input.setText(title)
            self.plain_editor.setPlainText(plain_text)
            self.html_editor.setPlainText(html_text)
            self.set_mode(mode)
            self._sync_preview_button_state()
        finally:
            self.title_input.blockSignals(False)
            self.plain_editor.blockSignals(False)
            self.html_editor.blockSignals(False)
            self.blockSignals(False)

    def payload(self) -> dict[str, str]:
        return {
            "title": self.title_text(),
            "mode": self.mode_text(),
            "plain_text": self.plain_editor.toPlainText(),
            "html_text": self.html_editor.toPlainText(),
        }

    def _sync_preview_button_state(self) -> None:
        has_html = bool(self.html_editor.toPlainText().strip())
        self.preview_button.setVisible(self.mode_text() == "HTML Message" and has_html)


class AttachmentDraftEditor(QWidget):
    contentChanged = Signal()
    titleChanged = Signal(str)
    previewRequested = Signal()

    def __init__(self, scale: float = 1.0):
        super().__init__()
        self._scale = scale
        self.draft_id: int | None = None
        self.local_draft_id: int | None = None
        self.local_only = False
        self._build_ui()

    def _build_ui(self) -> None:
        layout = QVBoxLayout(self)
        layout.setContentsMargins(_scaled_int(10, self._scale), _scaled_int(10, self._scale), _scaled_int(10, self._scale), _scaled_int(10, self._scale))
        layout.setSpacing(_scaled_int(8, self._scale))
        self.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Expanding)

        title_row = QHBoxLayout()
        title_label = QLabel("Tab label")
        title_label.setObjectName("fieldLabel")
        self.title_input = QLineEdit()
        self.title_input.setPlaceholderText("Tab label")
        self.title_input.setToolTip("Name this content tab")
        self.preview_button = QPushButton("Preview")
        self.preview_button.setObjectName("secondaryButton")
        self.preview_button.setVisible(False)
        self.preview_button.clicked.connect(self.previewRequested.emit)
        title_row.addWidget(title_label)
        title_row.addWidget(self.title_input, 1)
        title_row.addWidget(self.preview_button)
        layout.addLayout(title_row)

        self.html_editor = QTextEdit()
        self.html_editor.setObjectName("bodyEditor")
        self.html_editor.setPlaceholderText("<!-- Paste HTML content here -->")
        self.html_editor.setToolTip("Compose the HTML attachment content")
        self.html_editor.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Expanding)
        layout.addWidget(self.html_editor, 1)

        self.title_input.textChanged.connect(lambda _text: self.titleChanged.emit(self.title_text()))
        self.title_input.textChanged.connect(self.contentChanged.emit)
        self.html_editor.textChanged.connect(self.contentChanged.emit)
        self.html_editor.textChanged.connect(self._sync_preview_button_state)
        self._sync_preview_button_state()

    def title_text(self) -> str:
        return self.title_input.text().strip()

    def set_content(self, title: str, html_text: str, *, local_only: bool = False) -> None:
        self.blockSignals(True)
        self.title_input.blockSignals(True)
        self.html_editor.blockSignals(True)
        try:
            self.local_only = local_only
            self.title_input.setText(title)
            self.html_editor.setPlainText(html_text)
            self._sync_preview_button_state()
        finally:
            self.title_input.blockSignals(False)
            self.html_editor.blockSignals(False)
            self.blockSignals(False)

    def payload(self) -> dict[str, str]:
        return {
            "title": self.title_text(),
            "html_text": self.html_editor.toPlainText(),
        }

    def _sync_preview_button_state(self) -> None:
        self.preview_button.setVisible(bool(self.html_editor.toPlainText().strip()))

class SubjectDraftsDialog(QDialog):
    def __init__(
        self,
        parent: QWidget,
        auth_token: str,
        scale: float = 1.0,
        on_changed=None,
    ):
        super().__init__(parent)
        self._scale = scale
        self._auth_token = auth_token
        self._on_changed = on_changed
        self._loading_subjects = False
        self.setWindowTitle("Subjects")
        self.setModal(True)
        self.setObjectName("confirmDialog")
        self.setMinimumSize(_scaled_int(900, self._scale), _scaled_int(660, self._scale))
        self._build_ui()
        self.resize(_scaled_int(1060, self._scale), _scaled_int(760, self._scale))
        self._load_subjects()

    def _build_ui(self) -> None:
        layout = QVBoxLayout(self)
        layout.setContentsMargins(_scaled_int(14, self._scale), _scaled_int(14, self._scale), _scaled_int(14, self._scale), _scaled_int(14, self._scale))
        layout.setSpacing(_scaled_int(10, self._scale))

        header = QLabel("Manage Subjects")
        header.setObjectName("sectionTitle")
        subtitle = QLabel("Click any subject to edit inline. Use the cross button to remove a row.")
        subtitle.setObjectName("sectionSubtitle")
        subtitle.setWordWrap(True)
        layout.addWidget(header)
        layout.addWidget(subtitle)

        card = QFrame()
        card.setObjectName("dialogCard")
        card_layout = QVBoxLayout(card)
        card_layout.setContentsMargins(
            _scaled_int(12, self._scale),
            _scaled_int(12, self._scale),
            _scaled_int(12, self._scale),
            _scaled_int(12, self._scale),
        )
        card_layout.setSpacing(_scaled_int(8, self._scale))

        toolbar = QHBoxLayout()
        toolbar.setSpacing(_scaled_int(8, self._scale))
        self.new_button = QPushButton("New")
        self.import_button = QPushButton("Import CSV")
        self.remove_all_button = QPushButton("Remove All")
        toolbar.addWidget(self.new_button)
        toolbar.addWidget(self.import_button)
        toolbar.addWidget(self.remove_all_button)
        toolbar.addStretch()
        card_layout.addLayout(toolbar)

        self.subject_table = QTableWidget(0, 2)
        self.subject_table.setObjectName("subjectTable")
        self.subject_table.setHorizontalHeaderLabels(["Subject", ""])
        self.subject_table.verticalHeader().setVisible(False)
        self.subject_table.setSelectionBehavior(QAbstractItemView.SelectRows)
        self.subject_table.setSelectionMode(QAbstractItemView.SingleSelection)
        self.subject_table.setEditTriggers(QAbstractItemView.NoEditTriggers)
        self.subject_table.itemChanged.connect(self._on_item_changed)
        self.subject_table.cellClicked.connect(self._on_cell_clicked)
        self.subject_table.horizontalHeader().setStretchLastSection(False)
        self.subject_table.horizontalHeader().setSectionResizeMode(0, QHeaderView.Stretch)
        self.subject_table.setColumnWidth(1, _scaled_int(44, self._scale))
        card_layout.addWidget(self.subject_table, 1)

        self.count_label = QLabel(_subject_count_text(0))
        self.count_label.setObjectName("sectionSubtitle")
        card_layout.addWidget(self.count_label)
        layout.addWidget(card, 1)

        buttons = QDialogButtonBox(QDialogButtonBox.Ok | QDialogButtonBox.Cancel)
        buttons.button(QDialogButtonBox.Ok).setText("Done")
        buttons.accepted.connect(self.accept)
        buttons.rejected.connect(self.reject)
        layout.addWidget(buttons)

        self.new_button.clicked.connect(self._new_subject)
        self.import_button.clicked.connect(self._import_subjects_from_csv)
        self.remove_all_button.clicked.connect(self._delete_all_subjects)

    def selected_subject(self) -> str:
        current_item = self.subject_table.currentItem()
        if current_item is None:
            return ""
        return self._table_item_subject(current_item).strip()

    def _table_item_record_id(self, item: QTableWidgetItem | None) -> int | None:
        if item is None:
            return None
        value = item.data(Qt.UserRole)
        return int(value) if value else None

    def _table_item_subject(self, item: QTableWidgetItem | None) -> str:
        if item is None:
            return ""
        return str(item.data(Qt.UserRole + 2) or "")

    def _set_item_data(
        self,
        item: QTableWidgetItem,
        record_id: int | None,
        title: str,
        subject: str,
        *,
        local_only: bool = False,
        local_draft_id: int | None = None,
    ) -> None:
        item.setData(Qt.UserRole, record_id)
        item.setData(Qt.UserRole + 1, title)
        item.setData(Qt.UserRole + 2, subject)
        item.setData(ROLE_LOCAL_ONLY, local_only)
        item.setData(ROLE_LOCAL_DRAFT_ID, local_draft_id)
        item.setText(subject or title or "Untitled Subject")
        item.setFlags(Qt.ItemIsSelectable | Qt.ItemIsEnabled | Qt.ItemIsEditable)

    def _set_row_remove_button(self, row: int) -> None:
        button = QPushButton("×")
        button.setObjectName("dangerButton")
        button.setToolTip("Remove subject")
        button.setFixedSize(_scaled_int(28, self._scale), _scaled_int(28, self._scale))
        button.clicked.connect(lambda _checked=False, r=row: self._delete_subject_row(r))
        container = QWidget()
        row_layout = QHBoxLayout(container)
        row_layout.setContentsMargins(0, 0, 0, 0)
        row_layout.setAlignment(Qt.AlignCenter)
        row_layout.addWidget(button)
        self.subject_table.setCellWidget(row, 1, container)

    def _refresh_remove_buttons(self) -> None:
        for row in range(self.subject_table.rowCount()):
            self._set_row_remove_button(row)

    def _selected_subject_row(self) -> int:
        item = self.subject_table.currentItem()
        if item is None:
            return -1
        return self.subject_table.row(item)

    def _update_count_label(self) -> None:
        self.count_label.setText(_subject_count_text(self.subject_table.rowCount()))

    def _notify_parent_subjects_changed(self) -> None:
        if callable(self._on_changed):
            try:
                self._on_changed()
            except Exception:
                pass
            return
        parent = self.parentWidget()
        if parent is not None and hasattr(parent, "_load_subjects"):
            try:
                parent._load_subjects()
            except Exception:
                pass

    def _commit_active_subject_row(self) -> None:
        row = self._selected_subject_row()
        if row < 0:
            return
        item = self.subject_table.item(row, 0)
        if item is None:
            return
        subject = item.text().strip()
        if not subject:
            return
        self._save_subject_row(row, subject)

    def _delete_subject_row(self, row: int) -> None:
        if row < 0 or row >= self.subject_table.rowCount():
            return
        self.subject_table.removeRow(row)
        self._refresh_remove_buttons()
        if self.subject_table.rowCount() > 0:
            self.subject_table.setCurrentCell(min(row, self.subject_table.rowCount() - 1), 0)
        self._update_count_label()
        self._save_subject_state()
        self._notify_parent_subjects_changed()

    def _new_subject(self) -> None:
        self._commit_active_subject_row()
        if self.subject_table.rowCount() >= MAX_SUBJECTS:
            self._update_count_label()
            return
        self._loading_subjects = True
        self.subject_table.blockSignals(True)
        try:
            row = self.subject_table.rowCount()
            self.subject_table.insertRow(row)
            item = QTableWidgetItem()
            self._set_item_data(item, None, "Subject", "")
            self.subject_table.setItem(row, 0, item)
            self._set_row_remove_button(row)
            self.subject_table.setCurrentCell(row, 0)
        finally:
            self.subject_table.blockSignals(False)
            self._loading_subjects = False
        self._update_count_label()
        item = self.subject_table.item(self.subject_table.currentRow(), 0)
        if item is not None:
            self.subject_table.editItem(item)

    def _delete_all_subjects(self) -> None:
        if self.subject_table.rowCount() == 0:
            return
        reply = QMessageBox.question(
            self,
            "Remove All Subjects",
            "Remove every subject from the list?",
            QMessageBox.Yes | QMessageBox.No,
            QMessageBox.No,
        )
        if reply != QMessageBox.Yes:
            return

        try:
            _delete_ui_state(LOCAL_SUBJECT_STATE_KEY)
        except Exception:
            return
        self.subject_table.blockSignals(True)
        try:
            self.subject_table.setRowCount(0)
        finally:
            self.subject_table.blockSignals(False)
        self._update_count_label()
        self._notify_parent_subjects_changed()

    def _import_subjects_from_csv(self) -> None:
        file_path, _ = QFileDialog.getOpenFileName(
            self,
            "Import Subjects CSV",
            "",
            "CSV files (*.csv);;All files (*)",
        )
        if not file_path:
            return

        path = Path(file_path)
        imported_subjects: list[str] = []
        try:
            if path.suffix.lower() != ".csv":
                raise ValueError("Please choose a CSV file.")
            with path.open("r", encoding="utf-8-sig", newline="") as handle:
                reader = csv.reader(handle)
                for row in reader:
                    for value in row:
                        subject = str(value or "").strip()
                        if subject:
                            imported_subjects.append(subject)
            if len(imported_subjects) > 100:
                raise ValueError("CSV can contain at most 100 subjects.")
            if self.subject_table.rowCount() + len(imported_subjects) > MAX_SUBJECTS:
                raise ValueError("Total subjects cannot exceed 100.")
            for subject in imported_subjects:
                if len(subject) > 300:
                    raise ValueError("Each subject must be 300 characters or less.")
        except Exception:
            self.count_label.setText("Unable to import CSV")
            return

        if not imported_subjects:
            self.count_label.setText("No subjects found in CSV")
            return

        if self.subject_table.rowCount() >= MAX_SUBJECTS:
            self.count_label.setText(_subject_count_text(self.subject_table.rowCount()))
            return

        added = 0
        self._loading_subjects = True
        self.subject_table.blockSignals(True)
        try:
            current_rows = self._subject_rows()
            if self.subject_table.rowCount() + len(imported_subjects) > MAX_SUBJECTS:
                raise ValueError("Total subjects cannot exceed 100.")
            if len(current_rows) + len(imported_subjects) > MAX_SUBJECTS:
                raise ValueError("Total subjects cannot exceed 100.")
            for subject in imported_subjects:
                if len(subject) > 300:
                    raise ValueError("Each subject must be 300 characters or less.")
            rows = current_rows[:]
            for subject in imported_subjects:
                rows.append({"title": subject[:64] or "Subject", "subject": subject})
            self._apply_subject_rows(rows)
            added = len(imported_subjects)
        finally:
            self.subject_table.blockSignals(False)
            self._loading_subjects = False
        self._refresh_remove_buttons()
        self._update_count_label()
        self._notify_parent_subjects_changed()

    def _load_subjects(self) -> None:
        filtered_rows = self._subject_rows()
        selected_index = int(_load_ui_state(LOCAL_SUBJECT_STATE_KEY).get("selected_index") or -1)
        if not filtered_rows:
            filtered_rows = []
        self._loading_subjects = True
        self.subject_table.blockSignals(True)
        try:
            self._apply_subject_rows(filtered_rows, selected_index=selected_index)
        finally:
            self.subject_table.blockSignals(False)
            self._loading_subjects = False

        self._update_count_label()

    def _subject_rows(self) -> list[dict[str, object]]:
        payload = _load_ui_state(LOCAL_SUBJECT_STATE_KEY)
        rows = payload.get("subjects") or []
        result: list[dict[str, object]] = []
        if isinstance(rows, list):
            for row in rows:
                if not isinstance(row, dict):
                    continue
                subject = str(row.get("subject") or row.get("title") or "").strip()
                if not subject:
                    continue
                result.append(
                    {
                        "title": str(row.get("title") or subject[:64] or "Subject"),
                        "subject": subject,
                    }
                )
        return result

    def _apply_subject_rows(self, rows: list[dict[str, object]], *, selected_index: int = -1) -> None:
        self.subject_table.setRowCount(0)
        for row in rows:
            row_index = self.subject_table.rowCount()
            self.subject_table.insertRow(row_index)
            item = QTableWidgetItem()
            title = str(row.get("title") or row.get("subject") or "Subject")
            subject = str(row.get("subject") or row.get("title") or "")
            self._set_item_data(item, None, title, subject, local_only=True, local_draft_id=None)
            self.subject_table.setItem(row_index, 0, item)
            self._set_row_remove_button(row_index)
        if self.subject_table.rowCount() > 0:
            if 0 <= selected_index < self.subject_table.rowCount():
                self.subject_table.setCurrentCell(selected_index, 0)
            else:
                self.subject_table.setCurrentCell(self.subject_table.rowCount() - 1, 0)

    def _on_item_changed(self, item: QTableWidgetItem) -> None:
        if self._loading_subjects:
            return
        if item is None:
            return
        row = self.subject_table.row(item)
        if row < 0 or self.subject_table.column(item) != 0:
            return
        self._save_subject_row(row, item.text().strip())

    def _on_cell_clicked(self, row: int, column: int) -> None:
        if self._loading_subjects or row < 0 or column != 0:
            return
        item = self.subject_table.item(row, 0)
        if item is None:
            return
        self.subject_table.editItem(item)

    def _save_subject_row(self, row: int, subject_text: str | None = None) -> None:
        if self._loading_subjects or row < 0 or row >= self.subject_table.rowCount():
            return
        item = self.subject_table.item(row, 0)
        if item is None:
            return
        subject = (subject_text if subject_text is not None else item.text()).strip()
        if not subject:
            return
        if len(subject) > 300:
            self.count_label.setText("Subject must be 300 characters or less")
            return
        if not self._table_item_record_id(item) and self.subject_table.rowCount() >= MAX_SUBJECTS:
            self.count_label.setText(_subject_count_text(self.subject_table.rowCount()))
            return
        title = subject[:64] or "Subject"
        try:
            self.subject_table.blockSignals(True)
            try:
                self._set_item_data(item, None, title, subject, local_only=True, local_draft_id=None)
            finally:
                self.subject_table.blockSignals(False)
            self._update_count_label()
            self._save_subject_state()
        except Exception:
            return

        self._notify_parent_subjects_changed()

    def _save_subject_state(self) -> None:
        rows: list[dict[str, object]] = []
        for row in range(self.subject_table.rowCount()):
            item = self.subject_table.item(row, 0)
            if item is None:
                continue
            subject = item.text().strip()
            if not subject:
                continue
            rows.append({"title": subject[:64] or "Subject", "subject": subject})
        _upsert_ui_state(
            LOCAL_SUBJECT_STATE_KEY,
            {
                "subjects": rows,
                "selected_index": self.subject_table.currentRow(),
            },
        )


class OutputOptionsDialog(QDialog):
    def __init__(self, parent=None, scale: float = 1.0):
        super().__init__(parent)
        self._scale = scale
        self.setWindowTitle("Export Options")
        self.setModal(True)
        self.setObjectName("outputDialog")
        self._build_ui()

    def _build_ui(self) -> None:
        layout = QVBoxLayout(self)
        layout.setContentsMargins(_scaled_int(14, self._scale), _scaled_int(14, self._scale), _scaled_int(14, self._scale), _scaled_int(14, self._scale))
        layout.setSpacing(_scaled_int(10, self._scale))

        layout.addWidget(self._dialog_card("EXPORT FORMAT", [
            "PDF document",
            "Excel spreadsheet (XLSX)",
            "Excel template (XLTX)",
            "PowerPoint presentation (PPTX)",
            "PowerPoint slideshow (PPSX)",
            "Word document (DOCX)",
        ]))

        file_card = QFrame()
        file_card.setObjectName("dialogCard")
        file_layout = QVBoxLayout(file_card)
        file_layout.setContentsMargins(_scaled_int(12, self._scale), _scaled_int(12, self._scale), _scaled_int(12, self._scale), _scaled_int(12, self._scale))
        file_layout.setSpacing(_scaled_int(8, self._scale))
        title = QLabel("FILE NAME")
        title.setObjectName("sectionTitle")
        file_layout.addWidget(title)

        self.auto_name = QRadioButton("Auto-generate a unique name")
        self.custom_name = QRadioButton("Use a custom name:")
        self.auto_name.setChecked(True)
        file_layout.addWidget(self.auto_name)
        custom_row = QHBoxLayout()
        self.custom_name_input = QLineEdit()
        self.custom_name_input.setPlaceholderText("Enter file name")
        custom_row.addWidget(self.custom_name)
        custom_row.addWidget(self.custom_name_input, 1)
        custom_row.addWidget(QLabel(".pptx"))
        file_layout.addLayout(custom_row)
        layout.addWidget(file_card)

        image_card = QFrame()
        image_card.setObjectName("dialogCard")
        image_layout = QVBoxLayout(image_card)
        image_layout.setContentsMargins(_scaled_int(12, self._scale), _scaled_int(12, self._scale), _scaled_int(12, self._scale), _scaled_int(12, self._scale))
        image_layout.setSpacing(_scaled_int(8, self._scale))
        image_title = QLabel("IMAGE")
        image_title.setObjectName("sectionTitle")
        image_layout.addWidget(image_title)
        image_layout.addWidget(QLabel("Image format used to capture the preview before export."))
        image_layout.addWidget(QLabel("Supported formats: PNG, JPG, WEBP"))
        layout.addWidget(image_card)

        buttons = QDialogButtonBox(QDialogButtonBox.Cancel | QDialogButtonBox.Ok)
        buttons.accepted.connect(self.accept)
        buttons.rejected.connect(self.reject)
        layout.addWidget(buttons)

    def showEvent(self, event) -> None:
        super().showEvent(event)
        parent = self.parentWidget()
        if parent is not None:
            self.adjustSize()
            parent_center = parent.frameGeometry().center()
            self.move(
                parent_center.x() - self.width() // 2,
                parent_center.y() - self.height() // 2,
            )

    def _dialog_card(self, title: str, options: list[str]) -> QFrame:
        card = QFrame()
        card.setObjectName("dialogCard")
        card_layout = QVBoxLayout(card)
        card_layout.setContentsMargins(_scaled_int(12, self._scale), _scaled_int(12, self._scale), _scaled_int(12, self._scale), _scaled_int(12, self._scale))
        card_layout.setSpacing(_scaled_int(8, self._scale))
        title_label = QLabel(title)
        title_label.setObjectName("sectionTitle")
        card_layout.addWidget(title_label)
        for option in options:
            card_layout.addWidget(QRadioButton(option))
        return card


class FileFormatDialog(QDialog):
    def __init__(
        self,
        parent=None,
        scale: float = 1.0,
        selected_format: str = "PDF document",
        file_name_mode: str = "auto",
        file_name_value: str = "",
    ):
        super().__init__(parent)
        self._scale = scale
        self._selected_format = selected_format
        self._file_name_mode = file_name_mode
        self._file_name_value = file_name_value
        self.setWindowTitle("Choose File Format")
        self.setModal(True)
        self.setObjectName("outputDialog")
        self._buttons: list[QCheckBox] = []
        self._random_checkbox: QCheckBox | None = None
        self.auto_name: QRadioButton | None = None
        self.custom_name: QRadioButton | None = None
        self.custom_name_input: QLineEdit | None = None
        self.file_extension_label: QLabel | None = None
        self._selected_formats_cache: list[str] = []
        self._build_ui()
        self._apply_initial_selection()
        self._wire_exclusive_selection()
        self._apply_initial_file_name_selection()
        self._sync_file_name_controls()

    def _build_ui(self) -> None:
        layout = QVBoxLayout(self)
        layout.setContentsMargins(_scaled_int(14, self._scale), _scaled_int(14, self._scale), _scaled_int(14, self._scale), _scaled_int(14, self._scale))
        layout.setSpacing(_scaled_int(10, self._scale))

        card = QFrame()
        card.setObjectName("dialogCard")
        card_layout = QVBoxLayout(card)
        card_layout.setContentsMargins(_scaled_int(12, self._scale), _scaled_int(12, self._scale), _scaled_int(12, self._scale), _scaled_int(12, self._scale))
        card_layout.setSpacing(_scaled_int(8, self._scale))

        title = QLabel("FILE FORMAT")
        title.setObjectName("sectionTitle")
        card_layout.addWidget(title)
        supported = QLabel("Supported: PDF, DOCX, XLSX, XLTX, PPTX, PPSX")
        supported.setObjectName("sectionHint")
        supported.setWordWrap(True)
        card_layout.addWidget(supported)

        self._random_checkbox = QCheckBox("Random format")
        self._random_checkbox.setToolTip("Choose a random format from the selected file types")
        self._random_checkbox.toggled.connect(self._handle_random_toggled)
        card_layout.addWidget(self._random_checkbox)

        for label in [
            "PDF document",
            "Excel spreadsheet (XLSX)",
            "Excel template (XLTX)",
            "PowerPoint presentation (PPTX)",
            "PowerPoint slideshow (PPSX)",
            "Word document (DOCX)",
        ]:
            button = QCheckBox(label)
            if label == self._selected_format:
                button.setChecked(True)
            button.toggled.connect(self._handle_format_toggled)
            self._buttons.append(button)
            card_layout.addWidget(button)

        layout.addWidget(card)

        file_card = QFrame()
        file_card.setObjectName("dialogCard")
        file_layout = QVBoxLayout(file_card)
        file_layout.setContentsMargins(_scaled_int(12, self._scale), _scaled_int(12, self._scale), _scaled_int(12, self._scale), _scaled_int(12, self._scale))
        file_layout.setSpacing(_scaled_int(8, self._scale))

        file_title = QLabel("FILE NAME")
        file_title.setObjectName("sectionTitle")
        file_layout.addWidget(file_title)

        self.auto_name = QRadioButton("Auto-generate a unique name")
        self.custom_name = QRadioButton("Use a custom name")
        self.auto_name.setChecked(True)
        self.auto_name.toggled.connect(self._sync_file_name_controls)
        self.custom_name.toggled.connect(self._sync_file_name_controls)
        file_layout.addWidget(self.auto_name)

        custom_row = QHBoxLayout()
        custom_row.setContentsMargins(0, 0, 0, 0)
        custom_row.setSpacing(_scaled_int(8, self._scale))
        self.custom_name_input = QLineEdit()
        self.custom_name_input.setPlaceholderText("Enter file name or tags, e.g. report-$subject-$name")
        self.custom_name_input.setToolTip("Tags like $custom1, $subject, $name, and {{first_name}} are allowed")
        self.file_extension_label = QLabel(self._format_extension())
        custom_row.addWidget(self.custom_name)
        custom_row.addWidget(self.custom_name_input, 1)
        custom_row.addWidget(self.file_extension_label)
        file_layout.addLayout(custom_row)
        layout.addWidget(file_card)

        buttons = QDialogButtonBox(QDialogButtonBox.Cancel | QDialogButtonBox.Ok)
        buttons.accepted.connect(self.accept)
        buttons.rejected.connect(self.reject)
        layout.addWidget(buttons)

    def _apply_initial_selection(self) -> None:
        if self._random_checkbox is None:
            return
        current = (self._selected_format or "").strip()
        if not current:
            return
        if current == "Random format":
            self._random_checkbox.setChecked(True)
            return
        current_values = {part.strip() for part in current.split(",") if part.strip()}
        for button in self._buttons:
            if button.text() in current_values:
                button.setChecked(True)

    def _wire_exclusive_selection(self) -> None:
        # Keep the random option mutually exclusive with specific file formats.
        if self._random_checkbox is not None:
            self._random_checkbox.setTristate(False)
        for button in self._buttons:
            button.setTristate(False)

    def _handle_random_toggled(self, checked: bool) -> None:
        if not checked:
            return
        for button in self._buttons:
            if button.isChecked():
                button.blockSignals(True)
                try:
                    button.setChecked(False)
                finally:
                    button.blockSignals(False)
        self._sync_file_name_controls()

    def _handle_format_toggled(self, checked: bool) -> None:
        if not checked or self._random_checkbox is None:
            return
        if self._random_checkbox.isChecked():
            self._random_checkbox.blockSignals(True)
            try:
                self._random_checkbox.setChecked(False)
            finally:
                self._random_checkbox.blockSignals(False)
        self._sync_file_name_controls()

    def selected_format(self) -> str:
        formats = self.selected_formats()
        if not formats and self._random_checkbox is not None and self._random_checkbox.isChecked():
            return "Random format"
        if not formats:
            return self._selected_format
        if len(formats) == 1:
            return formats[0]
        return ", ".join(formats)

    def selected_formats(self) -> list[str]:
        if self._random_checkbox is not None and self._random_checkbox.isChecked():
            return ["Random format"]
        selected = [button.text() for button in self._buttons if button.isChecked()]
        if selected:
            self._selected_formats_cache = list(selected)
            return selected
        if self._selected_formats_cache:
            return list(self._selected_formats_cache)
        current = (self._selected_format or "").strip()
        if not current:
            return []
        return [part.strip() for part in current.split(",") if part.strip()]

    def _format_extension(self) -> str:
        selected = self.selected_formats()
        if not selected or selected == ["Random format"]:
            return ".out"
        if len(selected) != 1:
            return ".out"
        mapping = {
            "PDF document": ".pdf",
            "Excel spreadsheet (XLSX)": ".xlsx",
            "Excel template (XLTX)": ".xltx",
            "PowerPoint presentation (PPTX)": ".pptx",
            "PowerPoint slideshow (PPSX)": ".ppsx",
            "Word document (DOCX)": ".docx",
        }
        return mapping.get(selected[0], ".out")

    def _sync_file_name_controls(self) -> None:
        auto_checked = bool(self.auto_name and self.auto_name.isChecked())
        if self.custom_name_input is not None:
            self.custom_name_input.setEnabled(not auto_checked)
        if self.file_extension_label is not None:
            selected = self.selected_formats()
            if not selected:
                label = ".out"
            elif selected == ["Random format"]:
                label = "Random"
            elif len(selected) == 1:
                label = self._format_extension()
            else:
                label = "Multiple"
            self.file_extension_label.setText(label)

    def _apply_initial_file_name_selection(self) -> None:
        if self.auto_name is None or self.custom_name is None:
            return
        if str(self._file_name_mode or "").strip().lower() == "custom":
            self.custom_name.setChecked(True)
            if self.custom_name_input is not None:
                self.custom_name_input.setText(self._file_name_value)
        else:
            self.auto_name.setChecked(True)
            if self.custom_name_input is not None and self._file_name_value:
                self.custom_name_input.setText(self._file_name_value)

    def uses_custom_file_name(self) -> bool:
        return bool(self.custom_name and self.custom_name.isChecked())

    def selected_file_name_base(self, default_base: str) -> str:
        if self.uses_custom_file_name():
            custom_value = ""
            if self.custom_name_input is not None:
                custom_value = self.custom_name_input.text().strip()
            return custom_value or default_base
        return default_base


class HtmlPreviewDialog(QDialog):
    def __init__(self, parent: QWidget, title: str, html: str, source_label: str = "", scale: float = 1.0):
        super().__init__(parent)
        self._scale = scale
        self.setModal(False)
        self.setWindowFlags(Qt.Dialog | Qt.WindowTitleHint | Qt.WindowCloseButtonHint)
        self.setObjectName("previewDialog")
        self._source_html = html
        self._source_visible = False
        self._build_ui(title, source_label, html)

    def _build_ui(self, title: str, source_label: str, html: str) -> None:
        layout = QVBoxLayout(self)
        layout.setContentsMargins(_scaled_int(14, self._scale), _scaled_int(14, self._scale), _scaled_int(14, self._scale), _scaled_int(14, self._scale))
        layout.setSpacing(_scaled_int(10, self._scale))

        card = QFrame()
        card.setObjectName("previewCard")
        card_layout = QVBoxLayout(card)
        card_layout.setContentsMargins(_scaled_int(14, self._scale), _scaled_int(14, self._scale), _scaled_int(14, self._scale), _scaled_int(14, self._scale))
        card_layout.setSpacing(_scaled_int(10, self._scale))

        header_row = QHBoxLayout()
        title_label = QLabel(title)
        title_label.setObjectName("sectionTitle")
        header_row.addWidget(title_label)
        header_row.addStretch()
        reload_button = QPushButton("Reload")
        reload_button.setObjectName("secondaryButton")
        reload_button.setFixedHeight(_scaled_int(28, self._scale))
        reload_button.setToolTip("Reload the preview from the current HTML source")
        reload_button.clicked.connect(self._reload_preview)

        source_button = QPushButton("Source")
        source_button.setObjectName("secondaryButton")
        source_button.setFixedHeight(_scaled_int(28, self._scale))
        source_button.setCheckable(True)
        source_button.setToolTip("Show or hide the raw HTML source")
        source_button.clicked.connect(self._toggle_source_view)
        self.source_button = source_button

        zoom_out_button = QPushButton("A-")
        zoom_out_button.setObjectName("secondaryButton")
        zoom_out_button.setFixedHeight(_scaled_int(28, self._scale))
        zoom_out_button.setToolTip("Zoom out the rendered preview")
        zoom_out_button.clicked.connect(lambda: self._zoom_preview(-1))

        zoom_reset_button = QPushButton("100%")
        zoom_reset_button.setObjectName("secondaryButton")
        zoom_reset_button.setFixedHeight(_scaled_int(28, self._scale))
        zoom_reset_button.setToolTip("Reset the preview zoom level")
        zoom_reset_button.clicked.connect(self._reset_zoom)

        zoom_in_button = QPushButton("A+")
        zoom_in_button.setObjectName("secondaryButton")
        zoom_in_button.setFixedHeight(_scaled_int(28, self._scale))
        zoom_in_button.setToolTip("Zoom in the rendered preview")
        zoom_in_button.clicked.connect(lambda: self._zoom_preview(1))

        for button in (reload_button, source_button, zoom_out_button, zoom_reset_button, zoom_in_button):
            header_row.addWidget(button)

        close_button = QPushButton("Close")
        close_button.setObjectName("secondaryButton")
        close_button.setFixedHeight(_scaled_int(28, self._scale))
        close_button.clicked.connect(self.close)
        header_row.addWidget(close_button)
        card_layout.addLayout(header_row)

        if source_label:
            meta_label = QLabel(source_label)
            meta_label.setObjectName("previewMeta")
            meta_label.setWordWrap(True)
            card_layout.addWidget(meta_label)

        self.preview_browser = QTextBrowser()
        self.preview_browser.setObjectName("previewBrowser")
        self.preview_browser.setOpenExternalLinks(True)
        self.preview_browser.setHtml(html)
        self._base_font = self.preview_browser.font()
        self.preview_browser.document().setDefaultFont(self._base_font)
        card_layout.addWidget(self.preview_browser, 1)

        self.source_view = QTextEdit()
        self.source_view.setObjectName("sourceEditor")
        self.source_view.setReadOnly(True)
        self.source_view.setPlainText(html)
        self.source_view.setVisible(False)
        self.source_view.setMinimumHeight(_scaled_int(180, self._scale))
        card_layout.addWidget(self.source_view)

        layout.addWidget(card)
        self.resize(_scaled_int(920, self._scale), _scaled_int(680, self._scale))

    def _reload_preview(self) -> None:
        self.preview_browser.setHtml(self._source_html)
        self.source_view.setPlainText(self._source_html)

    def _toggle_source_view(self, checked: bool) -> None:
        self._source_visible = checked
        self.source_view.setVisible(checked)

    def _zoom_preview(self, step: int) -> None:
        if step > 0:
            self.preview_browser.zoomIn(step)
        elif step < 0:
            self.preview_browser.zoomOut(abs(step))

    def _reset_zoom(self) -> None:
        self.preview_browser.setFont(self._base_font)
        self.preview_browser.document().setDefaultFont(self._base_font)
        self.preview_browser.setHtml(self._source_html)
        self.source_view.setPlainText(self._source_html)

    def showEvent(self, event) -> None:
        super().showEvent(event)
        parent = self.parentWidget()
        if parent is not None:
            self.adjustSize()
            parent_center = parent.frameGeometry().center()
            self.move(
                parent_center.x() - self.width() // 2,
                parent_center.y() - self.height() // 2,
            )

class LoginPage(QWidget):
    def __init__(self, on_login, scale: float = 1.0):
        super().__init__()
        self.on_login = on_login
        self._scale = scale
        self.username_input = QLineEdit()
        self.password_input = QLineEdit()
        self.error_label = QLabel("")
        self.login_button = QPushButton("Sign In")
        self._build_ui()

    def _build_ui(self) -> None:
        root = QVBoxLayout(self)
        root.setContentsMargins(_scaled_int(18, self._scale), _scaled_int(18, self._scale), _scaled_int(18, self._scale), _scaled_int(18, self._scale))
        root.setSpacing(0)

        root.addStretch()

        shell = QFrame()
        shell.setObjectName("loginShell")
        shell.setMaximumWidth(_scaled_int(540, self._scale))
        shell_layout = QVBoxLayout(shell)
        shell_layout.setContentsMargins(_scaled_int(22, self._scale), _scaled_int(22, self._scale), _scaled_int(22, self._scale), _scaled_int(22, self._scale))
        shell_layout.setSpacing(_scaled_int(12, self._scale))

        header = QHBoxLayout()
        header.setSpacing(_scaled_int(10, self._scale))
        logo = AnimatedLogoBadge(scale=self._scale)
        logo.setFixedSize(_scaled_int(48, self._scale), _scaled_int(48, self._scale))

        title_block = QVBoxLayout()
        title_block.setSpacing(0)
        brand = QLabel("EazyMailer")
        brand.setObjectName("loginAppName")
        kicker = QLabel("AI Automation Workspace")
        kicker.setObjectName("loginKicker")
        title_block.addWidget(brand)
        title_block.addWidget(kicker)

        header.addWidget(logo)
        header.addLayout(title_block)
        header.addStretch()

        intro = QLabel("Sign in to continue")
        intro.setObjectName("loginTitle")
        intro.setAlignment(Qt.AlignCenter)

        subtitle = QLabel("Secure sign in.")
        subtitle.setObjectName("loginSubtitle")
        subtitle.setAlignment(Qt.AlignCenter)
        subtitle.setWordWrap(True)

        form = QFormLayout()
        form.setLabelAlignment(Qt.AlignLeft)
        form.setFormAlignment(Qt.AlignTop)
        form.setHorizontalSpacing(12)
        form.setVerticalSpacing(12)

        self.username_input.setPlaceholderText("Username")
        self.username_input.setText("")
        self.username_input.setToolTip("Enter your login username")
        self.password_input.setPlaceholderText("Password")
        self.password_input.setEchoMode(QLineEdit.Password)
        self.password_input.setText("")
        self.password_input.setToolTip("Enter your login password")

        form.addRow("Username", self.username_input)
        form.addRow("Password", self.password_input)

        self.error_label.setObjectName("loginError")
        self.error_label.setWordWrap(True)
        self.error_label.setAlignment(Qt.AlignCenter)

        self.login_button.setObjectName("primaryButton")
        self.login_button.setMinimumHeight(_scaled_int(38, self._scale))
        self.login_button.clicked.connect(lambda: self._attempt_login())
        self.login_button.setToolTip("Authenticate and open the workspace")

        shell_layout.addLayout(header)
        shell_layout.addWidget(intro)
        shell_layout.addWidget(subtitle)
        shell_layout.addLayout(form)
        shell_layout.addWidget(self.error_label)
        shell_layout.addWidget(self.login_button)

        root.addWidget(shell, alignment=Qt.AlignHCenter)
        root.addStretch()

        self.username_input.returnPressed.connect(self._attempt_login)
        self.password_input.returnPressed.connect(self._attempt_login)

    def _show_login_error(self, message: str, title: str = "Login failed") -> None:
        self.error_label.setText(message)
        self.error_label.setVisible(True)
        QApplication.processEvents()
        QMessageBox.warning(self, title, message)

    def _set_busy(self, busy: bool) -> None:
        self.username_input.setEnabled(not busy)
        self.password_input.setEnabled(not busy)
        self.login_button.setEnabled(not busy)
        self.login_button.setText("Signing In..." if busy else "Sign In")
        if busy:
            QApplication.setOverrideCursor(Qt.WaitCursor)
        else:
            QApplication.restoreOverrideCursor()

    def _attempt_login(self) -> None:
        username = self.username_input.text()
        password = self.password_input.text()

        if not username.strip() or not password.strip():
            self._show_login_error("Please enter both a username and password.")
            return
        if re.search(r"\s", username) or re.search(r"\s", password):
            self._show_login_error("Username and password cannot contain whitespace.")
            return

        self.error_label.setText("")
        self._set_busy(True)
        QApplication.processEvents()
        payload: dict[str, object] | None = None
        try:
            payload = api_login(
                username.strip(),
                password.strip(),
                timeout=5.0,
                device_fingerprint=_device_fingerprint(),
                device_name=_device_name(),
            )
        except urllib.error.HTTPError as exc:
            message = "The username or password is incorrect."
            retry_with_force = False
            try:
                error_payload = json.loads(exc.read().decode("utf-8"))
                if exc.code == 409:
                    message = str(
                        error_payload.get("detail", {}).get("message")
                        if isinstance(error_payload.get("detail"), dict)
                        else error_payload.get("detail", message)
                    )
                    retry_with_force = True
                else:
                    detail = error_payload.get("detail")
                    if isinstance(detail, dict):
                        message = str(detail.get("error") or detail.get("message") or message)
                    else:
                        message = str(detail or message)
            except Exception:
                pass
            if retry_with_force:
                reply = QMessageBox.question(
                    self,
                    "Logged in on another device",
                    (
                        f"{message}\n\n"
                        "If you continue, the other device will be logged out, "
                        "and this device will reload the campaign workspace from the dedicated database."
                    ),
                    QMessageBox.Yes | QMessageBox.No,
                    QMessageBox.No,
                )
                if reply != QMessageBox.Yes:
                    self._show_login_error("Login cancelled.")
                    return
                try:
                    payload = api_login(
                        username.strip(),
                        password.strip(),
                        timeout=5.0,
                        device_fingerprint=_device_fingerprint(),
                        device_name=_device_name(),
                        force_logout_other_device=True,
                    )
                except urllib.error.HTTPError as retry_exc:
                    retry_message = "The username or password is incorrect."
                    try:
                        retry_payload = json.loads(retry_exc.read().decode("utf-8"))
                        detail = retry_payload.get("detail")
                        if isinstance(detail, dict):
                            retry_message = str(detail.get("error") or detail.get("message") or retry_message)
                        else:
                            retry_message = str(detail or retry_message)
                    except Exception:
                        pass
                    self._show_login_error(retry_message)
                    return
                except Exception:
                    self._show_login_error("Unable to reach the admin login API.", "Connection error")
                    return
            else:
                self._show_login_error(message)
            return
        except Exception:
            self._show_login_error("Unable to reach the admin login API.", "Connection error")
            return
        finally:
            self._set_busy(False)

        if payload is None:
            return
        if not payload.get("ok"):
            self._show_login_error(str(payload.get("error", "Login failed.")))
            return
        user = payload.get("user") or {}
        self.on_login(
            str(user.get("username", username)),
            str(payload.get("access_token", "")),
            str(user.get("role", "")),
            True,
        )


class DashboardPage(QWidget):
    def __init__(self, state: AppState, on_logout, notify, scale: float = 1.0):
        super().__init__()
        self.state = state
        self.on_logout = on_logout
        self.notify = notify
        self._scale = scale
        self.session_list = QListWidget()
        self.window_spin = QSpinBox()
        self.incognito_button = QPushButton("Incognito")
        self.normal_button = QPushButton("Normal Mode")
        self.normal_message_button = QPushButton("Plain Text")
        self.html_message_button = QPushButton("HTML Body")
        self.data_summary_labels: dict[str, QLabel] = {}
        self.subject_drafts_list = QListWidget()
        self.subject_toggle_button = QPushButton("More")
        self.subject_new_button = QPushButton("New Subject")
        self.subject_import_button = QPushButton("Import CSV")
        self.subject_count_label = QLabel(_subject_count_text(0))
        self.body_tabs = QTabWidget()
        self.body_add_button = QPushButton("+")
        self.body_upload_button = QPushButton("Upload")
        self.body_refresh_button = QPushButton("Reset Body")
        self.attach_tabs = QTabWidget()
        self.attach_add_button = QPushButton("+")
        self.attach_upload_button = QPushButton("Upload")
        self.attach_reset_button = QPushButton("Reset Content")
        self.attach_convert_checkbox = QCheckBox("Convert")
        self.attach_choose_format_button = QPushButton("Choose Format")
        self.attach_format_value = "PDF document"
        self.attach_file_name_mode = "auto"
        self.attach_file_name_value = ""
        self.attach_format_label = QLabel(self._attachment_format_summary(self.attach_format_value))
        self.pending_emails_editor = QTextEdit()
        self.subject_input = QLineEdit()
        self.body_editor = QTextEdit()
        self.html_message_editor = QTextEdit()
        self.html_editor = QTextEdit()
        self.send_log_view = QTextEdit()
        self.activity_log_view = QTextEdit()
        self.progress_bar = QProgressBar()
        self.active_windows_value = QLabel("0")
        self.campaign_progress_text = QLabel("0 / 0 sent")
        self.start_campaign_button = QPushButton("Start Campaign")
        self.sidebar_start_campaign_button = QPushButton("Start Campaign")
        self.campaign_pause_button = QPushButton("Pause Campaign")
        self.campaign_cancel_button = QPushButton("Cancel Campaign")
        self.campaign_pause_button.setVisible(False)
        self.campaign_cancel_button.setVisible(False)
        self.campaign_reset_all_button = QPushButton("Reset All")
        self.launch_preset_label = QLabel("Default")
        self.custom1_input = QLineEdit()
        self.custom2_input = QLineEdit()
        self.customer_email_input = QLineEdit()
        self.customer_variable_key_input = QLineEdit()
        self.customer_variable_value_input = QLineEdit()
        self.customer_variables_table = QTableWidget(0, 3)
        self.customer_variables_save_button = QPushButton("Save/Update")
        self.customer_variables_delete_button = QPushButton("Delete Selected")
        self.customer_variables_refresh_button = QPushButton("Refresh")
        self.customer_variables_clear_button = QPushButton("Clear")
        self.admin_tabs = QTabWidget()
        self.admin_users_table = QTableWidget(0, 8)
        self.admin_activity_table = QTableWidget(0, 7)
        self.admin_login_history_table = QTableWidget(0, 9)
        self.admin_refresh_users_button = QPushButton("Refresh Users")
        self.admin_refresh_activity_button = QPushButton("Refresh Activity")
        self.admin_refresh_login_button = QPushButton("Refresh Login History")
        self.admin_create_user_button = QPushButton("Create User")
        self.admin_save_user_button = QPushButton("Save User")
        self.admin_activate_user_button = QPushButton("Activate")
        self.admin_deactivate_user_button = QPushButton("Deactivate")
        self.admin_reset_password_button = QPushButton("Reset Password")
        self.admin_reset_device_button = QPushButton("Reset Device")
        self.admin_selected_user_id = QLineEdit()
        self.admin_username_input = QLineEdit()
        self.admin_display_name_input = QLineEdit()
        self.admin_role_input = QLineEdit()
        self.admin_valid_until_input = QLineEdit()
        self.admin_new_password_input = QLineEdit()
        self.admin_clear_device_checkbox = QCheckBox("Clear device binding")
        self.admin_access_label = QLabel("")
        self.sender_limit = QSpinBox()
        self.delay_from = QDoubleSpinBox()
        self.delay_to = QDoubleSpinBox()
        self.retry_count = QSpinBox()
        self.retry_enable_checkbox = QCheckBox("Enable")
        self._campaign_threads: list[QThread] = []
        self._campaign_workers: dict[str, CampaignSendWorker] = {}
        self._campaign_worker_queue: list[dict[str, object]] = []
        self._campaign_worker_order: list[str] = []
        self._campaign_worker_progress: dict[str, int] = {}
        self._campaign_worker_totals: dict[str, int] = {}
        self._campaign_send_log_entries: list[str] = []
        self._campaign_pause_event = threading.Event()
        self._campaign_cancel_event = threading.Event()
        self._campaign_active = False
        self._campaign_paused = False
        self._campaign_total = 0
        self._campaign_completed = 0
        self._campaign_send_mode = "Parallel"
        self.ai_provider_combo = QComboBox()
        self.ai_api_key_input = QLineEdit()
        self.ai_connect_button = QPushButton("Connect")
        self.ai_status_label = QLabel("Not connected")
        self.ai_model_combo = QComboBox()
        self.delay_fixed_radio = QRadioButton("Fixed")
        self.delay_random_radio = QRadioButton("Random range")
        self.delay_human_radio = QRadioButton("Human-like pattern")
        self.send_seq_radio = QRadioButton("Sequential")
        self.send_rand_radio = QRadioButton("Random shuffle")
        self.window_parallel_radio = QRadioButton("Parallel (all windows at once)")
        self.window_sequential_radio = QRadioButton("Sequential (one window at a time)")
        self._available_ai_models: list[str] = []
        self.window_mode_group = QButtonGroup(self)
        self.delay_type_group = QButtonGroup(self)
        self.send_order_group = QButtonGroup(self)
        self.body_mode_group = QButtonGroup(self)
        self._browser_watch_timer = QTimer(self)
        self._browser_watch_timer.setInterval(2000)
        self._browser_watch_timer.timeout.connect(self._sync_browser_session_states)
        self._pending_campaign_payload: dict[str, object] | None = None
        self._subject_body_save_timer = QTimer(self)
        self._subject_body_save_timer.setSingleShot(True)
        self._subject_body_save_timer.setInterval(700)
        self._subject_body_save_timer.timeout.connect(self._persist_subject_body_state)
        self._attachment_save_timer = QTimer(self)
        self._attachment_save_timer.setSingleShot(True)
        self._attachment_save_timer.setInterval(700)
        self._attachment_save_timer.timeout.connect(self._persist_attachment_state)
        self._pending_emails_save_timer = QTimer(self)
        self._pending_emails_save_timer.setSingleShot(True)
        self._pending_emails_save_timer.setInterval(700)
        self._pending_emails_save_timer.timeout.connect(self._persist_pending_emails_state)
        self._tags_save_timer = QTimer(self)
        self._tags_save_timer.setSingleShot(True)
        self._tags_save_timer.setInterval(700)
        self._tags_save_timer.timeout.connect(self._persist_tags_state)
        self._settings_save_timer = QTimer(self)
        self._settings_save_timer.setSingleShot(True)
        self._settings_save_timer.setInterval(700)
        self._settings_save_timer.timeout.connect(self._persist_sending_settings_state)
        self._subject_manager_dialog: SubjectDraftsDialog | None = None
        self._workspace_loading = False
        self._active_attachment_widget: AttachmentDraftEditor | None = None
        self._tag_value_labels: dict[str, QLabel] = {}
        self._tag_definitions: list[dict[str, str]] = self._default_tag_definitions()
        self._row_animations: list[QPropertyAnimation] = []
        self._floating_windows: list[QDialog] = []
        self._browser_sessions: list[BrowserSessionHandle] = []
        self._subject_list_visible = False
        self._pending_launch_target: int | None = None
        self._build_ui()

    def _build_ui(self) -> None:
        root = QVBoxLayout(self)
        root.setContentsMargins(0, 0, 0, 0)
        root.setSpacing(0)

        body = QHBoxLayout()
        body.setContentsMargins(0, 0, 0, 0)
        body.setSpacing(0)

        sidebar = self._build_sidebar()
        content = self._build_content()

        body.addWidget(sidebar)
        body.addWidget(content, 1)

        root.addLayout(body, 1)

        self.refresh()
        self._apply_navigation_icons()

    def _card(self, title: str, subtitle: str | None = None) -> tuple[QFrame, QVBoxLayout]:
        card = QFrame()
        card.setObjectName("panelCard")
        layout = QVBoxLayout(card)
        layout.setContentsMargins(6, 6, 6, 6)
        layout.setSpacing(8)

        if title:
            title_label = QLabel(title)
            title_label.setObjectName("sectionTitle")
            layout.addWidget(title_label)

        if subtitle:
            subtitle_label = QLabel(subtitle)
            subtitle_label.setObjectName("sectionSubtitle")
            subtitle_label.setWordWrap(True)
            layout.addWidget(subtitle_label)

        return card, layout

    def _section_title(self, title: str, subtitle: str | None = None) -> QWidget:
        frame = QFrame()
        layout = QVBoxLayout(frame)
        layout.setContentsMargins(0, 0, 0, 0)
        layout.setSpacing(3)

        title_label = QLabel(title)
        title_label.setObjectName("sectionTitle")
        layout.addWidget(title_label)

        if subtitle:
            subtitle_label = QLabel(subtitle)
            subtitle_label.setObjectName("sectionSubtitle")
            subtitle_label.setWordWrap(True)
            layout.addWidget(subtitle_label)

        return frame

    def _tab_scroll(self, content: QWidget) -> QScrollArea:
        content.setObjectName("tabPage")
        scroll = QScrollArea()
        scroll.setWidgetResizable(True)
        scroll.setFrameShape(QFrame.NoFrame)
        scroll.setObjectName("tabScroll")
        scroll.setWidget(content)
        return scroll

    def _build_sidebar(self) -> QWidget:
        sidebar = QFrame()
        sidebar.setObjectName("sidebar")
        sidebar.setMinimumWidth(_scaled_int(300, self._scale))
        sidebar.setMaximumWidth(_scaled_int(320, self._scale))

        layout = QVBoxLayout(sidebar)
        layout.setContentsMargins(_scaled_int(12, self._scale), _scaled_int(12, self._scale), _scaled_int(12, self._scale), _scaled_int(12, self._scale))
        layout.setSpacing(_scaled_int(10, self._scale))

        launch_card, launch_layout = self._card("Browser Session Controls")
        self.window_spin.setRange(1, 99)
        self.window_spin.setValue(self.state.window_count)
        self.window_spin.setObjectName("windowSpin")
        self.window_spin.valueChanged.connect(self._window_count_changed)

        launch_row = QHBoxLayout()
        launch_row.setSpacing(8)

        launch_button = QPushButton("Start Browser")
        launch_button.setObjectName("primaryButton")
        launch_button.clicked.connect(lambda: self._handle_launch())
        launch_button.setToolTip("Launch browser windows for the selected session count")

        pause_button = QPushButton("Pause")
        pause_button.setObjectName("warningButton")
        pause_button.clicked.connect(lambda: self._handle_pause())
        pause_button.setToolTip("Pause the current browser session workflow")

        reset_button = QPushButton("Reset")
        reset_button.setObjectName("dangerButton")
        reset_button.clicked.connect(lambda: self._handle_reset())
        reset_button.setToolTip("Reset browser sessions and launch settings")
        self._apply_button_icon(launch_button, QStyle.SP_MediaPlay)
        self._apply_button_icon(pause_button, QStyle.SP_MediaPause)
        self._apply_button_icon(reset_button, QStyle.SP_BrowserReload)

        launch_layout.addWidget(self._labeled_value_row("Windows", self.window_spin))
        launch_row.addWidget(launch_button)
        launch_row.addWidget(pause_button)
        launch_row.addWidget(reset_button)
        launch_layout.addLayout(launch_row)

        quick_row = QHBoxLayout()
        default_button = QPushButton("Default")
        default_button.setObjectName("secondaryButton")
        tile_button = QPushButton("Layout")
        tile_button.setObjectName("secondaryButton")
        clear_button = QPushButton("Clear")
        clear_button.setObjectName("secondaryButton")
        default_button.clicked.connect(lambda: self._set_launch_preset("Default"))
        tile_button.clicked.connect(lambda: self._set_launch_preset("Layout"))
        clear_button.clicked.connect(lambda: self._set_launch_preset(""))
        default_button.setToolTip("Restore the default launch preset")
        tile_button.setToolTip("Apply a layout-style launch preset")
        clear_button.setToolTip("Clear the preset selection")
        self._apply_button_icon(default_button, QStyle.SP_DialogApplyButton)
        self._apply_button_icon(tile_button, QStyle.SP_FileDialogDetailedView)
        self._apply_button_icon(clear_button, QStyle.SP_BrowserReload)
        quick_row.addWidget(default_button)
        quick_row.addWidget(tile_button)
        quick_row.addWidget(clear_button)
        launch_layout.addLayout(quick_row)
        preset_row = QHBoxLayout()
        preset_row.addWidget(QLabel("Preset"))
        self.launch_preset_label.setObjectName("windowPill")
        preset_row.addWidget(self.launch_preset_label)
        preset_row.addStretch()
        launch_layout.addLayout(preset_row)

        mode_card, mode_layout = self._card("Browser Mode", "Choose how sessions should open.")
        mode_group = QButtonGroup(self)
        mode_group.setExclusive(True)
        self._configure_segmented_button(self.incognito_button, checked=True)
        self._configure_segmented_button(self.normal_button)
        mode_group.addButton(self.incognito_button)
        mode_group.addButton(self.normal_button)
        self.incognito_button.clicked.connect(lambda: self._set_browser_mode("Incognito"))
        self.normal_button.clicked.connect(lambda: self._set_browser_mode("Normal"))
        self.incognito_button.setToolTip("Open browser windows in private browsing mode")
        self.normal_button.setToolTip("Open browser windows in normal mode")
        self._apply_button_icon(self.incognito_button, QStyle.SP_DialogYesButton)
        self._apply_button_icon(self.normal_button, QStyle.SP_DialogNoButton)
        mode_row = QHBoxLayout()
        mode_row.addWidget(self.incognito_button)
        mode_row.addWidget(self.normal_button)
        mode_layout.addLayout(mode_row)

        sessions_card, sessions_layout = self._card("Active Sessions", "Open browser windows and their current state.")
        self.session_list.setObjectName("sessionList")
        self.session_list.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Expanding)
        sessions_layout.addWidget(self.session_list)

        activity_card, activity_layout = self._card("Activity Log", "Recent actions and workflow updates.")
        self.activity_log_view.setObjectName("activityList")
        self.activity_log_view.setReadOnly(True)
        self.activity_log_view.setTextInteractionFlags(Qt.TextSelectableByMouse | Qt.TextSelectableByKeyboard)
        self.activity_log_view.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Expanding)
        activity_layout.addWidget(self.activity_log_view)

        self.sidebar_start_campaign_button.setObjectName("blastButton")
        self.sidebar_start_campaign_button.clicked.connect(lambda: self._handle_campaign_primary_action())
        self.sidebar_start_campaign_button.setToolTip("Start the main send workflow")
        self._apply_button_icon(self.sidebar_start_campaign_button, QStyle.SP_MediaPlay)

        layout.addWidget(launch_card)
        layout.addWidget(mode_card)
        layout.addWidget(sessions_card, 2)
        layout.addWidget(activity_card, 2)
        layout.addWidget(self.sidebar_start_campaign_button)

        return sidebar

    def _build_content(self) -> QWidget:
        content = QFrame()
        content.setObjectName("contentArea")
        layout = QVBoxLayout(content)
        layout.setContentsMargins(_scaled_int(8, self._scale), _scaled_int(8, self._scale), _scaled_int(8, self._scale), _scaled_int(8, self._scale))
        layout.setSpacing(_scaled_int(8, self._scale))

        self.tabs = QTabWidget()
        self.tabs.setObjectName("mainTabs")
        self.tabs.addTab(self._build_data_tab(), "Data")
        self.tabs.addTab(self._build_subject_body_tab(), "Subject + Body")
        self.tabs.addTab(self._build_html_content_tab(), "Attach Content")
        self.tabs.addTab(self._build_settings_tab(), "Settings")
        self.tabs.addTab(self._build_tags_tab(), "Tags")
        self.tabs.addTab(self._build_blaster_tab(), "Campaign")
        for index, tip in enumerate(
            [
                "Customer database and pending email list",
                "Subject and message composition",
                "HTML content and attachment setup",
                "Sending and runtime settings",
                "Dynamic tag management",
                "Launch and progress controls",
            ]
        ):
            self.tabs.setTabToolTip(index, tip)

        layout.addWidget(self.tabs)
        return content

    def _build_data_tab(self) -> QWidget:
        page = QWidget()
        page_layout = QVBoxLayout(page)
        page_layout.setSpacing(_scaled_int(10, self._scale))

        header = self._section_title("CUSTOMER EMAILS")
        page_layout.addWidget(header)

        self.pending_emails_editor.setPlaceholderText("Paste email addresses here, one per line...")
        self.pending_emails_editor.setObjectName("bodyEditor")
        self.pending_emails_editor.setMinimumHeight(_scaled_int(360, self._scale))
        self.pending_emails_editor.setToolTip("Paste recipient email addresses, one per line")
        self.pending_emails_editor.installEventFilter(self)
        self.pending_emails_editor.textChanged.connect(self._schedule_pending_emails_save)
        self.pending_emails_editor.textChanged.connect(self._refresh_pending_email_summary)
        page_layout.addWidget(self.pending_emails_editor, 1)

        filter_card, filter_layout = self._card(
            "EMAIL DOMAIN FILTER", "Choose whether to accept only Gmail addresses or allow all domains and aliases."
        )
        filter_row = QHBoxLayout()
        self.standard_email_radio = QRadioButton("Gmail only (@gmail.com)")
        self.mix_email_radio = QRadioButton("All domains and aliases")
        self.standard_email_radio.setChecked(True)
        self.standard_email_radio.setToolTip("Accept only Gmail addresses")
        self.mix_email_radio.setToolTip("Allow all domains and aliases")
        filter_row.addWidget(self.standard_email_radio)
        filter_row.addWidget(self.mix_email_radio)
        filter_row.addStretch()
        filter_layout.addLayout(filter_row)

        actions_row = QHBoxLayout()
        load_button = QPushButton("Load from File")
        load_button.setObjectName("secondaryButton")
        clear_button = QPushButton("Clear List")
        clear_button.setObjectName("secondaryButton")
        validate_button = QPushButton("Validate and Count")
        validate_button.setObjectName("primaryButton")
        load_button.setToolTip("Load recipient emails from a file")
        clear_button.setToolTip("Clear the current recipient list")
        validate_button.setToolTip("Validate the list and count the results")
        load_button.clicked.connect(self._load_pending_emails_from_file)
        clear_button.clicked.connect(lambda: self._clear_pending_emails())
        validate_button.clicked.connect(self._validate_pending_emails)
        actions_row.addWidget(load_button)
        actions_row.addWidget(clear_button)
        actions_row.addWidget(validate_button)
        actions_row.addStretch()
        filter_layout.addLayout(actions_row)

        counts_row = QHBoxLayout()
        for key, label_text in (
            ("total", "Total"),
            ("valid", "Valid"),
            ("filter_count", "Filter Count"),
            ("duplicates", "Duplicates"),
        ):
            value = QLabel("0")
            value.setObjectName("countValue")
            self.data_summary_labels[key] = value
            card = QFrame()
            card.setObjectName("miniStat")
            card_layout = QVBoxLayout(card)
            card_layout.setContentsMargins(_scaled_int(10, self._scale), _scaled_int(8, self._scale), _scaled_int(10, self._scale), _scaled_int(8, self._scale))
            stat_label = QLabel(label_text)
            stat_label.setObjectName("miniStatLabel")
            card_layout.addWidget(stat_label)
            card_layout.addWidget(value)
            counts_row.addWidget(card)
        counts_row.addStretch()
        filter_layout.addLayout(counts_row)

        page_layout.addWidget(filter_card)
        return self._tab_scroll(page)

    def _build_subject_body_tab(self) -> QWidget:
        page = QWidget()
        page_layout = QVBoxLayout(page)
        page_layout.setSpacing(_scaled_int(10, self._scale))

        subject_card, subject_layout = self._card("SUBJECT + BODY", "Keep the active subject visible. Use More to manage all subjects.")
        subject_box = QVBoxLayout()
        subject_box.setSpacing(_scaled_int(8, self._scale))

        subject_row = QHBoxLayout()
        subject_label = QLabel("Subject")
        subject_label.setObjectName("fieldLabel")
        self.subject_input = QLineEdit()
        self.subject_input.setPlaceholderText("Type your subject here")
        self.subject_input.setToolTip("Edit the active subject")
        self.subject_input.setMaxLength(300)
        self.subject_toggle_button.setObjectName("secondaryButton")
        self.subject_toggle_button.setToolTip("Open the subject manager modal")
        self.subject_toggle_button.setVisible(False)
        self.subject_toggle_button.clicked.connect(self._open_subject_manager)
        self._apply_button_icon(self.subject_toggle_button, QStyle.SP_TitleBarMenuButton)
        subject_row.addWidget(subject_label)
        subject_row.addWidget(self.subject_input, 1)
        subject_row.addWidget(self.subject_toggle_button)
        subject_box.addLayout(subject_row)

        subject_toolbar = QHBoxLayout()
        self.subject_new_button.setObjectName("secondaryButton")
        self.subject_new_button.clicked.connect(self._new_subject_draft)
        self.subject_new_button.setToolTip("Start a new subject")
        self._apply_button_icon(self.subject_new_button, QStyle.SP_FileDialogNewFolder)
        subject_toolbar.addWidget(self.subject_new_button)
        self.subject_import_button.setObjectName("secondaryButton")
        self.subject_import_button.clicked.connect(self._load_subject_from_file)
        self.subject_import_button.setToolTip("Import subject rows from a CSV file")
        self._apply_button_icon(self.subject_import_button, QStyle.SP_DialogOpenButton)
        subject_toolbar.addWidget(self.subject_import_button)
        subject_toolbar.addStretch()
        subject_box.addLayout(subject_toolbar)
        subject_box.addWidget(self.subject_count_label)
        subject_layout.addLayout(subject_box)

        body_card, body_layout = self._card("BODY TABS", "Each body opens in a closable browser-style tab.")
        body_header = QHBoxLayout()
        body_title = QLabel("Bodies")
        body_title.setObjectName("sectionSubtitle")
        body_header.addWidget(body_title)
        body_header.addStretch()
        self.body_add_button.setObjectName("secondaryButton")
        self.body_add_button.setFixedWidth(_scaled_int(34, self._scale))
        self.body_add_button.setToolTip("Add a new body")
        self.body_add_button.clicked.connect(self._new_body_draft_tab)
        self._apply_button_icon(self.body_add_button, QStyle.SP_FileDialogNewFolder)
        self.body_upload_button.setObjectName("secondaryButton")
        self.body_upload_button.clicked.connect(self._upload_body_files)
        self.body_upload_button.setToolTip("Upload CSV text bodies or HTML body files")
        self._apply_button_icon(self.body_upload_button, QStyle.SP_DialogOpenButton)
        self.body_refresh_button.setObjectName("secondaryButton")
        self.body_refresh_button.clicked.connect(self.load_user_workspace)
        self.body_refresh_button.setToolTip("Reset the body workspace to a single active body tab")
        self._apply_button_icon(self.body_refresh_button, QStyle.SP_BrowserReload)
        body_header.addWidget(self.body_add_button)
        body_header.addWidget(self.body_upload_button)
        body_header.addWidget(self.body_refresh_button)
        body_layout.addLayout(body_header)

        self.body_tabs.setTabsClosable(True)
        self.body_tabs.setMovable(True)
        self.body_tabs.setDocumentMode(True)
        self.body_tabs.setUsesScrollButtons(True)
        self.body_tabs.tabCloseRequested.connect(self._remove_body_draft_tab)
        self.body_tabs.currentChanged.connect(self._on_body_tab_changed)
        self.body_tabs.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Expanding)
        body_layout.addWidget(self.body_tabs, 1)

        page_layout.addWidget(subject_card)
        page_layout.addWidget(body_card, 1)
        self.subject_input.textChanged.connect(lambda _text: self._schedule_subject_body_save())
        return self._tab_scroll(page)

    def _build_html_content_tab(self) -> QWidget:
        page = QWidget()
        layout = QVBoxLayout(page)
        layout.setSpacing(_scaled_int(10, self._scale))

        content_card, content_layout = self._card("", None)
        content_header = QHBoxLayout()
        content_header.setContentsMargins(0, 0, 0, 0)
        content_header.setSpacing(_scaled_int(8, self._scale))
        content_label = QLabel("ATTACHMENT CONTENT")
        content_label.setObjectName("sectionTitle")
        content_header.addWidget(content_label)
        content_header.addStretch()
        content_hint = QLabel("Upload HTML or paste HTML code, then choose an output format before sending.")
        content_hint.setObjectName("sectionHint")
        content_hint.setWordWrap(True)
        content_layout.addWidget(content_hint)
        self.attach_add_button.setObjectName("secondaryButton")
        self.attach_add_button.setFixedWidth(_scaled_int(34, self._scale))
        self.attach_add_button.setToolTip("Add a new HTML attachment tab")
        self.attach_add_button.clicked.connect(self._new_attachment_draft_tab)
        self._apply_button_icon(self.attach_add_button, QStyle.SP_FileDialogNewFolder)
        self.attach_upload_button.setObjectName("secondaryButton")
        self.attach_upload_button.clicked.connect(self._upload_attachment_files)
        self.attach_upload_button.setToolTip("Upload HTML files and create tabs")
        self._apply_button_icon(self.attach_upload_button, QStyle.SP_DialogOpenButton)
        self.attach_reset_button.setObjectName("secondaryButton")
        self.attach_reset_button.clicked.connect(self._reset_attachment_tabs)
        self.attach_reset_button.setToolTip("Reset attachments to a single blank tab")
        self._apply_button_icon(self.attach_reset_button, QStyle.SP_BrowserReload)
        content_header.addWidget(self.attach_add_button)
        content_header.addWidget(self.attach_upload_button)
        content_header.addWidget(self.attach_reset_button)
        content_layout.addLayout(content_header)

        self.attach_tabs.setTabsClosable(True)
        self.attach_tabs.setMovable(True)
        self.attach_tabs.setDocumentMode(True)
        self.attach_tabs.setUsesScrollButtons(True)
        self.attach_tabs.tabCloseRequested.connect(self._remove_attachment_draft_tab)
        self.attach_tabs.currentChanged.connect(self._on_attachment_tab_changed)
        self.attach_tabs.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Expanding)
        content_layout.addWidget(self.attach_tabs, 1)

        export_card, export_layout = self._card("ATTACHMENT FORMAT OPTIONS", "Choose the attachment file format and preview behavior for the active HTML content.")
        export_form = QFormLayout()
        export_form.setLabelAlignment(Qt.AlignLeft)

        self.attach_convert_checkbox.setChecked(True)
        self.attach_convert_checkbox.setToolTip("Enable file conversion options")
        self.attach_convert_checkbox.toggled.connect(self._sync_attachment_convert_controls)
        convert_row = QHBoxLayout()
        convert_row.setContentsMargins(0, 0, 0, 0)
        convert_row.setSpacing(_scaled_int(10, self._scale))
        convert_label = QLabel("Convert")
        convert_label.setObjectName("fieldLabel")
        convert_row.addWidget(convert_label)
        convert_row.addWidget(self.attach_convert_checkbox)
        convert_row.addStretch()
        export_form.addRow(self._wrap_layout(convert_row))

        self.attach_format_label.setObjectName("sectionSubtitle")
        self.attach_choose_format_button.setObjectName("secondaryButton")
        self.attach_choose_format_button.clicked.connect(self._choose_attachment_file_format)
        self.attach_choose_format_button.setToolTip("Open a modal to choose one or more export file formats")
        self._apply_button_icon(self.attach_choose_format_button, QStyle.SP_DialogApplyButton)
        format_row = QHBoxLayout()
        format_row.setContentsMargins(0, 0, 0, 0)
        format_row.setSpacing(_scaled_int(8, self._scale))
        format_row.addWidget(self.attach_choose_format_button)
        format_row.addWidget(self.attach_format_label)
        format_row.addStretch()
        self.attach_format_row_widget = self._wrap_layout(format_row)
        export_form.addRow("File format", self.attach_format_row_widget)
        self.attach_format_row_label = export_form.labelForField(self.attach_format_row_widget)

        export_layout.addLayout(export_form)
        self._sync_attachment_convert_controls(self.attach_convert_checkbox.isChecked())

        layout.addWidget(content_card, 1)
        layout.addWidget(export_card)
        if self.attach_tabs.count() == 0:
            self._add_attachment_draft_tab(select=True)
        self._sync_active_attachment_widget_refs()
        self._update_attachment_tab_controls()
        return self._tab_scroll(page)

    def _build_settings_tab(self) -> QWidget:
        page = QWidget()
        layout = QVBoxLayout(page)
        layout.setSpacing(_scaled_int(10, self._scale))

        header = self._section_title("SENDING SETTINGS")
        layout.addWidget(header)

        send_card, send_layout = self._card("Sending Settings", "Tune the execution behavior for each sender.")
        form = QFormLayout()
        form.setLabelAlignment(Qt.AlignLeft)
        self.sender_limit.setRange(1, 5000)
        self.sender_limit.setValue(300)
        self.sender_limit.setObjectName("windowSpin")
        self.sender_limit.setToolTip("Set the maximum emails per sender")
        self.delay_from.setDecimals(1)
        self.delay_from.setRange(0.0, 60.0)
        self.delay_from.setSingleStep(0.1)
        self.delay_from.setValue(0.5)
        self.delay_to.setDecimals(1)
        self.delay_to.setRange(0.0, 60.0)
        self.delay_to.setSingleStep(0.1)
        self.delay_to.setValue(1.0)
        self.retry_count.setRange(0, 20)
        self.retry_count.setValue(3)
        self.delay_from.setToolTip("Minimum delay between emails")
        self.delay_to.setToolTip("Maximum delay between emails")
        self.retry_count.setToolTip("Number of retry attempts for failed sends")

        form.addRow("Per-sender limit", self.sender_limit)
        delay_row = QHBoxLayout()
        delay_row.addWidget(QLabel("Delay between emails"))
        delay_row.addWidget(self.delay_from)
        delay_row.addWidget(QLabel("to"))
        delay_row.addWidget(self.delay_to)
        delay_row.addWidget(QLabel("seconds (random range)"))
        delay_holder = QFrame()
        delay_holder.setLayout(delay_row)
        form.addRow(delay_holder)
        send_layout.addLayout(form)

        delay_type_row = QHBoxLayout()
        self.delay_fixed_radio.setToolTip("Use the same delay for every send")
        self.delay_random_radio.setToolTip("Use a random delay within the range")
        self.delay_human_radio.setToolTip("Use a human-like delay pattern")
        self.delay_type_group = QButtonGroup(self)
        self.delay_type_group.setExclusive(True)
        for button in (self.delay_fixed_radio, self.delay_random_radio, self.delay_human_radio):
            self.delay_type_group.addButton(button)
            delay_type_row.addWidget(button)
        delay_type_row.addStretch()
        send_layout.addWidget(self._labeled_value_row("Delay type", self._wrap_layout(delay_type_row)))

        retry_row = QHBoxLayout()
        self.retry_enable_checkbox.setChecked(True)
        self.retry_enable_checkbox.setToolTip("Enable retry handling for failed sends")
        retry_row.addWidget(self.retry_enable_checkbox)
        retry_row.addWidget(self.retry_count)
        retry_row.addWidget(QLabel("retries"))
        retry_row.addStretch()
        send_layout.addWidget(self._labeled_value_row("Retry failed sends", self._wrap_layout(retry_row)))

        order_row = QHBoxLayout()
        self.send_seq_radio.setToolTip("Send in list order")
        self.send_rand_radio.setToolTip("Shuffle the send order")
        self.send_order_group = QButtonGroup(self)
        self.send_order_group.setExclusive(True)
        for button in (self.send_seq_radio, self.send_rand_radio):
            self.send_order_group.addButton(button)
            order_row.addWidget(button)
        order_row.addStretch()
        send_layout.addWidget(self._labeled_value_row("Email send order", self._wrap_layout(order_row)))

        window_mode_row = QHBoxLayout()
        self.window_parallel_radio.setToolTip("Launch and send in parallel")
        self.window_sequential_radio.setToolTip("Rotate through windows one at a time")
        self.window_mode_group = QButtonGroup(self)
        self.window_mode_group.setExclusive(True)
        for button in (self.window_parallel_radio, self.window_sequential_radio):
            self.window_mode_group.addButton(button)
            window_mode_row.addWidget(button)
        window_mode_row.addStretch()
        send_layout.addWidget(self._labeled_value_row("Window send mode", self._wrap_layout(window_mode_row)))

        layout.addWidget(send_card)

        ai_card, ai_layout = self._card("AI ASSISTANT", "Connect an AI provider to unlock model selection.")
        ai_form = QFormLayout()
        ai_form.setLabelAlignment(Qt.AlignLeft)
        ai_form.setFieldGrowthPolicy(QFormLayout.AllNonFixedFieldsGrow)
        ai_form.setHorizontalSpacing(_scaled_int(16, self._scale))
        ai_form.setVerticalSpacing(_scaled_int(12, self._scale))

        self.ai_provider_combo.clear()
        self.ai_provider_combo.addItems(["ChatGPT", "Claude", "DeepSeek"])
        self.ai_provider_combo.setCurrentText(self.state.ai_provider)
        self.ai_provider_combo.setToolTip("Select the AI provider")
        self.ai_provider_combo.currentTextChanged.connect(lambda _value: self._on_ai_provider_changed())
        self.ai_provider_combo.setMinimumWidth(_scaled_int(240, self._scale))
        self.ai_provider_combo.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Fixed)

        self.ai_api_key_input.setPlaceholderText("Enter API key")
        self.ai_api_key_input.setEchoMode(QLineEdit.Password)
        self.ai_api_key_input.setToolTip("Enter the provider API key")
        self.ai_api_key_input.setMinimumWidth(_scaled_int(240, self._scale))
        self.ai_api_key_input.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Fixed)

        self.ai_model_combo.setEnabled(False)
        self.ai_model_combo.setToolTip("Choose a model after connecting")
        self.ai_model_combo.currentTextChanged.connect(lambda _value: self._on_ai_model_changed())
        self.ai_model_combo.setMinimumWidth(_scaled_int(240, self._scale))
        self.ai_model_combo.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Fixed)

        self.ai_connect_button.setObjectName("primaryButton")
        self.ai_connect_button.clicked.connect(lambda: self._connect_ai_provider())
        self.ai_connect_button.setToolTip("Validate the API key and connect")
        self.ai_connect_button.setMinimumWidth(_scaled_int(140, self._scale))
        self.ai_connect_button.setCursor(Qt.PointingHandCursor)

        self.ai_status_label.setObjectName("windowPill")
        self.ai_status_label.setAlignment(Qt.AlignCenter)
        self.ai_status_label.setMinimumHeight(_scaled_int(32, self._scale))
        self.ai_status_label.setMinimumWidth(_scaled_int(140, self._scale))
        self.ai_status_label.setText("Not connected")

        ai_form.addRow("Provider", self.ai_provider_combo)
        ai_form.addRow("API key", self.ai_api_key_input)
        ai_form.addRow("", self.ai_connect_button)
        ai_form.addRow("Model", self.ai_model_combo)
        ai_form.addRow("Status", self.ai_status_label)
        ai_layout.addLayout(ai_form)
        layout.addWidget(ai_card)

        save_button = QPushButton("Save Settings")
        save_button.setObjectName("primaryButton")
        save_button.clicked.connect(lambda: self._save_sending_settings())
        save_button.setToolTip("Save the current sending settings")
        self._apply_button_icon(save_button, QStyle.SP_DialogSaveButton)
        layout.addWidget(save_button, alignment=Qt.AlignLeft)
        layout.addStretch()

        self.sender_limit.valueChanged.connect(lambda _value: self._schedule_sending_settings_save())
        self.delay_from.valueChanged.connect(lambda _value: self._schedule_sending_settings_save())
        self.delay_to.valueChanged.connect(lambda _value: self._schedule_sending_settings_save())
        self.retry_count.valueChanged.connect(lambda _value: self._schedule_sending_settings_save())
        self.retry_enable_checkbox.toggled.connect(lambda _value: self._schedule_sending_settings_save())
        self.delay_type_group.buttonToggled.connect(lambda *_args: self._schedule_sending_settings_save())
        self.send_order_group.buttonToggled.connect(lambda *_args: self._schedule_sending_settings_save())
        self.window_mode_group.buttonToggled.connect(lambda *_args: self._schedule_sending_settings_save())
        self.ai_provider_combo.currentTextChanged.connect(lambda _value: self._schedule_sending_settings_save())
        self.ai_api_key_input.textEdited.connect(lambda _value: self._on_ai_api_key_changed())
        self.ai_model_combo.currentTextChanged.connect(lambda _value: self._schedule_sending_settings_save())

        self._refresh_ai_models()
        self._sync_ai_connection_ui()

        return self._tab_scroll(page)

    def _refresh_ai_models(self) -> None:
        self.ai_model_combo.blockSignals(True)
        self.ai_model_combo.clear()
        self.ai_model_combo.addItems(self._available_ai_models)
        if self.state.ai_model and self.state.ai_model in self._available_ai_models:
            self.ai_model_combo.setCurrentText(self.state.ai_model)
        elif self._available_ai_models:
            self.ai_model_combo.setCurrentIndex(0)
            self.state.ai_model = self.ai_model_combo.currentText()
        self.ai_model_combo.blockSignals(False)

    def _sync_ai_connection_ui(self) -> None:
        connected = self.state.ai_connected
        self.ai_model_combo.setEnabled(connected and bool(self._available_ai_models))
        self.ai_connect_button.setText("Reconnect" if connected else "Connect")
        self.ai_connect_button.setEnabled(True)
        if connected:
            self.ai_status_label.setText("Connected")
            self.ai_status_label.setStyleSheet(
                "QLabel { background:#12351e; color:#7dff9a; border:1px solid #1e7a3a; border-radius:8px; padding:4px 10px; font-weight:700; }"
            )
        else:
            self.ai_status_label.setText("Not connected")
            self.ai_status_label.setStyleSheet(
                "QLabel { background:#2a1d1d; color:#f0a0a0; border:1px solid #7a3131; border-radius:8px; padding:4px 10px; font-weight:700; }"
            )
        self.ai_model_combo.blockSignals(True)
        if connected and self.state.ai_model:
            index = self.ai_model_combo.findText(self.state.ai_model)
            if index >= 0:
                self.ai_model_combo.setCurrentIndex(index)
        self.ai_model_combo.blockSignals(False)

    def _current_sending_settings_payload(self) -> dict[str, object]:
        delay_type = "Random range"
        if self.delay_fixed_radio.isChecked():
            delay_type = "Fixed"
        elif self.delay_human_radio.isChecked():
            delay_type = "Human-like pattern"

        email_send_order = "Sequential"
        if self.send_rand_radio.isChecked():
            email_send_order = "Random shuffle"

        window_send_mode = "Parallel"
        if self.window_sequential_radio.isChecked():
            window_send_mode = "Sequential"

        return {
            "sender_limit": int(self.sender_limit.value()),
            "delay_from": float(self.delay_from.value()),
            "delay_to": float(self.delay_to.value()),
            "retry_count": int(self.retry_count.value()),
            "retry_enabled": bool(self.retry_enable_checkbox.isChecked()),
            "delay_type": delay_type,
            "email_send_order": email_send_order,
            "window_send_mode": window_send_mode,
            "ai_provider": self.ai_provider_combo.currentText().strip() or "ChatGPT",
            "ai_api_key": self.ai_api_key_input.text(),
            "ai_model": self.ai_model_combo.currentText().strip(),
            "ai_connected": bool(self.state.ai_connected),
            "available_models": list(self._available_ai_models),
        }

    def _apply_sending_settings_payload(self, payload: dict[str, object]) -> None:
        def _as_int(value: object, default: int) -> int:
            try:
                return int(float(value))
            except Exception:
                return default

        def _as_float(value: object, default: float) -> float:
            try:
                return float(value)
            except Exception:
                return default

        def _as_bool(value: object, default: bool) -> bool:
            if isinstance(value, bool):
                return value
            if isinstance(value, (int, float)):
                return bool(value)
            if isinstance(value, str):
                lowered = value.strip().lower()
                if lowered in {"1", "true", "yes", "on"}:
                    return True
                if lowered in {"0", "false", "no", "off"}:
                    return False
            return default

        def _block(widget: QWidget, value: Callable[[], None]) -> None:
            widget.blockSignals(True)
            try:
                value()
            finally:
                widget.blockSignals(False)

        sender_limit = _as_int(payload.get("sender_limit"), 300)
        delay_from = _as_float(payload.get("delay_from"), 0.5)
        delay_to = _as_float(payload.get("delay_to"), 1.0)
        retry_count = _as_int(payload.get("retry_count"), 3)
        retry_enabled = _as_bool(payload.get("retry_enabled"), True)
        delay_type = str(payload.get("delay_type") or "Random range")
        email_send_order = str(payload.get("email_send_order") or "Sequential")
        window_send_mode = str(payload.get("window_send_mode") or "Parallel")
        ai_provider = str(payload.get("ai_provider") or "ChatGPT")
        ai_api_key = str(payload.get("ai_api_key") or "")
        ai_model = str(payload.get("ai_model") or "")
        ai_connected = _as_bool(payload.get("ai_connected"), False)
        available_models_raw = payload.get("available_models")
        available_models = available_models_raw if isinstance(available_models_raw, list) else []

        _block(self.sender_limit, lambda: self.sender_limit.setValue(sender_limit))
        _block(self.delay_from, lambda: self.delay_from.setValue(delay_from))
        _block(self.delay_to, lambda: self.delay_to.setValue(delay_to))
        _block(self.retry_count, lambda: self.retry_count.setValue(retry_count))
        _block(self.retry_enable_checkbox, lambda: self.retry_enable_checkbox.setChecked(retry_enabled))

        _block(self.delay_fixed_radio, lambda: self.delay_fixed_radio.setChecked(delay_type == "Fixed"))
        _block(self.delay_random_radio, lambda: self.delay_random_radio.setChecked(delay_type == "Random range"))
        _block(self.delay_human_radio, lambda: self.delay_human_radio.setChecked(delay_type == "Human-like pattern"))
        if not any((self.delay_fixed_radio.isChecked(), self.delay_random_radio.isChecked(), self.delay_human_radio.isChecked())):
            self.delay_random_radio.setChecked(True)

        _block(self.send_seq_radio, lambda: self.send_seq_radio.setChecked(email_send_order != "Random shuffle"))
        _block(self.send_rand_radio, lambda: self.send_rand_radio.setChecked(email_send_order == "Random shuffle"))

        _block(self.window_parallel_radio, lambda: self.window_parallel_radio.setChecked(window_send_mode != "Sequential"))
        _block(self.window_sequential_radio, lambda: self.window_sequential_radio.setChecked(window_send_mode == "Sequential"))

        _block(self.ai_provider_combo, lambda: self.ai_provider_combo.setCurrentText(ai_provider))
        _block(self.ai_api_key_input, lambda: self.ai_api_key_input.setText(ai_api_key))

        self.state.sender_limit = sender_limit
        self.state.delay_from = delay_from
        self.state.delay_to = delay_to
        self.state.retry_count = retry_count
        self.state.retry_enabled = retry_enabled
        self.state.delay_type = delay_type
        self.state.email_send_order = email_send_order
        self.state.window_send_mode = window_send_mode
        self.state.ai_provider = ai_provider
        self.state.ai_api_key = ai_api_key
        self.state.ai_model = ai_model
        self.state.ai_connected = ai_connected
        self._available_ai_models = [str(item) for item in available_models if str(item).strip()]
        self.state.ai_available_models = list(self._available_ai_models)

        self._refresh_ai_models()
        self._sync_ai_connection_ui()

    def _load_sending_settings_state(self, mysql_settings_map: dict[str, object] | None = None) -> None:
        local_payload = _load_ui_state(LOCAL_SETTINGS_STATE_KEY)
        payload: dict[str, object] = {}
        if isinstance(mysql_settings_map, dict):
            payload.update(mysql_settings_map)
        if isinstance(local_payload, dict):
            payload.update(local_payload)
        self._apply_sending_settings_payload(payload)

    def _schedule_sending_settings_save(self) -> None:
        if self._workspace_loading:
            return
        self._settings_save_timer.start()

    def _persist_sending_settings_state(self) -> None:
        if self._workspace_loading:
            return
        payload = self._current_sending_settings_payload()
        _upsert_ui_state(LOCAL_SETTINGS_STATE_KEY, payload)
        self.state.sender_limit = int(payload["sender_limit"])
        self.state.delay_from = float(payload["delay_from"])
        self.state.delay_to = float(payload["delay_to"])
        self.state.retry_count = int(payload["retry_count"])
        self.state.retry_enabled = bool(payload["retry_enabled"])
        self.state.delay_type = str(payload["delay_type"])
        self.state.email_send_order = str(payload["email_send_order"])
        self.state.window_send_mode = str(payload["window_send_mode"])
        self.state.ai_provider = str(payload["ai_provider"])
        self.state.ai_api_key = str(payload["ai_api_key"])
        self.state.ai_model = str(payload["ai_model"])
        self.state.ai_connected = bool(payload["ai_connected"])
        self.state.ai_available_models = [str(item) for item in payload.get("available_models") or []]
        self._available_ai_models = list(self.state.ai_available_models)
        if self.state.logged_in and self.state.username:
            try:
                for key, value in payload.items():
                    upsert_setting(self.state.username, str(key), value)
            except Exception:
                pass

    def _save_sending_settings(self) -> None:
        self._persist_sending_settings_state()
        self._log_action("Saved sending settings")
        self.notify("Settings saved")

    def _set_ai_busy(self, busy: bool) -> None:
        self.ai_provider_combo.setEnabled(not busy)
        self.ai_api_key_input.setEnabled(not busy)
        self.ai_model_combo.setEnabled((not busy) and self.state.ai_connected and bool(self._available_ai_models))
        self.ai_connect_button.setEnabled(not busy)
        if busy:
            self.ai_connect_button.setText("Connecting...")
            self.ai_connect_button.setStyleSheet(
                "QPushButton#primaryButton { background:#a87917; color:#ffffff; border:none; }"
            )
            QApplication.setOverrideCursor(Qt.WaitCursor)
        else:
            self.ai_connect_button.setStyleSheet("")
            QApplication.restoreOverrideCursor()

    def _on_ai_provider_changed(self) -> None:
        self.state.ai_provider = self.ai_provider_combo.currentText() or "ChatGPT"
        self.state.ai_connected = False
        self.state.ai_model = ""
        self._available_ai_models = []
        self._refresh_ai_models()
        self._sync_ai_connection_ui()
        self._log_action(f"AI provider selected: {self.state.ai_provider}")
        self._schedule_sending_settings_save()

    def _on_ai_api_key_changed(self) -> None:
        api_key = self.ai_api_key_input.text()
        self.state.ai_api_key = api_key
        self.state.ai_connected = False
        self.state.ai_model = ""
        self._available_ai_models = []
        self._refresh_ai_models()
        self._sync_ai_connection_ui()
        self._schedule_sending_settings_save()

    def _connect_ai_provider(self) -> None:
        provider = self.ai_provider_combo.currentText() or "ChatGPT"
        api_key = self.ai_api_key_input.text().strip()
        self.state.ai_api_key = api_key
        if not api_key:
            self.state.ai_connected = False
            self.state.ai_model = ""
            self._available_ai_models = []
            self._refresh_ai_models()
            self._sync_ai_connection_ui()
            self.notify("Enter an API key to connect")
            self._persist_sending_settings_state()
            return

        self._set_ai_busy(True)
        QApplication.processEvents()
        self.state.ai_provider = provider
        self.state.ai_connected = False
        self.state.ai_model = ""
        self._available_ai_models = []
        self._refresh_ai_models()
        self._sync_ai_connection_ui()

        try:
            models = AIValidationWorker(provider, api_key)._fetch_models()
        except Exception as exc:
            self.state.ai_connected = False
            self.state.ai_model = ""
            self._available_ai_models = []
            self._refresh_ai_models()
            self._sync_ai_connection_ui()
            self._set_ai_busy(False)
            self._log_action(f"AI connection failed: {exc}")
            self._persist_sending_settings_state()
            self.notify("AI connection failed")
            return

        self.state.ai_provider = provider
        self.state.ai_connected = True
        self._available_ai_models = models
        self._refresh_ai_models()
        if self._available_ai_models:
            if not self.state.ai_model or self.state.ai_model not in self._available_ai_models:
                self.state.ai_model = self._available_ai_models[0]
            self.ai_model_combo.setCurrentText(self.state.ai_model)
        self._sync_ai_connection_ui()
        self._set_ai_busy(False)
        self._log_action(f"Connected AI provider: {provider}")
        self._persist_sending_settings_state()
        self.notify(f"{provider} connected")

    def _on_ai_model_changed(self) -> None:
        if not self.ai_model_combo.isEnabled():
            return
        self.state.ai_model = self.ai_model_combo.currentText().strip()
        if self.state.ai_model:
            self._log_action(f"AI model selected: {self.state.ai_model}")
            self._schedule_sending_settings_save()

    def _build_blaster_tab(self) -> QWidget:
        page = QWidget()
        layout = QVBoxLayout(page)
        layout.setSpacing(_scaled_int(10, self._scale))

        header = self._section_title("CAMPAIGN SENDING CONTROL", "Send emails from all open Gmail windows.")
        layout.addWidget(header)

        controls_card, controls_layout = self._card("Active Gmail Windows")
        windows_row = QHBoxLayout()
        windows_row.addWidget(QLabel("Active Gmail Windows:"))
        self.active_windows_value.setObjectName("windowPill")
        self.active_windows_value.setToolTip("Count of active browser windows")
        windows_row.addWidget(self.active_windows_value)
        windows_row.addStretch()
        refresh_button = QPushButton("Refresh")
        refresh_button.setObjectName("secondaryButton")
        refresh_button.clicked.connect(lambda: self._log_action("Refreshed active Gmail windows"))
        refresh_button.setToolTip("Refresh the active Gmail window count")
        self._apply_button_icon(refresh_button, QStyle.SP_BrowserReload)
        windows_row.addWidget(refresh_button)
        self.campaign_reset_all_button.setObjectName("secondaryButton")
        self.campaign_reset_all_button.clicked.connect(self._reset_campaign_form_state)
        self.campaign_reset_all_button.setToolTip("Clear recipient data, subject/body, and attachment content")
        self._apply_button_icon(self.campaign_reset_all_button, QStyle.SP_BrowserReload)
        windows_row.addWidget(self.campaign_reset_all_button)
        controls_layout.addLayout(windows_row)

        controls_layout.addWidget(QLabel("Campaign Progress:"))
        self.progress_bar.setRange(0, 100)
        if not self._campaign_active:
            self.progress_bar.setValue(0)
        controls_layout.addWidget(self.progress_bar)
        self.campaign_progress_text.setObjectName("sectionHint")
        controls_layout.addWidget(self.campaign_progress_text)
        layout.addWidget(controls_card)

        self.start_campaign_button.setObjectName("blastButton")
        self.start_campaign_button.setMinimumHeight(_scaled_int(54, self._scale))
        self.start_campaign_button.clicked.connect(lambda: self._handle_campaign_primary_action())
        self.start_campaign_button.setToolTip("Start the email sending workflow")
        self._apply_button_icon(self.start_campaign_button, QStyle.SP_MediaPlay)
        layout.addWidget(self.start_campaign_button)

        campaign_action_row = QHBoxLayout()
        self.campaign_pause_button.setObjectName("warningButton")
        self.campaign_pause_button.clicked.connect(lambda: self._handle_campaign_pause_resume())
        self.campaign_pause_button.setToolTip("Pause or resume the active campaign")
        self._apply_button_icon(self.campaign_pause_button, QStyle.SP_MediaPause)
        self.campaign_cancel_button.setObjectName("dangerButton")
        self.campaign_cancel_button.clicked.connect(lambda: self._handle_campaign_cancel())
        self.campaign_cancel_button.setToolTip("Cancel the active campaign")
        self._apply_button_icon(self.campaign_cancel_button, QStyle.SP_DialogCancelButton)
        campaign_action_row.addWidget(self.campaign_pause_button)
        campaign_action_row.addWidget(self.campaign_cancel_button)
        campaign_action_row.addStretch()
        layout.addLayout(campaign_action_row)

        log_card, log_layout = self._card("SEND LOG")
        self.send_log_view.setReadOnly(True)
        self.send_log_view.setTextInteractionFlags(Qt.TextSelectableByMouse | Qt.TextSelectableByKeyboard)
        self.send_log_view.setObjectName("activityList")
        log_layout.addWidget(self.send_log_view)
        layout.addWidget(log_card, 1)

        return self._tab_scroll(page)

    def _build_admin_tab(self) -> QWidget:
        page = QWidget()
        layout = QVBoxLayout(page)
        layout.setSpacing(_scaled_int(10, self._scale))

        layout.addWidget(self._section_title("ADMIN CONSOLE", "Manage users, validity windows, device locks, and audit logs."))
        self.admin_access_label.setObjectName("sectionSubtitle")
        self.admin_access_label.setWordWrap(True)
        layout.addWidget(self.admin_access_label)

        users_card, users_layout = self._card("USERS", "Select a user, edit details, or create a new account.")
        users_header = QHBoxLayout()
        self.admin_refresh_users_button.setObjectName("secondaryButton")
        self.admin_refresh_users_button.clicked.connect(self._admin_refresh_users)
        self._apply_button_icon(self.admin_refresh_users_button, QStyle.SP_BrowserReload)
        users_header.addWidget(self.admin_refresh_users_button)
        users_header.addStretch()
        users_layout.addLayout(users_header)

        self.admin_users_table.setHorizontalHeaderLabels(
            ["ID", "Username", "Display Name", "Role", "Active", "Valid Until", "Device", "Last Login"]
        )
        self.admin_users_table.setSelectionBehavior(QAbstractItemView.SelectRows)
        self.admin_users_table.setSelectionMode(QAbstractItemView.SingleSelection)
        self.admin_users_table.setEditTriggers(QAbstractItemView.NoEditTriggers)
        self.admin_users_table.horizontalHeader().setStretchLastSection(True)
        self.admin_users_table.horizontalHeader().setSectionResizeMode(QHeaderView.ResizeToContents)
        self.admin_users_table.horizontalHeader().setSectionResizeMode(2, QHeaderView.Stretch)
        self.admin_users_table.horizontalHeader().setSectionResizeMode(6, QHeaderView.Stretch)
        self.admin_users_table.horizontalHeader().setSectionResizeMode(7, QHeaderView.Stretch)
        self.admin_users_table.itemSelectionChanged.connect(self._admin_load_selected_user)
        users_layout.addWidget(self.admin_users_table, 1)

        form_card, form_layout = self._card("USER DETAILS", "Edit the selected user or prepare a new one.")
        form = QFormLayout()
        form.setLabelAlignment(Qt.AlignLeft)
        self.admin_selected_user_id.setReadOnly(True)
        self.admin_username_input.setPlaceholderText("username")
        self.admin_display_name_input.setPlaceholderText("display name")
        self.admin_role_input.setPlaceholderText("admin or user")
        self.admin_valid_until_input.setPlaceholderText("YYYY-MM-DD HH:MM:SS or blank")
        self.admin_new_password_input.setPlaceholderText("new password")
        self.admin_new_password_input.setEchoMode(QLineEdit.Password)
        form.addRow("User ID", self.admin_selected_user_id)
        form.addRow("Username", self.admin_username_input)
        form.addRow("Display Name", self.admin_display_name_input)
        form.addRow("Role", self.admin_role_input)
        form.addRow("Login Valid Until", self.admin_valid_until_input)
        form.addRow("Reset Password", self.admin_new_password_input)
        form.addRow("", self.admin_clear_device_checkbox)
        form_layout.addLayout(form)

        actions_row = QHBoxLayout()
        self.admin_create_user_button.setObjectName("primaryButton")
        self.admin_create_user_button.clicked.connect(self._admin_create_user)
        self._apply_button_icon(self.admin_create_user_button, QStyle.SP_FileDialogNewFolder)
        self.admin_save_user_button.setObjectName("secondaryButton")
        self.admin_save_user_button.clicked.connect(self._admin_save_user)
        self._apply_button_icon(self.admin_save_user_button, QStyle.SP_DialogSaveButton)
        self.admin_activate_user_button.setObjectName("secondaryButton")
        self.admin_activate_user_button.clicked.connect(lambda: self._admin_toggle_active(True))
        self._apply_button_icon(self.admin_activate_user_button, QStyle.SP_MediaPlay)
        self.admin_deactivate_user_button.setObjectName("secondaryButton")
        self.admin_deactivate_user_button.clicked.connect(lambda: self._admin_toggle_active(False))
        self._apply_button_icon(self.admin_deactivate_user_button, QStyle.SP_MediaPause)
        self.admin_reset_password_button.setObjectName("secondaryButton")
        self.admin_reset_password_button.clicked.connect(self._admin_reset_password)
        self._apply_button_icon(self.admin_reset_password_button, QStyle.SP_BrowserReload)
        self.admin_reset_device_button.setObjectName("secondaryButton")
        self.admin_reset_device_button.clicked.connect(self._admin_reset_device)
        self._apply_button_icon(self.admin_reset_device_button, QStyle.SP_BrowserReload)
        for button in (
            self.admin_create_user_button,
            self.admin_save_user_button,
            self.admin_activate_user_button,
            self.admin_deactivate_user_button,
            self.admin_reset_password_button,
            self.admin_reset_device_button,
        ):
            actions_row.addWidget(button)
        actions_row.addStretch()
        form_layout.addLayout(actions_row)
        users_layout.addWidget(form_card)

        activity_card, activity_layout = self._card("ACTIVITY", "Recent actions for all users.")
        activity_header = QHBoxLayout()
        self.admin_refresh_activity_button.setObjectName("secondaryButton")
        self.admin_refresh_activity_button.clicked.connect(self._admin_refresh_activity)
        self._apply_button_icon(self.admin_refresh_activity_button, QStyle.SP_BrowserReload)
        activity_header.addWidget(self.admin_refresh_activity_button)
        activity_header.addStretch()
        activity_layout.addLayout(activity_header)
        self.admin_activity_table.setHorizontalHeaderLabels(
            ["ID", "Username", "Category", "Action", "Details", "IP", "Created"]
        )
        self.admin_activity_table.setSelectionBehavior(QAbstractItemView.SelectRows)
        self.admin_activity_table.setSelectionMode(QAbstractItemView.SingleSelection)
        self.admin_activity_table.setEditTriggers(QAbstractItemView.NoEditTriggers)
        self.admin_activity_table.horizontalHeader().setStretchLastSection(True)
        self.admin_activity_table.horizontalHeader().setSectionResizeMode(QHeaderView.ResizeToContents)
        self.admin_activity_table.horizontalHeader().setSectionResizeMode(4, QHeaderView.Stretch)
        activity_layout.addWidget(self.admin_activity_table, 1)

        login_card, login_layout = self._card("LOGIN HISTORY", "Login attempts, device bindings, and access validity.")
        login_header = QHBoxLayout()
        self.admin_refresh_login_button.setObjectName("secondaryButton")
        self.admin_refresh_login_button.clicked.connect(self._admin_refresh_login_history)
        self._apply_button_icon(self.admin_refresh_login_button, QStyle.SP_BrowserReload)
        login_header.addWidget(self.admin_refresh_login_button)
        login_header.addStretch()
        login_layout.addLayout(login_header)
        self.admin_login_history_table.setHorizontalHeaderLabels(
            ["ID", "User ID", "Username", "Success", "IP", "Device Fingerprint", "Device Name", "Agent", "Created"]
        )
        self.admin_login_history_table.setSelectionBehavior(QAbstractItemView.SelectRows)
        self.admin_login_history_table.setSelectionMode(QAbstractItemView.SingleSelection)
        self.admin_login_history_table.setEditTriggers(QAbstractItemView.NoEditTriggers)
        self.admin_login_history_table.horizontalHeader().setStretchLastSection(True)
        self.admin_login_history_table.horizontalHeader().setSectionResizeMode(QHeaderView.ResizeToContents)
        self.admin_login_history_table.horizontalHeader().setSectionResizeMode(5, QHeaderView.Stretch)
        self.admin_login_history_table.horizontalHeader().setSectionResizeMode(6, QHeaderView.Stretch)
        self.admin_login_history_table.horizontalHeader().setSectionResizeMode(7, QHeaderView.Stretch)
        login_layout.addWidget(self.admin_login_history_table, 1)

        self.admin_tabs = QTabWidget()
        self.admin_tabs.addTab(users_card, "Users")
        self.admin_tabs.addTab(activity_card, "Activity")
        self.admin_tabs.addTab(login_card, "Login History")
        self._apply_navigation_icons()
        layout.addWidget(self.admin_tabs, 1)
        return self._tab_scroll(page)

    def _admin_table_item(self, value: object) -> QTableWidgetItem:
        item = QTableWidgetItem("" if value is None else str(value))
        item.setFlags(item.flags() & ~Qt.ItemIsEditable)
        return item

    def _admin_format_datetime(self, value: object) -> str:
        if value in {None, ""}:
            return ""
        try:
            if isinstance(value, datetime):
                dt = value
            else:
                text = str(value).strip().replace("Z", "+00:00")
                dt = datetime.fromisoformat(text)
            if dt.tzinfo is not None:
                dt = dt.astimezone(timezone.utc).replace(tzinfo=None)
            return dt.strftime("%Y-%m-%d %H:%M:%S")
        except Exception:
            return str(value)

    def _admin_parse_datetime(self, value: str) -> datetime | None:
        text = value.strip()
        if not text:
            return None
        try:
            return datetime.fromisoformat(text.replace("Z", "+00:00"))
        except Exception:
            return None

    def _admin_current_user_id(self) -> int | None:
        current_item = self.admin_users_table.currentItem()
        if current_item is None:
            return None
        row = current_item.row()
        id_item = self.admin_users_table.item(row, 0)
        if id_item is None:
            return None
        try:
            return int(id_item.text().strip())
        except Exception:
            return None

    def _admin_clear_form(self) -> None:
        for widget in (
            self.admin_selected_user_id,
            self.admin_username_input,
            self.admin_display_name_input,
            self.admin_role_input,
            self.admin_valid_until_input,
            self.admin_new_password_input,
        ):
            widget.clear()
        self.admin_clear_device_checkbox.setChecked(False)

    def _admin_set_form_from_user(self, user: dict[str, object]) -> None:
        self.admin_selected_user_id.setText(str(user.get("id") or ""))
        self.admin_username_input.setText(str(user.get("username") or ""))
        self.admin_display_name_input.setText(str(user.get("display_name") or ""))
        self.admin_role_input.setText(str(user.get("role") or "user"))
        self.admin_valid_until_input.setText(self._admin_format_datetime(user.get("login_valid_until")))
        self.admin_new_password_input.clear()
        self.admin_clear_device_checkbox.setChecked(False)

    def _admin_refresh_users(self) -> None:
        if not self.state.logged_in or not self.state.auth_token:
            return
        try:
            payload = list_admin_users(self.state.auth_token)
            users = payload.get("users") or []
            if not isinstance(users, list):
                users = []
            self.admin_users_table.setRowCount(0)
            for row_index, user in enumerate(users):
                if not isinstance(user, dict):
                    continue
                self.admin_users_table.insertRow(row_index)
                values = [
                    user.get("id"),
                    user.get("username"),
                    user.get("display_name"),
                    user.get("role"),
                    "Yes" if bool(user.get("is_active", True)) else "No",
                    self._admin_format_datetime(user.get("login_valid_until")),
                    user.get("device_name") or user.get("device_fingerprint"),
                    user.get("last_login_at") or user.get("last_login_ip"),
                ]
                for col, value in enumerate(values):
                    self.admin_users_table.setItem(row_index, col, self._admin_table_item(value))
            if self.admin_users_table.rowCount() > 0 and self.admin_users_table.currentRow() < 0:
                self.admin_users_table.selectRow(0)
        except Exception as exc:
            self._log_action(f"Failed to refresh admin users: {exc}")

    def _admin_refresh_activity(self) -> None:
        if not self.state.logged_in or not self.state.auth_token:
            return
        try:
            payload = list_admin_activity(self.state.auth_token, limit=200)
            rows = payload.get("activity") or []
            if not isinstance(rows, list):
                rows = []
            self.admin_activity_table.setRowCount(0)
            for row_index, row in enumerate(rows):
                if not isinstance(row, dict):
                    continue
                self.admin_activity_table.insertRow(row_index)
                values = [
                    row.get("id"),
                    row.get("username"),
                    row.get("category"),
                    row.get("action"),
                    row.get("details_json"),
                    row.get("ip_address"),
                    row.get("created_at"),
                ]
                for col, value in enumerate(values):
                    self.admin_activity_table.setItem(row_index, col, self._admin_table_item(value))
        except Exception as exc:
            self._log_action(f"Failed to refresh admin activity: {exc}")

    def _admin_refresh_login_history(self) -> None:
        if not self.state.logged_in or not self.state.auth_token:
            return
        try:
            payload = list_admin_login_history(self.state.auth_token, limit=200)
            rows = payload.get("history") or []
            if not isinstance(rows, list):
                rows = []
            self.admin_login_history_table.setRowCount(0)
            for row_index, row in enumerate(rows):
                if not isinstance(row, dict):
                    continue
                self.admin_login_history_table.insertRow(row_index)
                values = [
                    row.get("id"),
                    row.get("user_id"),
                    row.get("username"),
                    "Yes" if bool(row.get("success")) else "No",
                    row.get("ip_address"),
                    row.get("device_fingerprint"),
                    row.get("device_name"),
                    row.get("user_agent"),
                    row.get("created_at"),
                ]
                for col, value in enumerate(values):
                    self.admin_login_history_table.setItem(row_index, col, self._admin_table_item(value))
        except Exception as exc:
            self._log_action(f"Failed to refresh login history: {exc}")

    def _admin_load_selected_user(self) -> None:
        user_id = self._admin_current_user_id()
        if user_id is None or not self.state.logged_in or not self.state.auth_token:
            return
        try:
            payload = get_admin_user(self.state.auth_token, user_id)
            user = payload.get("user") or {}
            if isinstance(user, dict):
                self._admin_set_form_from_user(user)
        except Exception as exc:
            self._log_action(f"Failed to load selected admin user: {exc}")

    def _admin_create_user(self) -> None:
        if not self.state.logged_in or not self.state.auth_token:
            return
        password = self.admin_new_password_input.text().strip()
        if not password:
            self.notify("Enter a password to create the user")
            return
        payload = {
            "username": self.admin_username_input.text().strip(),
            "password": password,
            "display_name": self.admin_display_name_input.text().strip(),
            "role": self.admin_role_input.text().strip() or "user",
            "is_active": True,
            "login_valid_until": self._admin_parse_datetime(self.admin_valid_until_input.text().strip()),
        }
        try:
            create_admin_user(self.state.auth_token, payload)
            self.notify("User created")
            self._admin_refresh_users()
            self._admin_clear_form()
        except Exception as exc:
            self._log_action(f"Failed to create admin user: {exc}")
            self.notify("Unable to create user")

    def _admin_save_user(self) -> None:
        user_id = self._admin_current_user_id()
        if user_id is None or not self.state.logged_in or not self.state.auth_token:
            self.notify("Select a user first")
            return
        payload: dict[str, object] = {
            "username": self.admin_username_input.text().strip(),
            "display_name": self.admin_display_name_input.text().strip(),
            "role": self.admin_role_input.text().strip() or "user",
            "login_valid_until": self._admin_parse_datetime(self.admin_valid_until_input.text().strip()),
            "clear_device_binding": self.admin_clear_device_checkbox.isChecked(),
        }
        password = self.admin_new_password_input.text().strip()
        if password:
            payload["reset_password"] = password
        try:
            update_admin_user(self.state.auth_token, user_id, payload)
            self.notify("User updated")
            self._admin_refresh_users()
        except Exception as exc:
            self._log_action(f"Failed to save admin user: {exc}")
            self.notify("Unable to save user")

    def _admin_toggle_active(self, is_active: bool) -> None:
        user_id = self._admin_current_user_id()
        if user_id is None or not self.state.logged_in or not self.state.auth_token:
            return
        try:
            update_admin_user(self.state.auth_token, user_id, {"is_active": is_active})
            self._admin_refresh_users()
            self.notify("User activated" if is_active else "User deactivated")
        except Exception as exc:
            self._log_action(f"Failed to change user state: {exc}")

    def _admin_reset_password(self) -> None:
        user_id = self._admin_current_user_id()
        password = self.admin_new_password_input.text().strip()
        if user_id is None or not self.state.logged_in or not self.state.auth_token:
            return
        if not password:
            self.notify("Enter a new password first")
            return
        try:
            reset_admin_user_password(self.state.auth_token, user_id, password)
            self.notify("Password reset")
            self._admin_refresh_users()
        except Exception as exc:
            self._log_action(f"Failed to reset password: {exc}")

    def _admin_reset_device(self) -> None:
        user_id = self._admin_current_user_id()
        if user_id is None or not self.state.logged_in or not self.state.auth_token:
            return
        try:
            reset_admin_user_device(self.state.auth_token, user_id)
            self.notify("Device binding cleared")
            self._admin_refresh_users()
        except Exception as exc:
            self._log_action(f"Failed to reset device binding: {exc}")

    def _sync_admin_access(self) -> None:
        is_admin = self.state.logged_in and self.state.role == "admin"
        index = self.tabs.indexOf(self.admin_tabs)
        if index >= 0:
            self.tabs.setTabEnabled(index, is_admin)
        self.admin_access_label.setText(
            "Admin features are enabled for this account."
            if is_admin
            else "Admin features are disabled for this account."
        )
        for widget in (
            self.admin_users_table,
            self.admin_activity_table,
            self.admin_login_history_table,
            self.admin_refresh_users_button,
            self.admin_refresh_activity_button,
            self.admin_refresh_login_button,
            self.admin_create_user_button,
            self.admin_save_user_button,
            self.admin_activate_user_button,
            self.admin_deactivate_user_button,
            self.admin_reset_password_button,
            self.admin_reset_device_button,
            self.admin_selected_user_id,
            self.admin_username_input,
            self.admin_display_name_input,
            self.admin_role_input,
            self.admin_valid_until_input,
            self.admin_new_password_input,
            self.admin_clear_device_checkbox,
        ):
            widget.setEnabled(is_admin)

    def _build_tags_tab(self) -> QWidget:
        page = QWidget()
        layout = QVBoxLayout(page)
        layout.setSpacing(_scaled_int(10, self._scale))

        header = self._section_title(
            "DYNAMIC TAGS",
            "Use these tags in Subject, Body, attachment HTML, and custom file names. They generate random values when sending.",
        )
        layout.addWidget(header)

        grid_card, grid_layout = self._card("", None)
        grid_layout.setContentsMargins(_scaled_int(6, self._scale), _scaled_int(6, self._scale), _scaled_int(6, self._scale), _scaled_int(6, self._scale))
        tag_scroll = QScrollArea()
        tag_scroll.setWidgetResizable(True)
        tag_scroll.setFrameShape(QFrame.NoFrame)
        tag_scroll.setHorizontalScrollBarPolicy(Qt.ScrollBarAlwaysOff)
        tag_scroll.setVerticalScrollBarPolicy(Qt.ScrollBarAsNeeded)
        tag_scroll.setMaximumHeight(_scaled_int(430, self._scale))

        tag_host = QWidget()
        tag_host_layout = QGridLayout(tag_host)
        tag_host_layout.setContentsMargins(0, 0, 0, 0)
        tag_host_layout.setSpacing(_scaled_int(8, self._scale))
        self._tag_value_labels.clear()
        values = self._tag_sample_values()
        for index, definition in enumerate(self._tag_definitions):
            token = definition["token"]
            title = definition["title"]
            description = definition["description"]
            value = values.get(token, definition["default_value"])
            tag_host_layout.addWidget(self._tag_card(title, token, description, value), index // 3, index % 3)
        tag_scroll.setWidget(tag_host)
        grid_layout.addWidget(tag_scroll)
        layout.addWidget(grid_card)

        actions_row = QHBoxLayout()
        regenerate_button = QPushButton("Regenerate All")
        regenerate_button.setObjectName("secondaryButton")
        ai_button = QPushButton("Generate with AI")
        ai_button.setObjectName("warningButton")
        reset_button = QPushButton("Reset to Default")
        reset_button.setObjectName("secondaryButton")
        regenerate_button.setToolTip("Regenerate all sample tag values")
        ai_button.setToolTip("Generate new tag ideas with AI")
        reset_button.setToolTip("Restore the default tag set")
        regenerate_button.clicked.connect(lambda: self._regenerate_tags())
        ai_button.clicked.connect(lambda: self._log_action("Generated tags with AI"))
        reset_button.clicked.connect(lambda: self._reset_tags_to_default())
        self._apply_button_icon(regenerate_button, QStyle.SP_BrowserReload)
        self._apply_button_icon(ai_button, QStyle.SP_ComputerIcon)
        self._apply_button_icon(reset_button, QStyle.SP_DialogResetButton)
        actions_row.addWidget(regenerate_button)
        actions_row.addWidget(ai_button)
        actions_row.addWidget(reset_button)
        actions_row.addStretch()
        layout.addLayout(actions_row)

        manual_card, manual_layout = self._card("MANUAL CUSTOM TAGS", "Values entered here replace $custom1 and $custom2.")
        manual_layout.addWidget(QLabel("For example, you can add a phone number, email address, or any other text."))
        self.custom1_input.setPlaceholderText("$custom1 =")
        self.custom2_input.setPlaceholderText("$custom2 =")
        self.custom1_input.setToolTip("Define the first manual custom tag value")
        self.custom2_input.setToolTip("Define the second manual custom tag value")
        self.custom1_input.textEdited.connect(lambda _value: self._schedule_tags_save())
        self.custom2_input.textEdited.connect(lambda _value: self._schedule_tags_save())
        manual_layout.addWidget(self._line_with_copy(self.custom1_input))
        manual_layout.addWidget(self._line_with_copy(self.custom2_input))
        layout.addWidget(manual_card)
        self._load_tags_state()
        layout.addStretch()
        return self._tab_scroll(page)

    def _default_tag_definitions(self) -> list[dict[str, str]]:
        return [
            {"title": "VOIS", "token": "$random4", "description": "4 char alphanumeric uppercase", "default_value": "VOIS"},
            {"title": "V3EAJO", "token": "$random6", "description": "6 char alphanumeric uppercase", "default_value": "V3EAJO"},
            {"title": "0G2639Q", "token": "$random8", "description": "8 char alphanumeric uppercase", "default_value": "0G2639Q"},
            {"title": "I51VNI166P", "token": "$random10", "description": "10 char alphanumeric uppercase", "default_value": "I51VNI166P"},
            {"title": "I520Z0QQQ7CN", "token": "$random12", "description": "12 char alphanumeric uppercase", "default_value": "I520Z0QQQ7CN"},
            {"title": "VIK", "token": "$word3", "description": "3-letter uppercase word", "default_value": "VIK"},
            {"title": "POPE", "token": "$word4", "description": "4-letter word pattern", "default_value": "POPE"},
            {"title": "FABEQ", "token": "$word5", "description": "5-letter uppercase word", "default_value": "FABEQ"},
            {"title": "7575", "token": "$num4", "description": "4 digit number", "default_value": "7575"},
            {"title": "640296", "token": "$num6", "description": "6 digit number", "default_value": "640296"},
            {"title": "45250809", "token": "$num8", "description": "8 digit number", "default_value": "45250809"},
            {"title": "1-800-181-7889", "token": "$phone", "description": "Phone-style sample", "default_value": "1-800-181-7889"},
            {"title": "QLRZ", "token": "$word4a", "description": "Alternative 4-letter word", "default_value": "QLRZ"},
            {"title": "MOTION", "token": "$word6", "description": "6-letter uppercase word", "default_value": "MOTION"},
            {"title": "7A4F", "token": "$mix4", "description": "Mixed 4-character token", "default_value": "7A4F"},
            {"title": "DX8M2P", "token": "$mix6", "description": "Mixed 6-character token", "default_value": "DX8M2P"},
            {"title": "0314", "token": "$day4", "description": "4 digit day code", "default_value": "0314"},
            {"title": "202608", "token": "$ym6", "description": "Year-month code", "default_value": "202608"},
            {"title": "94-221-88", "token": "$id9", "description": "Structured numeric token", "default_value": "94-221-88"},
            {"title": "support@ezymailer.com", "token": "$email", "description": "Email address sample", "default_value": "support@ezymailer.com"},
            {"title": "https://ezymailer.app", "token": "$url", "description": "Website URL sample", "default_value": "https://ezymailer.app"},
            {"title": "Alice Johnson", "token": "$name", "description": "Full name sample", "default_value": "Alice Johnson"},
            {"title": "Seattle", "token": "$city", "description": "City sample", "default_value": "Seattle"},
            {"title": "hello-world", "token": "$slug", "description": "Slug sample", "default_value": "hello-world"},
        ]

    def _configure_segmented_button(self, button: QPushButton, checked: bool = False) -> None:
        button.setCheckable(True)
        button.setChecked(checked)
        button.setCursor(Qt.PointingHandCursor)

    def _standard_icon(self, pixmap: QStyle.StandardPixmap):
        return self.style().standardIcon(pixmap)

    def _apply_button_icon(self, button: QPushButton, pixmap: QStyle.StandardPixmap) -> None:
        button.setIcon(self._standard_icon(pixmap))
        button.setIconSize(QSize(_scaled_int(14, self._scale), _scaled_int(14, self._scale)))

    def _apply_tab_icon(self, tab_widget: QTabWidget, index: int, pixmap: QStyle.StandardPixmap) -> None:
        tab_widget.setTabIcon(index, self._standard_icon(pixmap))

    def _apply_navigation_icons(self) -> None:
        if hasattr(self, "tabs"):
            self.tabs.setIconSize(QSize(_scaled_int(16, self._scale), _scaled_int(16, self._scale)))
            self._apply_tab_icon(self.tabs, 0, QStyle.SP_DirIcon)
            self._apply_tab_icon(self.tabs, 1, QStyle.SP_FileIcon)
            self._apply_tab_icon(self.tabs, 2, QStyle.SP_FileDialogContentsView)
            self._apply_tab_icon(self.tabs, 3, QStyle.SP_FileDialogDetailedView)
            self._apply_tab_icon(self.tabs, 4, QStyle.SP_DialogApplyButton)
            self._apply_tab_icon(self.tabs, 5, QStyle.SP_MediaPlay)
        if hasattr(self, "body_tabs"):
            self.body_tabs.setIconSize(QSize(_scaled_int(14, self._scale), _scaled_int(14, self._scale)))
        if hasattr(self, "attach_tabs"):
            self.attach_tabs.setIconSize(QSize(_scaled_int(14, self._scale), _scaled_int(14, self._scale)))
        if hasattr(self, "admin_tabs"):
            self.admin_tabs.setIconSize(QSize(_scaled_int(16, self._scale), _scaled_int(16, self._scale)))
            self._apply_tab_icon(self.admin_tabs, 0, QStyle.SP_DirHomeIcon)
            self._apply_tab_icon(self.admin_tabs, 1, QStyle.SP_FileDialogListView)
            self._apply_tab_icon(self.admin_tabs, 2, QStyle.SP_MessageBoxInformation)
        for button, icon in (
            (getattr(self, "sidebar_start_campaign_button", None), QStyle.SP_MediaPlay),
            (getattr(self, "start_campaign_button", None), QStyle.SP_MediaPlay),
            (getattr(self, "campaign_pause_button", None), QStyle.SP_MediaPause),
            (getattr(self, "campaign_cancel_button", None), QStyle.SP_DialogCancelButton),
            (getattr(self, "campaign_reset_all_button", None), QStyle.SP_BrowserReload),
            (getattr(self, "launch_preset_label", None), None),
        ):
            if isinstance(button, QPushButton) and icon is not None:
                self._apply_button_icon(button, icon)

    def _labeled_value_row(self, label_text: str, widget: QWidget) -> QWidget:
        row = QFrame()
        row_layout = QHBoxLayout(row)
        row_layout.setContentsMargins(0, 0, 0, 0)
        row_layout.setSpacing(_scaled_int(10, self._scale))

        label = QLabel(label_text)
        label.setObjectName("fieldLabel")
        row_layout.addWidget(label)
        row_layout.addWidget(widget, 1)
        return row

    def _wrap_layout(self, layout: QHBoxLayout) -> QWidget:
        holder = QFrame()
        holder.setLayout(layout)
        return holder

    def _line_with_copy(self, line_edit: QLineEdit) -> QWidget:
        row = QFrame()
        row_layout = QHBoxLayout(row)
        row_layout.setContentsMargins(0, 0, 0, 0)
        row_layout.setSpacing(_scaled_int(10, self._scale))

        line_edit.setObjectName("searchInput")
        copy_button = QPushButton("Copy")
        copy_button.setObjectName("secondaryButton")
        token = line_edit.placeholderText().split("=", 1)[0].strip() or line_edit.text().strip()
        copy_button.clicked.connect(lambda _checked=False, token=token: self._copy_to_clipboard(token, "custom tag"))
        line_edit.setToolTip("Edit the manual custom tag value")
        copy_button.setToolTip("Copy the manual custom tag token")
        row_layout.addWidget(line_edit, 1)
        row_layout.addWidget(copy_button)
        return row

    def _copy_to_clipboard(self, text: str, label: str = "text") -> None:
        token = str(text or "").strip()
        if not token:
            return
        clipboard = QGuiApplication.clipboard()
        try:
            clipboard.clear(mode=QClipboard.Clipboard)
        except Exception:
            pass
        clipboard.setText(token, mode=QClipboard.Clipboard)
        try:
            clipboard.setText(token)
        except Exception:
            pass
        QApplication.processEvents()
        self._log_action(f"Copied {label}: {token}")
        self.notify(f"Copied {token}")

    def _tag_sample_values(self) -> dict[str, str]:
        values = dict(self.state.tag_samples or {})
        for definition in self._tag_definitions:
            token = definition["token"]
            if token not in values:
                values[token] = self._generate_random_tag_value(token, definition["default_value"])
        return values

    def _dynamic_tag_values(self) -> dict[str, str]:
        return {
            definition["token"]: self._generate_random_tag_value(definition["token"], definition["default_value"])
            for definition in self._tag_definitions
        }

    def _resolved_tag_values(self, recipient: str = "", subject: str = "", tag_values: dict[str, str] | None = None) -> dict[str, str]:
        replacements = {
            "$custom1": self.custom1_input.text().strip(),
            "$custom2": self.custom2_input.text().strip(),
            "$subject": subject.strip(),
        }
        replacements.update(tag_values or self._dynamic_tag_values())
        replacements.update(self._recipient_tag_values(recipient))
        replacements.update(self._customer_variable_values(recipient))
        return replacements

    def _generate_random_tag_value(self, token: str, default_value: str) -> str:
        token = token.strip()
        if token.startswith("$random"):
            length_text = token.removeprefix("$random")
            try:
                length = max(3, min(24, int(length_text)))
            except Exception:
                length = 8
            alphabet = string.ascii_uppercase + string.digits
            return "".join(secrets.choice(alphabet) for _ in range(length))
        if token.startswith("$word"):
            length_text = token.removeprefix("$word")
            try:
                length = max(3, min(16, int(length_text)))
            except Exception:
                length = 5
            alphabet = string.ascii_uppercase
            return "".join(secrets.choice(alphabet) for _ in range(length))
        if token.startswith("$num"):
            length_text = token.removeprefix("$num")
            try:
                length = max(2, min(12, int(length_text)))
            except Exception:
                length = 4
            return "".join(secrets.choice(string.digits) for _ in range(length))
        if token.startswith("$mix"):
            length_text = token.removeprefix("$mix")
            try:
                length = max(2, min(16, int(length_text)))
            except Exception:
                length = 6
            alphabet = string.ascii_uppercase + string.digits
            return "".join(secrets.choice(alphabet) for _ in range(length))
        if token.startswith("$day"):
            return "".join(secrets.choice(string.digits) for _ in range(4))
        if token.startswith("$ym"):
            today = QDateTime.currentDateTime().date()
            return f"{today.year()}{secrets.randbelow(12) + 1:02d}"
        if token.startswith("$id"):
            return f"{''.join(secrets.choice(string.digits) for _ in range(2))}-{''.join(secrets.choice(string.digits) for _ in range(3))}-{''.join(secrets.choice(string.digits) for _ in range(2))}"
        if token == "$phone":
            return f"1-{''.join(secrets.choice(string.digits) for _ in range(3))}-{''.join(secrets.choice(string.digits) for _ in range(3))}-{''.join(secrets.choice(string.digits) for _ in range(4))}"
        if token == "$email":
            return f"support{secrets.randbelow(9000) + 1000}@ezymailer.com"
        if token == "$url":
            return f"https://ezymailer-{secrets.randbelow(9000) + 1000}.app"
        if token == "$name":
            first_names = ["Alice", "Maya", "Jordan", "Sam", "Taylor", "Riley"]
            last_names = ["Johnson", "Patel", "Smith", "Brown", "Lee", "Walker"]
            return f"{secrets.choice(first_names)} {secrets.choice(last_names)}"
        if token == "$city":
            cities = ["Seattle", "Austin", "Denver", "Miami", "Chennai", "Berlin"]
            return secrets.choice(cities)
        if token == "$slug":
            words = ["hello", "launch", "email", "campaign", "update", "ready"]
            return f"{secrets.choice(words)}-{secrets.choice(words)}"
        return default_value

    def _current_tag_state(self) -> dict[str, object]:
        samples = self._tag_sample_values()
        return {
            "custom1": self.custom1_input.text().strip(),
            "custom2": self.custom2_input.text().strip(),
            "samples": samples,
        }

    def _apply_tag_state(self, payload: dict[str, object]) -> None:
        custom1 = str(payload.get("custom1") or "")
        custom2 = str(payload.get("custom2") or "")
        samples_raw = payload.get("samples") or {}
        samples: dict[str, str] = {}
        if isinstance(samples_raw, dict):
            for key, value in samples_raw.items():
                key_text = str(key).strip()
                value_text = str(value).strip()
                if key_text:
                    samples[key_text] = value_text
        self.custom1_input.blockSignals(True)
        self.custom2_input.blockSignals(True)
        try:
            self.custom1_input.setText(custom1)
            self.custom2_input.setText(custom2)
        finally:
            self.custom1_input.blockSignals(False)
            self.custom2_input.blockSignals(False)
        self.state.custom_tag_1 = custom1
        self.state.custom_tag_2 = custom2
        merged_samples = self._tag_sample_values()
        merged_samples.update(samples)
        self.state.tag_samples = merged_samples
        for token, label in self._tag_value_labels.items():
            label.setText(merged_samples.get(token, label.text()))

    def _schedule_tags_save(self) -> None:
        if self._workspace_loading:
            return
        self._tags_save_timer.start()

    def _persist_tags_state(self) -> None:
        if self._workspace_loading:
            return
        try:
            payload = self._current_tag_state()
            _upsert_tag_state(payload)
            self.state.custom_tag_1 = str(payload.get("custom1") or "")
            self.state.custom_tag_2 = str(payload.get("custom2") or "")
            self.state.tag_samples = dict(payload.get("samples") or {})
            if self.state.logged_in and self.state.auth_token:
                try:
                    api_save_tags(self.state.auth_token, payload)
                except Exception:
                    pass
        except Exception as exc:
            self._log_action(f"Failed to save tags: {exc}")

    def _load_tags_state(self) -> None:
        payload: dict[str, object] = {}
        if self.state.logged_in and self.state.auth_token:
            try:
                api_payload = api_get_tags(self.state.auth_token)
                if isinstance(api_payload, dict):
                    remote = api_payload.get("tag_state")
                    if isinstance(remote, dict):
                        payload = remote
            except Exception:
                payload = {}
        if not payload:
            payload = _load_tag_state()
        if not payload:
            payload = self._current_tag_state()
        self._apply_tag_state(payload)

    def _regenerate_tags(self) -> None:
        payload = self._current_tag_state()
        samples = {}
        for definition in self._tag_definitions:
            token = definition["token"]
            default_value = definition["default_value"]
            samples[token] = self._generate_random_tag_value(token, default_value)
        payload["samples"] = samples
        self._apply_tag_state(payload)
        self._persist_tags_state()
        self._log_action("Regenerated tag values")
        self.notify("Tags regenerated")

    def _reset_tags_to_default(self) -> None:
        payload = {
            "custom1": "",
            "custom2": "",
            "samples": self._dynamic_tag_values(),
        }
        self._apply_tag_state(payload)
        self._persist_tags_state()
        self._log_action("Randomized tag values")
        self.notify("Tags randomized")

    def _recipient_tag_values(self, recipient: str) -> dict[str, str]:
        recipient = (recipient or "").strip()
        if not recipient:
            return {}

        local_part = recipient.split("@", 1)[0].strip()
        if not local_part:
            return {
                "$email": recipient,
                "$recipient": recipient,
                "$recipient_email": recipient,
            }

        parts = [part for part in re.split(r"[._+-]+", local_part) if part]
        first_name = parts[0].strip().title() if parts else local_part.title()
        last_name = parts[1].strip().title() if len(parts) > 1 else ""
        full_name = " ".join(part.strip().title() for part in parts) if parts else local_part.title()
        return {
            "$email": recipient,
            "$recipient": recipient,
            "$recipient_email": recipient,
            "$recipient_local_part": local_part,
            "$username": local_part,
            "$first_name": first_name,
            "$last_name": last_name,
            "$full_name": full_name,
        }

    def _customer_variable_values(self, recipient: str) -> dict[str, str]:
        recipient = (recipient or "").strip().lower()
        if not recipient:
            return {}

        payload: dict[str, object] = {}
        if self.state.logged_in and self.state.auth_token:
            try:
                api_payload = api_get_customer_variables(self.state.auth_token, recipient)
                if isinstance(api_payload, dict):
                    items = api_payload.get("items")
                    if isinstance(items, list):
                        for item in items:
                            if not isinstance(item, dict):
                                continue
                            if str(item.get("email") or "").strip().lower() != recipient:
                                continue
                            variables = item.get("variables")
                            if isinstance(variables, dict):
                                payload.update(variables)
                                break
                    else:
                        variables = api_payload.get("variables")
                        if isinstance(variables, dict):
                            payload.update(variables)
            except Exception:
                pass

        if not payload:
            payload.update(_load_customer_variables(recipient))

        normalized: dict[str, str] = {}
        for key, value in payload.items():
            key_text = str(key).strip()
            value_text = str(value).strip()
            if not key_text or not value_text:
                continue
            normalized[key_text] = value_text
            normalized[key_text.lower()] = value_text
            if not key_text.startswith("$"):
                normalized[f"${key_text}"] = value_text
        return normalized

    def _apply_tags_to_text(
        self,
        text: str,
        recipient: str = "",
        subject: str = "",
        tag_values: dict[str, str] | None = None,
    ) -> str:
        result = text or ""
        replacements = self._resolved_tag_values(recipient, subject, tag_values)

        brace_lookup = {
            key.lstrip("$").lower(): value
            for key, value in replacements.items()
            if value and key
        }

        def _replace_brace(match: re.Match[str]) -> str:
            token = match.group(1).strip().lower()
            return brace_lookup.get(token, match.group(0))

        result = re.sub(r"\{\{\s*([A-Za-z0-9_]+)\s*\}\}", _replace_brace, result)

        for token, value in sorted(replacements.items(), key=lambda item: len(item[0]), reverse=True):
            if token.startswith("$") and value:
                result = result.replace(token, value)
        return result

    def _customer_variable_records(self) -> dict[str, dict[str, str]]:
        records: dict[str, dict[str, str]] = {}
        if self.state.logged_in and self.state.auth_token:
            try:
                api_payload = api_get_customer_variables(self.state.auth_token)
                if isinstance(api_payload, dict):
                    items = api_payload.get("items")
                    if isinstance(items, list):
                        for item in items:
                            if not isinstance(item, dict):
                                continue
                            email = str(item.get("email") or "").strip().lower()
                            variables = item.get("variables")
                            if not email or not isinstance(variables, dict):
                                continue
                            merged = records.setdefault(email, {})
                            for key, value in variables.items():
                                key_text = str(key).strip()
                                value_text = str(value).strip()
                                if key_text and value_text:
                                    merged[key_text] = value_text
            except Exception:
                pass

        if records:
            return records

        return {
            email: {str(key): str(value) for key, value in variables.items() if str(key).strip() and str(value).strip()}
            for email, variables in _load_all_customer_variables().items()
        }

    def _refresh_customer_variables_table(self) -> None:
        if not hasattr(self, "customer_variables_table"):
            return
        records = self._customer_variable_records()
        rows: list[tuple[str, str, str]] = []
        for email in sorted(records.keys()):
            variables = records[email]
            for key in sorted(variables.keys()):
                rows.append((email, key, variables[key]))

        self.customer_variables_table.blockSignals(True)
        try:
            self.customer_variables_table.setRowCount(0)
            for row_index, (email, key, value) in enumerate(rows):
                self.customer_variables_table.insertRow(row_index)
                for col, text in enumerate((email, key, value)):
                    item = QTableWidgetItem(text)
                    item.setFlags(Qt.ItemIsSelectable | Qt.ItemIsEnabled)
                    self.customer_variables_table.setItem(row_index, col, item)
        finally:
            self.customer_variables_table.blockSignals(False)
        if self.customer_variables_table.rowCount() > 0:
            self.customer_variables_table.selectRow(0)
            self._populate_customer_variable_form_from_selection()
        else:
            self._clear_customer_variable_form()

    def _populate_customer_variable_form_from_selection(self) -> None:
        row = self.customer_variables_table.currentRow() if hasattr(self, "customer_variables_table") else -1
        if row < 0:
            return
        email_item = self.customer_variables_table.item(row, 0)
        key_item = self.customer_variables_table.item(row, 1)
        value_item = self.customer_variables_table.item(row, 2)
        if email_item is None or key_item is None or value_item is None:
            return
        self.customer_email_input.setText(email_item.text().strip())
        self.customer_variable_key_input.setText(key_item.text().strip())
        self.customer_variable_value_input.setText(value_item.text().strip())

    def _clear_customer_variable_form(self) -> None:
        self.customer_email_input.clear()
        self.customer_variable_key_input.clear()
        self.customer_variable_value_input.clear()
        if hasattr(self, "customer_variables_table"):
            self.customer_variables_table.clearSelection()

    def _persist_customer_variable_record(self) -> None:
        email = self.customer_email_input.text().strip().lower()
        key = self.customer_variable_key_input.text().strip()
        value = self.customer_variable_value_input.text().strip()
        if not email or not key or not value:
            self.notify("Enter email, variable name, and value")
            return

        current = self._customer_variable_records().get(email, {})
        current[key] = value
        _upsert_customer_variables(email, current)
        if self.state.logged_in and self.state.auth_token:
            try:
                api_save_customer_variables(self.state.auth_token, email, current)
            except Exception:
                pass
        self._refresh_customer_variables_table()
        self._log_action(f"Saved customer variable {key} for {email}")
        self.notify("Customer variable saved")

    def _delete_selected_customer_variable_record(self) -> None:
        row = self.customer_variables_table.currentRow()
        if row < 0:
            self.notify("Select a customer variable first")
            return
        email_item = self.customer_variables_table.item(row, 0)
        key_item = self.customer_variables_table.item(row, 1)
        if email_item is None or key_item is None:
            return
        email = email_item.text().strip().lower()
        key = key_item.text().strip()
        records = self._customer_variable_records()
        current = records.get(email, {})
        if key in current:
            current.pop(key, None)
        if current:
            _upsert_customer_variables(email, current)
            if self.state.logged_in and self.state.auth_token:
                try:
                    api_save_customer_variables(self.state.auth_token, email, current)
                except Exception:
                    pass
        else:
            _delete_customer_variables(email)
            if self.state.logged_in and self.state.auth_token:
                try:
                    api_delete_customer_variables(self.state.auth_token, email)
                except Exception:
                    pass
        self._refresh_customer_variables_table()
        self._clear_customer_variable_form()
        self._log_action(f"Deleted customer variable {key} for {email}")
        self.notify("Customer variable deleted")

    def _load_customer_variables_state(self) -> None:
        self._refresh_customer_variables_table()

    def _open_output_options_dialog(self, preview: bool = False) -> None:
        dialog = OutputOptionsDialog(self.window(), scale=self._scale)
        if dialog.exec() == QDialog.Accepted:
            self._log_action("Opened export options")
            if preview:
                self._log_action("Exported HTML and opened preview")
                self.notify("HTML export preview ready")
                self._preview_html_content(title="Converted HTML Preview")
            else:
                self.notify("Export options confirmed")

    def _choose_attachment_file_format(self) -> None:
        dialog = FileFormatDialog(
            self.window(),
            scale=self._scale,
            selected_format=getattr(self, "attach_format_value", "PDF document") or "PDF document",
            file_name_mode=getattr(self, "attach_file_name_mode", "auto"),
            file_name_value=getattr(self, "attach_file_name_value", ""),
        )
        if dialog.exec() == QDialog.Accepted:
            selected_formats = self._normalize_attachment_format_values(dialog.selected_formats())
            if dialog._random_checkbox is not None and dialog._random_checkbox.isChecked():
                self.attach_format_value = "Random format"
            else:
                self.attach_format_value = ", ".join(selected_formats) if selected_formats else "PDF document"
            self.attach_file_name_mode = "custom" if dialog.uses_custom_file_name() else "auto"
            self.attach_file_name_value = dialog.custom_name_input.text().strip() if dialog.uses_custom_file_name() and dialog.custom_name_input is not None else ""
            self.attach_format_label.setText(self._attachment_format_summary(self.attach_format_value))
            self._log_action(f"Selected attachment format: {self.attach_format_value}")
            self.notify(f"Format set to {self.attach_format_value}")

    def _sync_attachment_convert_controls(self, checked: bool) -> None:
        row_widget = getattr(self, "attach_format_row_widget", None)
        row_label = getattr(self, "attach_format_row_label", None)
        if row_widget is not None:
            row_widget.setVisible(checked)
        if row_label is not None:
            row_label.setVisible(checked)
        self.attach_choose_format_button.setEnabled(checked)
        self.attach_format_label.setEnabled(checked)

    def _attachment_format_summary(self, selected: str | list[str]) -> str:
        summary_map = {
            "PDF document": "PDF",
            "Excel spreadsheet (XLSX)": "XLSX",
            "Excel template (XLTX)": "XLTX",
            "PowerPoint presentation (PPTX)": "PPTX",
            "PowerPoint slideshow (PPSX)": "PPSX",
            "Word document (DOCX)": "DOCX",
        }
        if isinstance(selected, list):
            parts = [str(item).strip() for item in selected if str(item).strip()]
        else:
            selected = (selected or "").strip()
            if not selected:
                return "PDF"
            if selected == "Random format":
                return "Random"
            parts = [part.strip() for part in selected.split(",") if part.strip()]
        if not parts:
            return "PDF"
        if len(parts) == 1:
            return summary_map.get(parts[0], parts[0])
        short_parts = [summary_map.get(part, part) for part in parts]
        return f"Multiple: {', '.join(short_parts)}"

    def _normalize_attachment_format_values(self, selected: str | list[str] | tuple[str, ...]) -> list[str]:
        if isinstance(selected, (list, tuple)):
            parts = [str(item).strip() for item in selected if str(item).strip()]
        else:
            selected = (selected or "").strip()
            if not selected:
                return ["PDF document"]
            if selected == "Random format":
                return ["Random format"]
            parts = [part.strip() for part in selected.split(",") if part.strip()]
        mapping = {
            "PDF document": "PDF document",
            "Excel spreadsheet (XLSX)": "Excel spreadsheet (XLSX)",
            "Excel template (XLTX)": "Excel template (XLTX)",
            "PowerPoint presentation (PPTX)": "PowerPoint presentation (PPTX)",
            "PowerPoint slideshow (PPSX)": "PowerPoint slideshow (PPSX)",
            "Word document (DOCX)": "Word document (DOCX)",
        }
        normalized = [mapping.get(part, part) for part in parts if mapping.get(part, part)]
        return normalized or ["PDF document"]

    def _normalize_attachment_format_value(self, selected: str) -> str:
        normalized = self._normalize_attachment_format_values(selected)
        if not normalized:
            return "PDF document"
        if normalized == ["Random format"]:
            return "Random format"
        return normalized[0]

    def _attachment_format_extension(self, format_value: str) -> str:
        format_value = self._normalize_attachment_format_value(format_value)
        mapping = {
            "PDF document": ".pdf",
            "Excel spreadsheet (XLSX)": ".xlsx",
            "Excel template (XLTX)": ".xltx",
            "PowerPoint presentation (PPTX)": ".pptx",
            "PowerPoint slideshow (PPSX)": ".ppsx",
            "Word document (DOCX)": ".docx",
        }
        return mapping.get(format_value, ".out")

    def _attachment_default_name_base(self, recipient: str, subject: str) -> str:
        alphabet = string.ascii_uppercase + string.digits
        target_length = secrets.randbelow(9) + 8
        segments = min(max(secrets.randbelow(3) + 2, 2), max(target_length // 2, 2))
        lengths = [2] * segments
        remaining = target_length - (2 * segments)
        for _ in range(remaining):
            lengths[secrets.randbelow(segments)] += 1
        parts = ["".join(secrets.choice(alphabet) for _ in range(length)) for length in lengths]
        return "-".join(parts)

    def _attachment_output_name_base(
        self,
        recipient: str,
        subject: str,
        format_value: str,
        file_name_mode: str = "auto",
        file_name_value: str = "",
        resolve_tags: bool = True,
    ) -> str:
        default_base = self._attachment_default_name_base(recipient, subject)
        if str(file_name_mode or "").strip().lower() == "custom":
            raw_value = self._apply_tags_to_text(file_name_value or "", recipient, subject) if resolve_tags else (file_name_value or "")
            custom = self._sanitize_filename(raw_value)
            if custom:
                return custom
        return default_base

    def _wrap_text_as_html(self, text: str, subject: str = "") -> str:
        safe_subject = html.escape(subject.strip())
        safe_text = html.escape(text).replace("\n", "<br>")
        return f"""
        <html>
          <body style="margin:0; padding:24px; background:#1e1e1e; color:#d4d4d4; font-family:'Segoe UI Variable Text','Segoe UI',sans-serif;">
            <div style="max-width: 820px; margin: 0 auto;">
              {f'<div style="color:#569cd6; font-size:12px; letter-spacing:.4px; font-weight:700; text-transform:uppercase; margin-bottom:10px;">{safe_subject}</div>' if safe_subject else ''}
              <div style="white-space:normal; line-height:1.7; font-size:14px;">{safe_text}</div>
            </div>
          </body>
        </html>
        """

    def _build_preview_dialog(self, title: str, html_content: str, source_label: str) -> HtmlPreviewDialog:
        dialog = HtmlPreviewDialog(self.window(), title, html_content, source_label, scale=self._scale)
        self._floating_windows.append(dialog)
        dialog.finished.connect(lambda _result, d=dialog: self._remove_floating_window(d))
        dialog.destroyed.connect(lambda *_args, d=dialog: self._remove_floating_window(d))
        return dialog

    def _remove_floating_window(self, dialog: QDialog) -> None:
        if dialog in self._floating_windows:
            self._floating_windows.remove(dialog)

    def _preview_subject_body(self) -> None:
        subject = self._apply_tags_to_text(self.subject_input.text().strip()) or "Subject Preview"
        current_body = self._current_body_widget()
        if current_body is not None:
            body_payload = current_body.payload()
            if body_payload["mode"] == "HTML Message":
                html_content = self._apply_tags_to_text(body_payload["html_text"].strip())
                source = "Previewing the HTML message content."
                if not html_content:
                    html_content = "<html><body style='background:#1e1e1e; color:#d4d4d4; font-family:Segoe UI;'>No HTML content available.</body></html>"
            else:
                body_text = self._apply_tags_to_text(body_payload["plain_text"].strip())
                source = "Previewing the plain-text message as HTML."
                if not body_text:
                    body_text = "No message body available."
                html_content = self._wrap_text_as_html(body_text, subject)
        else:
            html_content = "<html><body style='background:#1e1e1e; color:#d4d4d4; font-family:Segoe UI;'>No message body available.</body></html>"
            source = "No active body is available."

        dialog = self._build_preview_dialog("Message Preview", html_content, source)
        dialog.show()
        dialog.raise_()
        dialog.activateWindow()
        self._log_action("Opened message preview")
        self.notify("Message preview opened")

    def _preview_body_editor_html(self, widget: BodyDraftEditor | None) -> None:
        if widget is None or widget.mode_text() != "HTML Message":
            self.notify("Switch to HTML body first")
            return

        html_content = self._apply_tags_to_text(widget.html_editor.toPlainText().strip())
        if not html_content:
            self.notify("Add HTML content first")
            return

        title = widget.title_text() or "HTML Body Preview"
        dialog = self._build_preview_dialog(title, html_content, "Previewing the selected HTML body.")
        dialog.show()
        dialog.raise_()
        dialog.activateWindow()
        self._log_action(f"Opened body preview: {title}")
        self.notify("Body preview opened")

    def _preview_html_content(self, title: str = "HTML Preview") -> None:
        current_widget = self._current_attachment_widget()
        html_content = ""
        if current_widget is not None:
            html_content = self._apply_tags_to_text(current_widget.html_editor.toPlainText().strip())
        elif hasattr(self, "html_editor") and isinstance(self.html_editor, QTextEdit):
            html_content = self._apply_tags_to_text(self.html_editor.toPlainText().strip())
        if not html_content:
            html_content = "<html><body style='background:#1e1e1e; color:#d4d4d4; font-family:Segoe UI;'>No HTML template available.</body></html>"
        dialog = self._build_preview_dialog(
            title,
            html_content,
            "Previewing the HTML template in a separate window.",
        )
        dialog.show()
        dialog.raise_()
        dialog.activateWindow()
        self._log_action(f"Opened {title.lower()}")
        self.notify("HTML preview opened")

    def _tag_card(self, title: str, token: str, description: str, value: str) -> QWidget:
        card = QFrame()
        card.setObjectName("panelCard")
        layout = QVBoxLayout(card)
        layout.setContentsMargins(_scaled_int(12, self._scale), _scaled_int(12, self._scale), _scaled_int(12, self._scale), _scaled_int(12, self._scale))
        layout.setSpacing(_scaled_int(6, self._scale))

        header = QHBoxLayout()
        title_label = QLabel(title)
        title_label.setObjectName("sectionTitle")
        token_label = QLabel(token)
        token_label.setObjectName("windowPill")
        header.addWidget(title_label)
        header.addStretch()
        header.addWidget(token_label)
        layout.addLayout(header)

        row = QHBoxLayout()
        token_value = QLabel(value or token)
        token_value.setObjectName("sectionSubtitle")
        copy_button = QPushButton("Copy")
        copy_button.setObjectName("secondaryButton")
        copy_button.clicked.connect(lambda _checked=False, token=token: self._copy_to_clipboard(token, "tag"))
        token_value.setToolTip(f"Token value for {token}")
        copy_button.setToolTip(f"Copy {token} to clipboard")
        row.addWidget(token_value)
        row.addStretch()
        row.addWidget(copy_button)
        layout.addLayout(row)

        desc = QLabel(description)
        desc.setObjectName("sectionHint")
        desc.setWordWrap(True)
        layout.addWidget(desc)

        self._tag_value_labels[token] = token_value

        return card

    def _set_browser_mode(self, mode: str) -> None:
        self.state.browser_mode = mode
        self.incognito_button.setChecked(mode == "Incognito")
        self.normal_button.setChecked(mode == "Normal")
        self._persist_browser_state()
        self._log_action(f"Browser mode set to {mode}")
        self.notify(f"Browser mode changed to {mode}")

    def _set_body_mode(self, mode: str) -> None:
        self.state.body_mode = mode
        self.normal_message_button.setChecked(mode == "Normal Message")
        self.html_message_button.setChecked(mode == "HTML Message")
        current_body = self._current_body_widget()
        if current_body is not None:
            current_body.set_mode(mode)
        if self.state.username:
            try:
                upsert_setting(self.state.username, "body_mode", mode)
            except Exception:
                pass
        self._schedule_subject_body_save()
        label = "Plain Text" if mode == "Normal Message" else "HTML Body"
        self._log_action(f"Body mode set to {label}")
        self.notify(f"Body mode changed to {label}")

    def _set_launch_preset(self, preset: str | None) -> None:
        self.state.launch_preset = preset or ""
        label = preset if preset else "None"
        self.launch_preset_label.setText(label)
        self._persist_browser_state()
        self._log_action(f"Launch preset set to {label}")
        self.notify(f"Launch preset updated: {label}")

    def _current_browser_state_payload(self) -> dict[str, object]:
        return {
            "browser_mode": self.state.browser_mode,
            "launch_preset": self.state.launch_preset,
            "window_count": int(self.state.window_count),
        }

    def _apply_browser_state_payload(self, payload: dict[str, object]) -> None:
        browser_mode = str(payload.get("browser_mode") or "Incognito")
        launch_preset = str(payload.get("launch_preset") or "Default")
        try:
            window_count = max(1, int(payload.get("window_count") or 1))
        except Exception:
            window_count = 1

        self.state.browser_mode = browser_mode if browser_mode in {"Incognito", "Normal"} else "Incognito"
        self.state.launch_preset = launch_preset
        self.state.window_count = window_count

    def _load_browser_state(self) -> None:
        payload = _load_ui_state(LOCAL_BROWSER_STATE_KEY)
        if not isinstance(payload, dict):
            payload = {}
        self._apply_browser_state_payload(payload)

    def _persist_browser_state(self) -> None:
        if self._workspace_loading:
            return
        _upsert_ui_state(LOCAL_BROWSER_STATE_KEY, self._current_browser_state_payload())

    def _window_count_changed(self, value: int) -> None:
        self.state.window_count = max(1, value)

    def _browser_binary(self) -> Path | None:
        env_binary = os.getenv("EZYM_MAILER_BROWSER_BINARY", "").strip()
        if env_binary:
            candidate = Path(env_binary).expanduser()
            if candidate.exists():
                return candidate
        return _installed_browser_binary() or _cached_browser_binary(_browser_cache_dir())

    def _browser_launch_rect(self, index: int, total: int) -> tuple[int, int, int, int]:
        screen = self.window().screen() or QApplication.primaryScreen()
        if screen is None:
            return (80, 80, 1280, 800)

        geometry = screen.availableGeometry()
        margin = _scaled_int(18, self._scale)
        gap = _scaled_int(14, self._scale)
        cols = max(1, ceil(sqrt(total)))
        rows = max(1, ceil(total / cols))

        usable_width = max(800, geometry.width() - (margin * 2) - (gap * (cols - 1)))
        usable_height = max(600, geometry.height() - (margin * 2) - (gap * (rows - 1)))
        cell_width = max(360, usable_width // cols)
        cell_height = max(280, usable_height // rows)

        row = (index - 1) // cols
        col = (index - 1) % cols
        x = geometry.x() + margin + (col * (cell_width + gap))
        y = geometry.y() + margin + (row * (cell_height + gap))
        return x, y, cell_width, cell_height

    def _create_browser_profile_dir(self, index: int) -> Path:
        prefix = f"ezymailer-{self.state.username or 'guest'}-{index}-"
        profile_dir = Path(tempfile.mkdtemp(prefix=prefix))
        self._seed_browser_profile_dir(profile_dir)
        return profile_dir

    def _reserve_debug_port(self) -> int:
        with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as sock:
            sock.bind(("127.0.0.1", 0))
            return int(sock.getsockname()[1])

    def _seed_browser_profile_dir(self, profile_dir: Path) -> None:
        """Pre-populate a fresh Chrome profile to suppress first-run prompts."""
        try:
            profile_dir.mkdir(parents=True, exist_ok=True)
            (profile_dir / "First Run").touch(exist_ok=True)
            local_state_path = profile_dir / "Local State"
            if not local_state_path.exists():
                local_state_path.write_text(
                    json.dumps(
                        {
                            "user_experience_metrics": {
                                "reporting_enabled": False,
                            },
                            "browser": {
                                "check_default_browser": False,
                            },
                        },
                        indent=2,
                    ),
                    encoding="utf-8",
                )
        except Exception:
            pass

    def _cleanup_browser_profile_dir(self, profile_dir: Path | None) -> None:
        if profile_dir is None:
            return
        try:
            shutil.rmtree(profile_dir, ignore_errors=True)
        except Exception:
            pass

    def _launch_browser_process(self, index: int) -> BrowserSessionHandle:
        binary = self._browser_binary()
        if binary is None:
            raise RuntimeError("No supported Edge, Chrome, or Chromium browser was found for this app.")

        incognito = self.state.browser_mode == "Incognito"
        browser_name = _browser_product_name(binary)
        browser_slug = browser_name.lower().replace(" ", "-")
        session_id = f"{browser_slug}-{QDateTime.currentMSecsSinceEpoch()}-{index}"
        title = f"{browser_name} Window {index}"
        profile_dir = self._create_browser_profile_dir(index)
        debug_port = self._reserve_debug_port()
        args = [str(binary), "--new-window", f"--user-data-dir={profile_dir}"]
        # Suppress Chrome's first-run welcome dialog and default-browser prompt.
        args.extend(
            [
                "--no-first-run",
                "--no-default-browser-check",
                "--disable-infobars",
                "--disable-default-apps",
                "--disable-features=ChromeWhatsNewUI",
                f"--remote-debugging-port={debug_port}",
                "--remote-allow-origins=*",
            ]
        )
        if incognito:
            args.append(_browser_private_flag(browser_name))
        x, y, width, height = self._browser_launch_rect(index, max(1, self.window_spin.value()))
        args.append(f"--window-position={x},{y}")
        args.append(f"--window-size={width},{height}")
        args.append("https://www.google.com/")
        process = subprocess.Popen(
            args,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
            start_new_session=True,
        )
        return BrowserSessionHandle(
            session_id=session_id,
            title=title,
            mode="Incognito" if incognito else "Normal",
            browser_name=browser_name,
            process=process,
            status="Running",
            profile_dir=profile_dir,
            debug_port=debug_port,
        )

    def _sync_browser_session_states(self) -> None:
        removed_sessions: list[BrowserSessionHandle] = []
        alive_sessions: list[BrowserSessionHandle] = []
        for session in self._browser_sessions:
            if self._browser_session_is_alive(session):
                if session.status != "Paused":
                    session.status = "Running"
                alive_sessions.append(session)
                continue
            removed_sessions.append(session)

        if removed_sessions:
            self._browser_sessions = alive_sessions
            self._sync_session_state_from_handles()
            self._refresh_sessions()
            for session in removed_sessions:
                try:
                    _upsert_local_browser_session(
                        self.state.username,
                        session.session_id,
                        session.title,
                        session.browser_name,
                        self.state.browser_mode,
                        "Closed",
                        None,
                        self.state.launch_preset or "Default",
                        {
                            "browser_mode": self.state.browser_mode,
                            "profile_dir": str(session.profile_dir) if session.profile_dir else "",
                        },
                    )
                except Exception:
                    pass
                self._cleanup_browser_profile_dir(session.profile_dir)
                self._log_action(f"Browser window closed: {session.title}")
        else:
            self._sync_session_state_from_handles()

        if not self._browser_sessions:
            self._browser_watch_timer.stop()

    def _browser_session_is_alive(self, session: BrowserSessionHandle) -> bool:
        """Treat a browser as alive when its CDP endpoint is still serving.

        On Windows, Edge may hand the window to another process and exit the
        Popen handle. The debugging endpoint is the reliable signal because it
        is also what campaign workers use to connect to the existing window.
        """
        if session.process is not None and session.process.poll() is None:
            return True
        if session.debug_port is None:
            return False
        try:
            with socket.create_connection(("127.0.0.1", session.debug_port), timeout=0.25):
                return True
        except OSError:
            return False

    def _usable_browser_sessions(self) -> list[BrowserSessionHandle]:
        self._sync_browser_session_states()
        return [
            session
            for session in self._browser_sessions
            if session.debug_port is not None and self._browser_session_is_alive(session)
        ]

    def _sync_session_state_from_handles(self) -> None:
        self.state.active_sessions = [
            f"{session.title} - {session.status}" for session in self._browser_sessions
        ]
        self.state.window_count = len(self._browser_sessions)
        self.active_windows_value.setText(str(self.state.window_count))

    def _terminate_browser_sessions(self, log_reason: str | None = None) -> None:
        for session in self._browser_sessions:
            process = session.process
            if process is not None and process.poll() is None:
                try:
                    process.terminate()
                    process.wait(timeout=3)
                except Exception:
                    try:
                        process.kill()
                    except Exception:
                        pass
            try:
                _upsert_local_browser_session(
                    self.state.username,
                    session.session_id,
                    session.title,
                    session.browser_name,
                    session.mode,
                    "Closed",
                    None,
                    self.state.launch_preset or "Default",
                    {
                        "browser_mode": self.state.browser_mode,
                        "profile_dir": str(session.profile_dir) if session.profile_dir else "",
                    },
                )
            except Exception:
                pass
            self._cleanup_browser_profile_dir(session.profile_dir)
            if log_reason:
                self._log_action(f"{log_reason}: {session.title}")
        self._browser_sessions.clear()
        self._sync_session_state_from_handles()
        self._refresh_sessions()
        self._browser_watch_timer.stop()

    def _handle_launch(self) -> None:
        title = "Confirm Launch"
        prompt = (
            f"Launch {self.window_spin.value()} browser window(s) using {self.state.browser_mode} mode "
            f"and the {self.launch_preset_label.text()} preset?"
        )
        confirm = ConfirmDialog(self.window(), title, prompt, scale=self._scale)
        if confirm.exec() != QDialog.Accepted:
            self.notify("Launch cancelled")
            return
        target = max(1, self.window_spin.value())
        self._terminate_browser_sessions()
        if self._browser_binary() is None:
            self._pending_launch_target = target
            self._log_action("No installed browser found; downloading Chromium before launch")
            self.notify("Preparing browser runtime")
            try:
                self.window()._start_windows_browser_bootstrap()
            except Exception:
                pass
            return
        self._log_action(f"Preparing {target} browser window(s)")
        self.notify(f"Launching {target} browser window(s)")
        self._show_launch_loader(
            "Launching browser windows",
            "Applying browser mode and launch preset.",
        )
        QTimer.singleShot(900, lambda t=target: self._complete_launch(t))

    def _resume_pending_launch(self) -> None:
        target = self._pending_launch_target
        if target is None:
            return
        self._pending_launch_target = None
        self._log_action(f"Browser runtime ready; launching {target} browser window(s)")
        self.notify(f"Launching {target} browser window(s)")
        self._show_launch_loader(
            "Launching browser windows",
            "Applying browser mode and launch preset.",
        )
        QTimer.singleShot(900, lambda t=target: self._complete_launch(t))

    def _handle_pause(self) -> None:
        if not self._browser_sessions:
            self.notify("No browser windows are currently open")
            return
        for session in self._browser_sessions:
            session.status = "Paused"
        self._sync_session_state_from_handles()
        self._refresh_sessions()
        self._browser_watch_timer.start()
        for session in self._browser_sessions:
            try:
                _upsert_local_browser_session(
                    self.state.username,
                    session.session_id,
                    session.title,
                    session.browser_name,
                    session.mode,
                    session.status,
                    session.process.pid if session.process is not None else None,
                    self.state.launch_preset or "Default",
                    {
                        "browser_mode": self.state.browser_mode,
                        "profile_dir": str(session.profile_dir) if session.profile_dir else "",
                    },
                )
            except Exception:
                pass
        self._log_action("Paused active browser sessions")
        self.notify("Sessions paused")

    def _handle_stop(self) -> None:
        self._handle_pause()

    def _handle_reset(self) -> None:
        self._terminate_browser_sessions()
        self.state.window_count = 1
        self.state.browser_mode = "Incognito"
        self.state.body_mode = "Normal Message"
        self.state.launch_preset = "Default"
        self.state.ai_provider = "ChatGPT"
        self.state.ai_api_key = ""
        self.state.ai_model = ""
        self.state.ai_connected = False
        self.state.ai_available_models = []
        self._available_ai_models = []
        self.window_spin.setValue(1)
        self.ai_provider_combo.blockSignals(True)
        self.ai_provider_combo.setCurrentText("ChatGPT")
        self.ai_provider_combo.blockSignals(False)
        self.ai_api_key_input.clear()
        self._persist_browser_state()
        self._refresh_controls()
        self._log_action("Reset workspace to defaults")
        self.notify("Workspace reset to defaults")

    def _current_campaign_payload(self) -> dict[str, object]:
        recipients = self._extract_email_candidates(self.pending_emails_editor.toPlainText())
        subject = self.subject_input.text().strip()
        subjects = [
            self._subject_item_subject(self.subject_drafts_list.item(index)).strip()
            for index in range(self.subject_drafts_list.count())
        ]
        subjects = [item for item in subjects if item]
        if subject and subject not in subjects:
            subjects.insert(0, subject)
        current_body = self._current_body_widget()
        body_text = ""
        body_html = ""
        if current_body is not None:
            payload = current_body.payload()
            body_text = str(payload.get("plain_text") or "").strip()
            body_html = str(payload.get("html_text") or "").strip()
        attachment_widget = self._current_attachment_widget()
        attachment_html = attachment_widget.html_editor.toPlainText().strip() if attachment_widget is not None else ""
        return {
            "recipients": recipients,
            "subject": subject,
            "subjects": subjects,
            "body_text": body_text,
            "body_html": body_html,
            "attachment_html": attachment_html,
            "attachment_format": self.attach_format_value,
            "attachment_formats": self._normalize_attachment_format_values(self.attach_format_value),
            "attachment_file_name_mode": self.attach_file_name_mode,
            "attachment_file_name_value": self.attach_file_name_value,
        }

    def _campaign_missing_fields(self, payload: dict[str, object] | None = None) -> list[str]:
        payload = payload or self._current_campaign_payload()
        missing: list[str] = []
        recipients = payload.get("recipients") or []
        if not isinstance(recipients, list) or not recipients:
            missing.append("customer emails")
        if not str(payload.get("subject") or "").strip():
            missing.append("subject")
        if not (str(payload.get("body_text") or "").strip() or str(payload.get("body_html") or "").strip()):
            missing.append("body content")
        if not str(payload.get("attachment_html") or "").strip():
            missing.append("attachment content")
        return missing

    def _refresh_campaign_action_state(self) -> None:
        self._update_campaign_action_state()

    def _reset_campaign_form_state(self, confirm: bool = True) -> None:
        if confirm:
            reply = QMessageBox.question(
                self,
                "Reset All",
                "Clear customer emails, subjects + body, and attachment content?",
                QMessageBox.Yes | QMessageBox.No,
                QMessageBox.No,
            )
            if reply != QMessageBox.Yes:
                return

        self._workspace_loading = True
        try:
            if self._campaign_active:
                self._handle_campaign_cancel()
            self._pending_campaign_payload = None
            self.progress_bar.setValue(0)
            self.campaign_progress_text.setText("0 / 0 sent")
            self._campaign_send_log_entries.clear()
            self._campaign_total = 0
            self._campaign_completed = 0
            self._delete_campaign_workspace_state()
            self._clear_pending_emails()
            self.subject_drafts_list.blockSignals(True)
            try:
                self.subject_drafts_list.clear()
            finally:
                self.subject_drafts_list.blockSignals(False)
            self.subject_input.blockSignals(True)
            try:
                self.subject_input.clear()
            finally:
                self.subject_input.blockSignals(False)
            self.state.subject_text = ""
            self._clear_subject_selection()
            self._load_body_tabs_from_state()
            self._clear_subject_body()
            self._load_attachment_tabs_from_local()
            self._refresh_sessions()
            self._refresh_controls()
        finally:
            self._workspace_loading = False
        self._log_action("Reset campaign workspace to defaults")
        self.notify("Campaign form reset to defaults")
        self._refresh_campaign_action_state()

    def _delete_campaign_workspace_state(self) -> None:
        try:
            _delete_ui_state(LOCAL_PENDING_EMAILS_STATE_KEY)
            _delete_ui_state(LOCAL_SUBJECT_STATE_KEY)
            _delete_ui_state(LOCAL_BODY_STATE_KEY)
            _delete_attachment_state()
            _delete_local_drafts("attachment")
        except Exception as exc:
            self._log_action(f"Failed to clear campaign workspace state: {exc}")

    def _close_session(self, session_id: str) -> None:
        session = next((item for item in self._browser_sessions if item.session_id == session_id), None)
        if session is None:
            return
        process = session.process
        if process is not None and process.poll() is None:
            try:
                process.terminate()
                process.wait(timeout=3)
            except Exception:
                try:
                    process.kill()
                except Exception:
                    pass
        elif process is not None and sys.platform.startswith("win"):
            # Edge/Chrome can hand the window to a child process before the
            # launcher exits; terminate that process tree as well.
            try:
                subprocess.run(
                    ["taskkill", "/PID", str(process.pid), "/T", "/F"],
                    check=False,
                    stdout=subprocess.DEVNULL,
                    stderr=subprocess.DEVNULL,
                )
            except Exception:
                pass
        self._browser_sessions = [item for item in self._browser_sessions if item.session_id != session_id]
        self._sync_session_state_from_handles()
        self._refresh_sessions()
        try:
            _upsert_local_browser_session(
                self.state.username,
                session.session_id,
                session.title,
                session.browser_name,
                session.mode,
                "Closed",
                None,
                self.state.launch_preset or "Default",
                {
                    "browser_mode": self.state.browser_mode,
                    "profile_dir": str(session.profile_dir) if session.profile_dir else "",
                },
            )
        except Exception:
            pass
        self._cleanup_browser_profile_dir(session.profile_dir)
        self._log_action(f"Closed browser window {session.title}")
        self.notify(f"Closed {session.title}")

    def _append_campaign_send_log(self, message: str) -> None:
        self._campaign_send_log_entries.append(message)
        self._refresh_activity()

    def _update_campaign_action_state(self) -> None:
        can_start = self.state.logged_in and not self._campaign_active and not self._campaign_missing_fields()
        primary_label = "Start Campaign"
        if self._campaign_active:
            primary_label = "Resume Campaign" if self._campaign_paused else "Pause Campaign"
        for button in (self.start_campaign_button, self.sidebar_start_campaign_button):
            button.setText(primary_label)
            button.setEnabled(can_start if not self._campaign_active else True)
        self.campaign_pause_button.setVisible(self._campaign_active)
        self.campaign_pause_button.setEnabled(self._campaign_active)
        self.campaign_pause_button.setText("Resume Campaign" if self._campaign_paused else "Pause Campaign")
        self.campaign_cancel_button.setVisible(self._campaign_active)
        self.campaign_cancel_button.setEnabled(self._campaign_active)
        if self._campaign_total:
            remaining = max(self._campaign_total - self._campaign_completed, 0)
            self.campaign_progress_text.setText(f"{self._campaign_completed} / {self._campaign_total} sent • {remaining} remaining")
        else:
            self.campaign_progress_text.setText("0 / 0 sent")

    def _chunk_campaign_recipients(self, recipients: list[str], window_count: int, sender_limit: int) -> list[list[str]]:
        if not recipients or window_count <= 0:
            return []
        total = len(recipients)
        if total > window_count * sender_limit:
            raise RuntimeError(
                f"{total} recipient(s) exceed the current window capacity of {window_count * sender_limit}. "
                f"Increase browser windows or lower the recipient count."
            )
        chunks: list[list[str]] = []
        index = 0
        base, remainder = divmod(total, window_count)
        for window_index in range(window_count):
            size = base + (1 if window_index < remainder else 0)
            if size <= 0:
                continue
            next_index = min(total, index + min(size, sender_limit))
            chunks.append(recipients[index:next_index])
            index = next_index
        if index < total:
            chunks[-1].extend(recipients[index:])
        return chunks

    def _build_campaign_jobs(
        self,
        sessions: list[BrowserSessionHandle],
        recipients: list[str],
        subject_template: str | list[str],
        body_template: str,
        attachment_html: str,
        attachment_formats: list[str],
        file_name_mode: str,
        file_name_value: str,
    ) -> list[dict[str, object]]:
        chunks = self._chunk_campaign_recipients(
            recipients,
            len(sessions),
            max(1, int(getattr(self.state, "sender_limit", 300))),
        )
        jobs: list[dict[str, object]] = []
        for index, (session, chunk) in enumerate(zip(sessions, chunks), start=1):
            tasks: list[dict[str, str]] = []
            window_label = f"Window {index}"
            for recipient in chunk:
                tag_values = self._dynamic_tag_values()
                selected_subject = random.choice(subject_template) if isinstance(subject_template, list) and subject_template else str(subject_template)
                subject = self._apply_tags_to_text(selected_subject, recipient, tag_values=tag_values)
                body_text = self._apply_tags_to_text(body_template, recipient, subject, tag_values=tag_values)
                attachment_text = self._apply_tags_to_text(attachment_html, recipient, subject, tag_values=tag_values)
                file_name_text = self._apply_tags_to_text(file_name_value, recipient, subject, tag_values=tag_values)
                tasks.append(
                    {
                        "recipient": recipient,
                        "subject": subject,
                        "body_text": body_text,
                        "attachment_html": attachment_text,
                        "attachment_format": attachment_formats[0] if attachment_formats else self.attach_format_value,
                        "file_name_value": file_name_text,
                    }
                )
            jobs.append(
                {
                    "session": session,
                    "window_label": window_label,
                    "tasks": tasks,
                }
            )
        return jobs

    def _queue_campaign_job(self, job: dict[str, object]) -> None:
        session = job["session"]
        if isinstance(session, BrowserSessionHandle):
            session.send_completed = 0
            session.send_total = len(job["tasks"])
            self._refresh_sessions()
        thread = QThread(self)
        worker = CampaignSendWorker(
            session,
            job["tasks"],
            self._campaign_pause_event,
            self._campaign_cancel_event,
            self._send_compose_with_playwright,
            window_label=str(job["window_label"]),
            delay_mode=getattr(self.state, "delay_type", "Random range"),
            delay_from=float(getattr(self.state, "delay_from", 0.5)),
            delay_to=float(getattr(self.state, "delay_to", 1.0)),
            retry_count=int(getattr(self.state, "retry_count", 3)),
            retry_enabled=bool(getattr(self.state, "retry_enabled", True)),
            convert_enabled=bool(self.attach_convert_checkbox.isChecked()),
            attachment_formats=self._normalize_attachment_format_values(
                self._pending_campaign_payload.get("attachment_formats")
                or self._pending_campaign_payload.get("attachment_format")
                or self.attach_format_value
            ),
            file_name_mode=str(self._pending_campaign_payload.get("attachment_file_name_mode") or self.attach_file_name_mode or "auto"),
            window_send_mode=getattr(self.state, "window_send_mode", "Parallel"),
        )
        worker.moveToThread(thread)
        # Start the bootstrap worker directly on the QThread. This avoids the
        # queued-start race that can leave the modal loader indeterminate.
        thread.started.connect(worker.run, Qt.DirectConnection)
        worker.log.connect(self._append_campaign_send_log)
        worker.progress.connect(self._on_campaign_worker_progress)
        worker.finished.connect(self._on_campaign_worker_finished)
        worker.finished.connect(worker.deleteLater)
        worker.finished.connect(thread.quit)
        thread.finished.connect(thread.deleteLater)
        self._campaign_threads.append(thread)
        self._campaign_workers[str(job["session"].session_id)] = worker
        thread.start()

    def _start_next_campaign_worker(self) -> None:
        if self._campaign_cancel_event.is_set() or self._campaign_paused:
            return
        if not self._campaign_worker_queue:
            return
        job = self._campaign_worker_queue.pop(0)
        self._queue_campaign_job(job)

    def _on_campaign_worker_progress(self, session_id: str, window_label: str, completed: int, total: int) -> None:
        session = next((item for item in self._browser_sessions if item.session_id == session_id), None)
        if session is not None:
            session.send_completed = completed
            session.send_total = total
            self._refresh_sessions()
        self._campaign_worker_progress[session_id] = completed
        self._campaign_worker_totals[session_id] = total
        self._campaign_completed = sum(self._campaign_worker_progress.values())
        remaining = max(self._campaign_total - self._campaign_completed, 0)
        self.progress_bar.setValue(int((self._campaign_completed / max(self._campaign_total, 1)) * 100))
        self.campaign_progress_text.setText(f"{self._campaign_completed} / {self._campaign_total} sent • {remaining} remaining")

    def _finish_campaign_runtime(self, message: str | None = None, *, cancelled: bool = False) -> None:
        self._campaign_active = False
        self._campaign_paused = False
        self._campaign_cancel_event.clear()
        self._campaign_pause_event.clear()
        self._campaign_threads.clear()
        self._campaign_workers.clear()
        self._campaign_worker_queue.clear()
        self._campaign_worker_progress.clear()
        self._campaign_worker_totals.clear()
        if cancelled:
            self.notify("Campaign cancelled")
            self._log_action("Campaign cancelled")
        elif message:
            self.notify(message)
        self._pending_campaign_payload = None
        self._update_campaign_action_state()

    def _on_campaign_worker_finished(
        self,
        session_id: str,
        window_label: str,
        cancelled: bool,
        completed: int,
        total: int,
        error_message: str,
    ) -> None:
        self._campaign_workers.pop(session_id, None)
        self._campaign_worker_progress[session_id] = completed
        self._campaign_worker_totals[session_id] = total
        session = next((item for item in self._browser_sessions if item.session_id == session_id), None)
        if session is not None:
            session.send_completed = completed
            session.send_total = total
            self._refresh_sessions()
        self._campaign_completed = sum(self._campaign_worker_progress.values())
        self.progress_bar.setValue(int((self._campaign_completed / max(self._campaign_total, 1)) * 100))
        if error_message:
            self._log_action(f"{window_label} finished with an error: {error_message}")
        if self._campaign_cancel_event.is_set() or cancelled:
            if not self._campaign_workers:
                self._finish_campaign_runtime(cancelled=True)
            return
        if self.state.window_send_mode == "Sequential" and self._campaign_worker_queue and not self._campaign_paused:
            self._start_next_campaign_worker()
        elif not self._campaign_workers and not self._campaign_worker_queue:
            self._append_campaign_send_log(f"[{QDateTime.currentDateTime().toString('hh:mm:ss')}] Campaign complete")
            self._finish_campaign_runtime(message=f"Sent {self._campaign_completed} email(s)")

    def _handle_campaign_primary_action(self) -> None:
        if self._campaign_active:
            self._handle_campaign_pause_resume()
            return
        self._start_blast()

    def _handle_campaign_pause_resume(self) -> None:
        if not self._campaign_active:
            return
        if self._campaign_paused:
            self._campaign_paused = False
            self._campaign_pause_event.set()
            self._update_campaign_action_state()
            self._log_action("Campaign resumed")
            if self.state.window_send_mode == "Sequential" and not self._campaign_workers and self._campaign_worker_queue:
                self._start_next_campaign_worker()
            return
        self._campaign_paused = True
        self._campaign_pause_event.clear()
        self._update_campaign_action_state()
        self._log_action("Campaign paused")

    def _handle_campaign_cancel(self) -> None:
        if not self._campaign_active:
            return
        self._campaign_cancel_event.set()
        self._campaign_pause_event.set()
        self._campaign_worker_queue.clear()
        self._update_campaign_action_state()
        self._log_action("Campaign cancel requested")

    def _begin_campaign_send(self, payload: dict[str, object]) -> None:
        if self._campaign_active:
            self.notify("Campaign is already running")
            return

        recipients_raw = payload.get("recipients") or []
        recipients = [str(item).strip() for item in recipients_raw if str(item).strip()]
        if not recipients:
            self.notify("Add customer emails before starting the campaign")
            self._log_action("Campaign blocked: no recipients available")
            return

        subject_values = payload.get("subjects")
        subject_template: str | list[str]
        if isinstance(subject_values, list):
            subject_template = [str(item).strip() for item in subject_values if str(item).strip()]
        else:
            subject_template = str(payload.get("subject") or "").strip()
        if not subject_template:
            subject_template = str(payload.get("subject") or "").strip()
        body_template = self._campaign_body_text(str(payload.get("body_text") or ""), str(payload.get("body_html") or ""))
        if not body_template:
            self.notify("Add body content before starting the campaign")
            self._log_action("Campaign blocked: no body content available")
            return

        attachment_html = str(payload.get("attachment_html") or "").strip()
        attachment_formats = self._normalize_attachment_format_values(
            payload.get("attachment_formats") or payload.get("attachment_format") or self.attach_format_value or "PDF document"
        )
        file_name_mode = str(payload.get("attachment_file_name_mode") or self.attach_file_name_mode or "auto")
        file_name_value = str(payload.get("attachment_file_name_value") or self.attach_file_name_value or "")

        sessions = self._usable_browser_sessions()
        if not sessions:
            self.notify("Open browser windows first")
            self._log_action("Campaign blocked: no browser windows available")
            return

        ordered_recipients = list(recipients)
        if self.state.email_send_order == "Random shuffle":
            random.shuffle(ordered_recipients)

        total = len(ordered_recipients)
        window_count = min(len(sessions), total)
        per_sender_limit = max(1, int(getattr(self.state, "sender_limit", 300)))
        if total > window_count * per_sender_limit:
            QMessageBox.warning(
                self,
                "Campaign capacity exceeded",
                f"The campaign needs {total} sends, but {window_count} window(s) at {per_sender_limit} per sender can only handle {window_count * per_sender_limit}.",
            )
            self._log_action("Campaign blocked: sender capacity exceeded")
            self.notify("Increase browser windows or lower the recipient count")
            return

        self._pending_campaign_payload = payload
        self._campaign_send_log_entries.clear()
        self._campaign_pause_event.set()
        self._campaign_cancel_event.clear()
        self._campaign_active = True
        self._campaign_paused = False
        self._campaign_total = total
        self._campaign_completed = 0
        self._campaign_worker_queue = []
        self._campaign_worker_progress.clear()
        self._campaign_worker_totals.clear()
        self.progress_bar.setValue(0)
        self.campaign_progress_text.setText(f"0 / {total} sent • {total} remaining")
        self._update_campaign_action_state()

        jobs = self._build_campaign_jobs(
            sessions[:window_count],
            ordered_recipients,
            subject_template,
            body_template,
            attachment_html,
            attachment_formats,
            file_name_mode,
            file_name_value,
        )
        self._campaign_worker_queue = list(jobs)
        self._log_action(f"Campaign ready for {total} recipient(s) using {window_count} browser window(s)")
        self.notify("Sending campaign")

        if self.state.window_send_mode == "Sequential":
            self._start_next_campaign_worker()
        else:
            queue = list(self._campaign_worker_queue)
            self._campaign_worker_queue.clear()
            for job in queue:
                self._queue_campaign_job(job)

    def _start_blast(self) -> None:
        if not self.state.logged_in:
            self.notify("Sign in first to start a campaign")
            return

        payload = self._current_campaign_payload()
        missing = self._campaign_missing_fields(payload)
        if missing:
            joined = ", ".join(missing)
            QMessageBox.warning(
                self,
                "Missing campaign fields",
                f"Please complete these fields before starting the campaign:\n\n- " + "\n- ".join(missing),
            )
            self.notify(f"Add {joined} before starting the campaign")
            self._log_action(f"Campaign blocked: missing {joined}")
            self._refresh_campaign_action_state()
            return

        reply = QMessageBox.question(
            self,
            "Start Campaign",
            "Start sending the campaign now?",
            QMessageBox.Yes | QMessageBox.No,
            QMessageBox.No,
        )
        if reply != QMessageBox.Yes:
            self._log_action("Campaign start cancelled by user")
            return

        recipients = payload["recipients"]
        if not self._usable_browser_sessions():
            self._log_action("Campaign requested: launching browser windows first")
            self._pending_campaign_payload = payload
            self._handle_launch()
            return

        self._execute_campaign_send(payload)

    def _html_to_plain_text(self, value: str) -> str:
        cleaned = re.sub(r"(?is)<(script|style).*?>.*?</\1>", "", value or "")
        cleaned = re.sub(r"(?i)<br\s*/?>", "\n", cleaned)
        cleaned = re.sub(r"(?i)</p\s*>", "\n\n", cleaned)
        cleaned = re.sub(r"<[^>]+>", "", cleaned)
        cleaned = html.unescape(cleaned)
        cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
        return cleaned.strip()

    def _campaign_body_text(self, body_text: str, body_html: str) -> str:
        if body_text.strip():
            return body_text.strip()
        if body_html.strip():
            return self._html_to_plain_text(body_html)
        return ""

    def _gmail_compose_url(self, recipient: str, subject: str, body_text: str) -> str:
        return (
            "https://mail.google.com/mail/?view=cm&fs=1"
            f"&to={quote_plus(recipient)}"
            f"&su={quote_plus(subject)}"
            f"&body={quote_plus(body_text)}"
        )

    def _browser_cdp_url(self, session: BrowserSessionHandle) -> str:
        if session.debug_port is None:
            raise RuntimeError("Browser session does not expose a CDP port.")
        return f"http://127.0.0.1:{session.debug_port}"

    def _render_html_to_jpg(self, html_content: str, output_path: Path) -> None:
        ensure_external_dependencies()
        from PIL import Image
        from playwright.sync_api import sync_playwright

        html_content = html_content.strip() or "<html><body></body></html>"
        with sync_playwright() as playwright:
            binary = self._browser_binary()
            if binary is None:
                raise RuntimeError("No bundled Chromium runtime was found for this app.")
            browser = playwright.chromium.launch(
                headless=True,
                executable_path=str(binary),
                args=["--disable-gpu", "--no-sandbox", "--disable-dev-shm-usage"],
            )
            try:
                context = browser.new_context(viewport={"width": 1440, "height": 900}, device_scale_factor=1)
                context.route(
                    "**/*",
                    lambda route: route.abort()
                    if route.request.resource_type in {"image", "media", "font", "stylesheet", "xhr", "fetch", "script"}
                    else route.continue_(),
                )
                page = context.new_page()
                page.set_content(html_content, wait_until="domcontentloaded", timeout=15000)
                page.emulate_media(media="screen")
                page.wait_for_timeout(500)
                width = int(page.evaluate(
                    """
                    () => Math.max(
                        document.body.scrollWidth,
                        document.documentElement.scrollWidth,
                        window.innerWidth
                    )
                    """
                ) or 1440)
                height = int(page.evaluate(
                    """
                    () => Math.max(
                        document.body.scrollHeight,
                        document.documentElement.scrollHeight,
                        window.innerHeight
                    )
                    """
                ) or 900)
                width = max(800, min(width + 40, 2400))
                height = max(600, min(height + 40, 5000))
                page.set_viewport_size({"width": width, "height": height})
                page.wait_for_timeout(250)
                png_path = output_path.with_suffix(".png")
                page.screenshot(path=str(png_path), full_page=True)
            finally:
                browser.close()

        image = Image.open(png_path)
        try:
            if image.mode in {"RGBA", "LA"}:
                background = Image.new("RGB", image.size, (255, 255, 255))
                background.paste(image, mask=image.split()[-1])
                background.save(output_path, format="JPEG", quality=92, optimize=True)
            else:
                image.convert("RGB").save(output_path, format="JPEG", quality=92, optimize=True)
        finally:
            image.close()
            try:
                png_path.unlink(missing_ok=True)
            except Exception:
                pass
        return

    def _attachment_image_size(self, jpg_path: Path) -> tuple[int, int]:
        ensure_external_dependencies()
        from PIL import Image

        image = Image.open(jpg_path)
        try:
            width, height = image.size
            if width <= 0 or height <= 0:
                raise RuntimeError("Unable to read attachment image size.")
            return width, height
        finally:
            image.close()

    def _attachment_page_dimensions(self, image_width: int, image_height: int) -> tuple[float, float]:
        if image_width >= image_height:
            page_width_in = 13.333
            page_height_in = max(6.0, 13.333 * (float(image_height) / float(image_width)))
        else:
            page_height_in = 11.0
            page_width_in = max(6.0, 11.0 * (float(image_width) / float(image_height)))
        return page_width_in, page_height_in

    def _export_jpg_to_format(self, jpg_path: Path, format_value: str, output_path: Path) -> Path:
        ensure_external_dependencies()
        from docx import Document
        from docx.shared import Inches
        from openpyxl import Workbook
        from openpyxl.drawing.image import Image as OpenPyxlImage
        from reportlab.lib.units import inch as reportlab_inch
        from reportlab.lib.utils import ImageReader
        from reportlab.pdfgen import canvas
        from pptx import Presentation
        from pptx.util import Inches as PptxInches

        format_value = (format_value or "").strip()
        if format_value == "PDF document":
            image_width, image_height = self._attachment_image_size(jpg_path)
            page_width_in, page_height_in = self._attachment_page_dimensions(image_width, image_height)
            pdf = canvas.Canvas(
                str(output_path),
                pagesize=(page_width_in * reportlab_inch, page_height_in * reportlab_inch),
            )
            pdf.drawImage(
                ImageReader(str(jpg_path)),
                0,
                0,
                width=page_width_in * reportlab_inch,
                height=page_height_in * reportlab_inch,
            )
            pdf.showPage()
            pdf.save()
            return output_path

        if format_value == "Word document (DOCX)":
            image_width, image_height = self._attachment_image_size(jpg_path)
            page_width_in, page_height_in = self._attachment_page_dimensions(image_width, image_height)
            doc = Document()
            section = doc.sections[0]
            section.page_width = Inches(page_width_in)
            section.page_height = Inches(page_height_in)
            margin = Inches(0.15)
            section.top_margin = margin
            section.bottom_margin = margin
            section.left_margin = margin
            section.right_margin = margin
            available_width = section.page_width - section.left_margin - section.right_margin
            doc.add_picture(str(jpg_path), width=available_width)
            doc.save(str(output_path))
            return output_path

        if format_value == "Excel spreadsheet (XLSX)" or format_value == "Excel template (XLTX)":
            image_width, image_height = self._attachment_image_size(jpg_path)
            workbook = Workbook()
            workbook.template = format_value == "Excel template (XLTX)"
            sheet = workbook.active
            sheet.title = "Attachment"
            sheet.freeze_panes = "A1"
            sheet.sheet_view.zoomScale = 90
            sheet.page_setup.orientation = "landscape" if image_width >= image_height else "portrait"
            sheet.page_setup.fitToWidth = 1
            sheet.page_setup.fitToHeight = 0
            sheet.page_margins.left = 0.1
            sheet.page_margins.right = 0.1
            sheet.page_margins.top = 0.1
            sheet.page_margins.bottom = 0.1
            sheet.column_dimensions["A"].width = max(12, min(60, image_width / 7.0))
            sheet.row_dimensions[1].height = max(18, min(360, image_height * 0.75))
            image_copy = OpenPyxlImage(str(jpg_path))
            max_width = 1100
            max_height = 1600
            scale = min(1.0, float(max_width) / float(image_width), float(max_height) / float(image_height))
            image_copy.width = max(1, int(image_width * scale))
            image_copy.height = max(1, int(image_height * scale))
            sheet.add_image(image_copy, "A1")
            workbook.save(str(output_path))
            return output_path

        if format_value in {"PowerPoint presentation (PPTX)", "PowerPoint slideshow (PPSX)"}:
            presentation = Presentation()
            image_width, image_height = self._attachment_image_size(jpg_path)
            slide_width_in, slide_height_in = self._attachment_page_dimensions(image_width, image_height)
            slide_width = PptxInches(slide_width_in)
            slide_height = PptxInches(slide_height_in)
            presentation.slide_width = slide_width
            presentation.slide_height = slide_height
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            slide.shapes.add_picture(str(jpg_path), 0, 0, width=slide_width, height=slide_height)
            presentation.save(str(output_path))
            return output_path

        if format_value == "Random format":
            return self._export_jpg_to_format(jpg_path, "PDF document", output_path.with_suffix(".pdf"))

        raise RuntimeError(f"Unsupported attachment format: {format_value or 'unknown'}")

    def _compose_attachment_paths(
        self,
        recipient: str,
        subject: str,
        attachment_html: str,
        format_value: str | list[str],
        file_name_mode: str,
        file_name_value: str,
        convert_enabled: bool,
        already_resolved: bool = False,
    ) -> list[Path]:
        ensure_external_dependencies()
        temp_dir = Path(tempfile.mkdtemp(prefix="ezymailer-gmail-"))
        resolved_attachment_html = attachment_html if already_resolved else self._apply_tags_to_text(attachment_html, recipient, subject)
        base_name = self._attachment_output_name_base(
            recipient,
            subject,
            str(format_value),
            file_name_mode,
            file_name_value,
            resolve_tags=not already_resolved,
        )
        jpg_path = temp_dir / f"{base_name}.jpg"
        self._render_html_to_jpg(resolved_attachment_html, jpg_path)

        if not convert_enabled:
            html_path = temp_dir / f"{base_name}.html"
            html_path.write_text(resolved_attachment_html or "<html><body></body></html>", encoding="utf-8")
            try:
                jpg_path.unlink(missing_ok=True)
            except Exception:
                pass
            return [html_path]

        raw_formats = [str(item).strip() for item in format_value] if isinstance(format_value, list) else [str(format_value or "").strip()]
        if any(item == "Random format" for item in raw_formats):
            format_values = [secrets.choice([
                "PDF document",
                "Excel spreadsheet (XLSX)",
                "Excel template (XLTX)",
                "PowerPoint presentation (PPTX)",
                "PowerPoint slideshow (PPSX)",
                "Word document (DOCX)",
            ])]
        else:
            format_values = self._normalize_attachment_format_values(format_value)

        final_paths: list[Path] = []
        for selected_format in format_values:
            extension = self._attachment_format_extension(selected_format)
            final_path = temp_dir / f"{base_name}{extension}"
            self._export_jpg_to_format(jpg_path, selected_format, final_path)
            final_paths.append(final_path)
        return final_paths

    def _sanitize_filename(self, value: str) -> str:
        cleaned = re.sub(r"[^A-Za-z0-9._-]+", "-", value or "").strip("-_.")
        return cleaned[:48]

    def _gmail_page_from_session(self, page, timeout_ms: int = 30000):
        page.wait_for_load_state("domcontentloaded", timeout=timeout_ms)
        if "mail.google.com" not in (page.url or ""):
            page.goto("https://mail.google.com/mail/u/0/#inbox", wait_until="domcontentloaded", timeout=timeout_ms)
        return page

    def _send_compose_with_playwright(
        self,
        session: BrowserSessionHandle,
        recipient: str,
        subject: str,
        body_text: str,
        attachment_html: str,
        attachment_format: str | list[str],
        file_name_mode: str,
        file_name_value: str,
        convert_enabled: bool,
        already_resolved: bool = False,
        log_steps: bool = True,
    ) -> None:
        attachment_paths = self._compose_attachment_paths(
            recipient,
            subject,
            attachment_html,
            attachment_format,
            file_name_mode,
            file_name_value,
            convert_enabled,
            already_resolved=already_resolved,
        )
        try:
            if log_steps:
                self._log_action(f"Preparing attachment file for {recipient}")
            self._send_compose_with_playwright_attachment(
                session,
                recipient,
                subject,
                body_text,
                attachment_paths,
                log_steps=log_steps,
            )
        finally:
            try:
                if attachment_paths:
                    shutil.rmtree(attachment_paths[0].parent, ignore_errors=True)
            except Exception:
                pass

    def _send_compose_with_playwright_attachment(
        self,
        session: BrowserSessionHandle,
        recipient: str,
        subject: str,
        body_text: str,
        attachment_paths: list[Path],
        *,
        log_steps: bool = True,
    ) -> None:
        ensure_external_dependencies()
        from playwright.sync_api import TimeoutError as PlaywrightTimeoutError
        from playwright.sync_api import sync_playwright

        try:
            with sync_playwright() as playwright:
                browser = None
                last_exc: Exception | None = None
                for _attempt in range(20):
                    try:
                        browser = playwright.chromium.connect_over_cdp(self._browser_cdp_url(session))
                        break
                    except Exception as exc:
                        last_exc = exc
                        time.sleep(0.5)
                if browser is None:
                    if last_exc is not None:
                        raise last_exc
                    raise RuntimeError("Unable to connect to Gmail browser.")

                if not browser.contexts:
                    raise RuntimeError("No browser context was available for Gmail automation.")
                context = browser.contexts[0]
                page = context.pages[0] if context.pages else context.new_page()
                page.set_default_timeout(10000)
                page.set_default_navigation_timeout(15000)
                self._gmail_page_from_session(page)
                page.bring_to_front()
                if log_steps:
                    self._log_action("Gmail page ready")

                compose_clicked = False
                try:
                    compose_button = page.get_by_role("button", name=re.compile(r"^Compose$", re.IGNORECASE))
                    compose_button.first.click(timeout=8000)
                    compose_clicked = True
                except Exception:
                    try:
                        page.goto("https://mail.google.com/mail/u/0/#inbox?compose=new", wait_until="domcontentloaded", timeout=10000)
                        compose_clicked = True
                    except Exception:
                        pass
                if not compose_clicked:
                    raise RuntimeError("Unable to open Gmail compose window.")
                if log_steps:
                    self._log_action("Gmail compose opened")

                recipient_box = self._gmail_recipient_input(page)
                recipient_box.fill(recipient)
                recipient_box.press("Enter")
                if log_steps:
                    self._log_action("Recipient entered")

                subject_box = page.locator('input[name="subjectbox"], input[placeholder*="Subject"]').first
                subject_box.wait_for(state="visible", timeout=10000)
                subject_box.fill(subject)
                if log_steps:
                    self._log_action("Subject entered")

                body_box = page.locator('div[role="textbox"][aria-label*="Message Body"], div[contenteditable="true"][aria-label*="Message Body"]').first
                body_box.wait_for(state="visible", timeout=10000)
                body_box.click()
                body_box.fill(body_text)
                if log_steps:
                    self._log_action("Body entered")

                attach_candidates = page.locator('input[type="file"]')
                if attach_candidates.count() == 0:
                    attach_button = page.get_by_role("button", name=re.compile(r"Attach files", re.IGNORECASE))
                    attach_button.first.click(timeout=8000)
                    attach_candidates = page.locator('input[type="file"]')
                if log_steps:
                    self._log_action("Attachment picker ready")
                files_to_attach = [str(path) for path in attachment_paths if path.exists()]
                if not files_to_attach:
                    raise RuntimeError("No attachment files were generated.")
                attach_candidates.last.set_input_files(files_to_attach)
                if log_steps:
                    self._log_action(
                        "Attachment selected: " + ", ".join(Path(path).name for path in files_to_attach)
                    )

                page.wait_for_timeout(1000)
                send_button = page.locator('div[role="button"][aria-label^="Send"]')
                if send_button.count() == 0:
                    send_button = page.locator('div[role="button"][data-tooltip^="Send"]')
                if send_button.count() == 0:
                    send_button = page.get_by_role("button", name=re.compile(r"^Send$", re.IGNORECASE))
                send_button.first.click(timeout=10000)
                page.wait_for_timeout(1500)
                if log_steps:
                    self._log_action(f"Gmail send completed for {recipient}")
        except PlaywrightTimeoutError as exc:
            raise RuntimeError(f"Gmail automation timed out: {exc}") from exc

    def _gmail_recipient_input(self, page):
        """Focus Gmail's active compose recipient editor before filling it.

        Gmail keeps its recipient input hidden until the compose recipient area
        receives focus. The hidden input can still match the old selector, so
        visibility must be established by clicking the active compose region.
        """
        recipient_selectors = (
            'input[aria-label="To recipients"]',
            'input[aria-label^="To"]',
            '[role="combobox"][aria-label*="To"]',
            'input[name="to"]',
        )
        focus_selectors = (
            '[role="dialog"] [aria-label="To recipients"]',
            '[role="dialog"] [aria-label^="To"]',
            '[role="dialog"] [role="combobox"]',
            '[role="dialog"] div[aria-label="To"]',
            '[role="dialog"] .agP',
        )
        deadline = time.monotonic() + 15
        while time.monotonic() < deadline:
            for selector in recipient_selectors:
                locator = page.locator(selector)
                try:
                    count = locator.count()
                except Exception:
                    continue
                for index in range(count):
                    candidate = locator.nth(index)
                    try:
                        if candidate.is_visible():
                            candidate.click()
                            return candidate
                    except Exception:
                        continue

            for selector in focus_selectors:
                locator = page.locator(selector).last
                try:
                    if locator.is_visible():
                        locator.click(timeout=1000)
                        break
                except Exception:
                    continue
            page.wait_for_timeout(250)

        raise RuntimeError("Gmail compose recipient field did not become available.")

    def _maybe_send_with_playwright(
        self,
        session: BrowserSessionHandle,
        recipient: str,
        subject: str,
        body_text: str,
        attachment_html: str,
        attachment_format: str | list[str],
        file_name_mode: str,
        file_name_value: str,
    ) -> None:
        self._send_compose_with_playwright(
            session,
            recipient,
            subject,
            body_text,
            attachment_html,
            attachment_format,
            file_name_mode,
            file_name_value,
            convert_enabled=True,
        )

    def _execute_campaign_send(self, payload: dict[str, object]) -> None:
        self._begin_campaign_send(payload)

    def _show_launch_loader(self, title: str, subtitle: str) -> None:
        self.window().show_launch_loader(title, subtitle)

    def _current_body_widget(self) -> BodyDraftEditor | None:
        widget = self.body_tabs.currentWidget()
        if isinstance(widget, BodyDraftEditor):
            return widget
        return None

    def _sync_active_body_widget_refs(self) -> None:
        widget = self._current_body_widget()
        if widget is None:
            return
        self.body_editor = widget.plain_editor
        self.html_message_editor = widget.html_editor
        self.message_stack = widget.stack
        self.state.body_mode = widget.mode_text()
        self.normal_message_button = widget.plain_button
        self.html_message_button = widget.html_button
        self.body_mode_group = widget._plain_group

    def _update_body_tab_controls(self) -> None:
        limit_reached = self.body_tabs.count() >= MAX_BODY_TABS
        self.body_add_button.setEnabled(not limit_reached)
        self.body_add_button.setToolTip("Maximum of 50 bodies reached" if limit_reached else "Add a new body")

    def _update_attachment_tab_controls(self) -> None:
        limit_reached = self.attach_tabs.count() >= MAX_ATTACHMENT_TABS
        self.attach_add_button.setEnabled(not limit_reached)
        self.attach_add_button.setToolTip("Maximum of 50 content tabs reached" if limit_reached else "Add a new HTML content tab")

    def _set_subject_item_data(
        self,
        item: QListWidgetItem,
        record_id: int | None,
        title: str,
        subject: str,
        *,
        local_only: bool = False,
        local_draft_id: int | None = None,
    ) -> None:
        item.setData(Qt.UserRole, record_id)
        item.setData(Qt.UserRole + 1, title)
        item.setData(Qt.UserRole + 2, subject)
        item.setData(ROLE_LOCAL_ONLY, local_only)
        item.setData(ROLE_LOCAL_DRAFT_ID, local_draft_id)
        item.setText(title or subject or "Untitled Subject")

    def _selected_subject_item(self) -> QListWidgetItem | None:
        return self.subject_drafts_list.currentItem()

    def _clear_subject_selection(self) -> None:
        self.subject_drafts_list.blockSignals(True)
        try:
            self.subject_drafts_list.clearSelection()
        finally:
            self.subject_drafts_list.blockSignals(False)

    def _update_subject_toggle_visibility(self) -> None:
        self.subject_toggle_button.setVisible(self.subject_drafts_list.count() > 0)

    def _subject_item_subject(self, item: QListWidgetItem | None) -> str:
        if item is None:
            return ""
        return str(item.data(Qt.UserRole + 2) or "")

    def _subject_item_title(self, item: QListWidgetItem | None) -> str:
        if item is None:
            return ""
        return str(item.data(Qt.UserRole + 1) or item.text() or "")

    def _open_subject_manager(self) -> None:
        if not self.state.logged_in or not self.state.auth_token:
            self.notify("Sign in first to manage subjects")
            return

        if self._subject_manager_dialog is not None:
            if self._subject_manager_dialog.isVisible():
                try:
                    self._subject_manager_dialog.raise_()
                    self._subject_manager_dialog.activateWindow()
                except Exception:
                    pass
                return
            try:
                self._subject_manager_dialog.deleteLater()
            except Exception:
                pass
            self._subject_manager_dialog = None

        try:
            dialog = SubjectDraftsDialog(self, self.state.auth_token, scale=self._scale, on_changed=self._load_subjects)
            self._subject_manager_dialog = dialog
            dialog.setAttribute(Qt.WA_DeleteOnClose, True)
            dialog.finished.connect(lambda _result, d=dialog: self._on_subject_manager_finished(d))
            dialog.setWindowModality(Qt.ApplicationModal)
            dialog.show()
            dialog.raise_()
            dialog.activateWindow()
            QTimer.singleShot(0, dialog.raise_)
            QTimer.singleShot(0, dialog.activateWindow)
        except Exception as exc:
            self._subject_manager_dialog = None
            self._log_action(f"Failed to open subject manager: {exc}")
            self.notify("Unable to open subject manager")

    def _on_subject_manager_finished(self, dialog: SubjectDraftsDialog) -> None:
        if self._subject_manager_dialog is dialog:
            self._subject_manager_dialog = None
        subject = dialog.selected_subject()
        if subject:
            self.subject_input.blockSignals(True)
            try:
                self.subject_input.setText(subject)
            except Exception:
                pass
            finally:
                self.subject_input.blockSignals(False)
            self.state.subject_text = subject
            self._schedule_subject_body_save()
        self._load_subjects()

    def _on_subject_selection_changed(self, current: QListWidgetItem | None, previous: QListWidgetItem | None) -> None:
        if self._workspace_loading:
            return
        if current is None:
            return
        subject = self._subject_item_subject(current)
        if subject:
            self.subject_input.blockSignals(True)
            try:
                self.subject_input.setText(subject)
            finally:
                self.subject_input.blockSignals(False)
            self.state.subject_text = subject

    def _new_subject_draft(self) -> None:
        current_text = self.subject_input.text().strip()
        if self.subject_drafts_list.count() >= MAX_SUBJECTS:
            self.notify("Maximum 100 subjects allowed")
            self._update_subject_count_label()
            return
        if current_text:
            self._save_subject_draft(force_new=True)
        self._clear_subject_selection()
        self.subject_input.blockSignals(True)
        try:
            self.subject_input.clear()
        finally:
            self.subject_input.blockSignals(False)
        self.state.subject_text = ""
        self._update_subject_toggle_visibility()
        self._persist_subject_body_state()

    def _subject_draft_rows(self) -> list[dict[str, object]]:
        rows: list[dict[str, object]] = []
        for index in range(self.subject_drafts_list.count()):
            item = self.subject_drafts_list.item(index)
            rows.append(
                {
                    "title": self._subject_item_title(item),
                    "subject": self._subject_item_subject(item),
                }
            )
        return rows

    def _save_subject_draft(self, *, force_new: bool = False) -> None:
        if self._workspace_loading or not self.state.logged_in or not self.state.username:
            return

        subject = self.subject_input.text().strip()
        if not subject:
            self.notify("Enter a subject first")
            return
        if len(subject) > 300:
            self.notify("Subject must be 300 characters or less")
            return
        current_item = self._selected_subject_item()
        if current_item is None and self.subject_drafts_list.count() >= MAX_SUBJECTS:
            self.notify("Maximum 100 subjects allowed")
            self._update_subject_count_label()
            return

        self._show_subject_body_loader("Saving subject.")
        try:
            title = subject[:64]
            if current_item is not None and not force_new:
                self._set_subject_item_data(
                    current_item,
                    None,
                    title,
                    subject,
                    local_only=True,
                    local_draft_id=None,
                )
            else:
                new_item = QListWidgetItem()
                self._set_subject_item_data(
                    new_item,
                    None,
                    title,
                    subject,
                    local_only=True,
                    local_draft_id=None,
                )
                self.subject_drafts_list.addItem(new_item)
                self.subject_drafts_list.setCurrentItem(new_item)
            self.state.subject_text = subject
            self._update_subject_count_label()
            self._update_subject_toggle_visibility()
            self._persist_subject_body_state()
            self._log_action(f"Saved subject: {title}")
            self.notify("Subject saved")
        except Exception as exc:
            self._log_action(f"Failed to save subject: {exc}")
            self.notify("Unable to save subject")
        finally:
            self._hide_subject_body_loader()

    def _delete_subject_draft(self) -> None:
        item = self._selected_subject_item()
        if item is None:
            return
        title = self._subject_item_title(item)
        row = self.subject_drafts_list.row(item)
        self.subject_drafts_list.takeItem(row)
        self.subject_input.clear()
        self.state.subject_text = ""
        self._update_subject_count_label()
        self._update_subject_toggle_visibility()
        self._persist_subject_body_state()
        self._log_action(f"Removed subject: {title or 'Untitled'}")
        self.notify("Subject removed")

    def _create_body_draft_tab(self, record: dict[str, object] | None = None) -> BodyDraftEditor:
        widget = BodyDraftEditor(scale=self._scale)
        if record:
            mode = str(record.get("mode") or "Normal Message")
            title = str(record.get("title") or self._body_record_title(record, self.body_tabs.count()))
            body_text = str(record.get("plain_text") or record.get("body_text") or "")
            body_html = str(record.get("html_text") or record.get("body_html") or "")
            widget.set_content(title, mode, body_text, body_html or body_text, local_only=True)
        else:
            title = self._next_body_title()
            widget.set_content(title, "Normal Message", "Hello {{first_name}},\n\nThis is a body message.", "<div></div>")

        widget.titleChanged.connect(lambda title_text, w=widget: self._rename_body_tab(w, title_text))
        widget.contentChanged.connect(self._schedule_subject_body_save)
        widget.modeChanged.connect(lambda mode, w=widget: self._on_body_mode_changed(w, mode))
        widget.previewRequested.connect(lambda w=widget: self._preview_body_editor_html(w))
        return widget

    def _body_tab_label(self, widget: BodyDraftEditor, index: int) -> str:
        title = widget.title_text()
        if title:
            base = title
        else:
            base = f"Body {index + 1}"
        mode_label = "HTML" if widget.mode_text() == "HTML Message" else "Text"
        return f"{base} [{mode_label}]"

    def _body_record_title(self, record: dict[str, object], index: int) -> str:
        details = self._decode_setting_value(record.get("details_json"))
        title = str(record.get("title") or "").strip()
        source_file = ""
        if isinstance(details, dict):
            source_file = str(details.get("source_file") or "").strip()
        title_normalized = title.lower()
        if source_file or not title or self._looks_like_auto_body_title(title_normalized):
            return f"Body {index + 1}"
        return title

    def _looks_like_auto_body_title(self, title_normalized: str) -> bool:
        if not title_normalized:
            return True
        if title_normalized.startswith("body "):
            return False
        if "inv" in title_normalized:
            return True
        if "copy" in title_normalized and any(ch.isdigit() for ch in title_normalized):
            return True
        if title_normalized.startswith("text body") or title_normalized.startswith("html body"):
            return True
        return False

    def _refresh_body_tab_labels(self) -> None:
        for index in range(self.body_tabs.count()):
            widget = self.body_tabs.widget(index)
            if isinstance(widget, BodyDraftEditor):
                self.body_tabs.setTabText(index, self._body_tab_label(widget, index))

    def _rename_body_tab(self, widget: BodyDraftEditor, title: str) -> None:
        index = self.body_tabs.indexOf(widget)
        if index < 0:
            return
        self.body_tabs.setTabText(index, self._body_tab_label(widget, index))

    def _on_body_mode_changed(self, widget: BodyDraftEditor, mode: str) -> None:
        if widget is self._current_body_widget():
            self.state.body_mode = mode
            self.normal_message_button = widget.plain_button
            self.html_message_button = widget.html_button
        self._refresh_body_tab_labels()
        self._schedule_subject_body_save()

    def _add_body_draft_tab(self, record: dict[str, object] | None = None, select: bool = True) -> BodyDraftEditor | None:
        if self.body_tabs.count() >= MAX_BODY_TABS:
            self._update_body_tab_controls()
            self.notify("Maximum of 50 bodies reached")
            return None
        widget = self._create_body_draft_tab(record)
        tab_number = self.body_tabs.count() + 1
        tab_label = f"Body {tab_number} [{'HTML' if widget.mode_text() == 'HTML Message' else 'Text'}]"
        index = self.body_tabs.addTab(widget, tab_label)
        self.body_tabs.setTabIcon(index, self._standard_icon(QStyle.SP_FileIcon))
        self._refresh_body_tab_labels()
        self._update_body_tab_controls()
        if select:
            self.body_tabs.setCurrentIndex(index)
        return widget

    def _new_body_draft_tab(self) -> None:
        self._add_body_draft_tab(select=True)

    def _remove_body_draft_tab(self, index: int) -> None:
        widget = self.body_tabs.widget(index)
        if not isinstance(widget, BodyDraftEditor):
            return
        self.body_tabs.removeTab(index)
        if self.body_tabs.count() == 0:
            self._add_body_draft_tab(select=True)
        else:
            self._refresh_body_tab_labels()
        self._update_body_tab_controls()
        self._sync_active_body_widget_refs()
        self._persist_subject_body_state()
        self._log_action("Removed body tab")

    def _on_body_tab_changed(self, _index: int) -> None:
        self._sync_active_body_widget_refs()
        self._schedule_subject_body_save()

    def _current_attachment_widget(self) -> AttachmentDraftEditor | None:
        widget = self.attach_tabs.currentWidget()
        return widget if isinstance(widget, AttachmentDraftEditor) else None

    def _sync_active_attachment_widget_refs(self) -> None:
        if self.attach_tabs.count() == 0:
            self._add_attachment_draft_tab(select=True)
        widget = self._current_attachment_widget()
        self._active_attachment_widget = widget
        if widget is not None:
            self.html_editor = widget.html_editor
            self.state.html_template_text = widget.html_editor.toPlainText()
        self._update_attachment_tab_controls()

    def _attachment_tab_label(self, widget: AttachmentDraftEditor, index: int) -> str:
        title = widget.title_text() or f"Content {index + 1}"
        return f"{title} [HTML]"

    def _refresh_attachment_tab_labels(self) -> None:
        for index in range(self.attach_tabs.count()):
            widget = self.attach_tabs.widget(index)
            if isinstance(widget, AttachmentDraftEditor):
                self.attach_tabs.setTabText(index, self._attachment_tab_label(widget, index))

    def _rename_attachment_tab(self, widget: AttachmentDraftEditor, title: str) -> None:
        index = self.attach_tabs.indexOf(widget)
        if index < 0:
            return
        self.attach_tabs.setTabText(index, self._attachment_tab_label(widget, index))

    def _next_attachment_title(self) -> str:
        return f"Content {self.attach_tabs.count() + 1}"

    def _looks_like_auto_attachment_title(self, title_normalized: str) -> bool:
        if not title_normalized:
            return True
        if title_normalized.startswith("content "):
            return True
        if title_normalized.startswith("html content"):
            return True
        if "copy" in title_normalized and any(ch.isdigit() for ch in title_normalized):
            return True
        return False

    def _attachment_record_title(self, record: dict[str, object], index: int) -> str:
        details = self._decode_setting_value(record.get("details_json"))
        title = str(record.get("title") or "").strip()
        source_file = ""
        if isinstance(details, dict):
            source_file = str(details.get("source_file") or "").strip()
        title_normalized = title.lower()
        if source_file or not title or self._looks_like_auto_attachment_title(title_normalized):
            return f"Content {index + 1}"
        return title

    def _create_attachment_draft_tab(self, record: dict[str, object] | None = None) -> AttachmentDraftEditor:
        widget = AttachmentDraftEditor(scale=self._scale)
        if record:
            title = self._attachment_record_title(record, self.attach_tabs.count())
            html_text = str(record.get("html_text") or record.get("body_html") or record.get("body_text") or "")
            widget.set_content(title, html_text, local_only=True)
        else:
            title = self._next_attachment_title()
            widget.set_content(title, "")

        widget.titleChanged.connect(lambda title_text, w=widget: self._rename_attachment_tab(w, title_text))
        widget.titleChanged.connect(lambda _title, self=self: self._persist_attachment_state())
        widget.contentChanged.connect(self._persist_attachment_state)
        widget.previewRequested.connect(lambda w=widget: self._preview_attachment_editor_html(w))
        return widget

    def _add_attachment_draft_tab(self, record: dict[str, object] | None = None, select: bool = True) -> AttachmentDraftEditor | None:
        if self.attach_tabs.count() >= MAX_ATTACHMENT_TABS:
            self._update_attachment_tab_controls()
            self.notify("Maximum of 50 content tabs reached")
            return None
        widget = self._create_attachment_draft_tab(record)
        tab_number = self.attach_tabs.count() + 1
        index = self.attach_tabs.addTab(widget, f"Content {tab_number} [HTML]")
        self.attach_tabs.setTabIcon(index, self._standard_icon(QStyle.SP_FileDialogContentsView))
        self._refresh_attachment_tab_labels()
        self._update_attachment_tab_controls()
        if select:
            self.attach_tabs.setCurrentIndex(index)
        return widget

    def _new_attachment_draft_tab(self) -> None:
        self._add_attachment_draft_tab(select=True)

    def _remove_attachment_draft_tab(self, index: int) -> None:
        widget = self.attach_tabs.widget(index)
        if not isinstance(widget, AttachmentDraftEditor):
            return
        self.attach_tabs.removeTab(index)
        if self.attach_tabs.count() == 0:
            self._add_attachment_draft_tab(select=True)
        else:
            self._refresh_attachment_tab_labels()
        self._update_attachment_tab_controls()
        self._sync_active_attachment_widget_refs()
        self._persist_attachment_state()
        self._log_action("Removed attachment tab")

    def _on_attachment_tab_changed(self, _index: int) -> None:
        self._sync_active_attachment_widget_refs()
        self._persist_attachment_state()

    def _reset_attachment_tabs(self) -> None:
        self._show_subject_body_loader("Resetting attachment tabs.")
        try:
            _delete_attachment_state()
            _delete_local_drafts("attachment")
            self.attach_format_value = "PDF document"
            self.attach_file_name_mode = "auto"
            self.attach_file_name_value = ""
            self.attach_tabs.blockSignals(True)
            try:
                for index in reversed(range(self.attach_tabs.count())):
                    self.attach_tabs.removeTab(index)
            finally:
                self.attach_tabs.blockSignals(False)
            self._add_attachment_draft_tab(select=True)
            self._refresh_attachment_tab_labels()
            self._update_attachment_tab_controls()
            self._sync_active_attachment_widget_refs()
            self._log_action("Reset attachment tabs")
            self.notify("Content tabs reset")
        except Exception as exc:
            self._log_action(f"Failed to reset attachment tabs: {exc}")
            self.notify("Unable to reset content tabs")
        finally:
            self._hide_subject_body_loader()

    def _import_html_attachment_file(self, path: Path) -> int:
        html_text = self._read_text_template_file(path)
        if not html_text.strip():
            return 0
        title = self._next_attachment_title()
        self._create_and_store_attachment_tab(title=title, html_text=html_text, source_name=path.name)
        return 1

    def _create_and_store_attachment_tab(self, *, title: str, html_text: str, source_name: str = "") -> AttachmentDraftEditor | None:
        if self.attach_tabs.count() >= MAX_ATTACHMENT_TABS:
            return None
        record = {
            "title": title[:64] or "Content",
            "html_text": html_text,
        }
        widget = self._add_attachment_draft_tab(record, select=False)
        if widget is None:
            return None
        self._refresh_attachment_tab_labels()
        self._update_attachment_tab_controls()
        self._persist_attachment_state()
        return widget

    def _upload_attachment_files(self) -> None:
        file_paths, _ = QFileDialog.getOpenFileNames(
            self,
            "Upload attachment content",
            "",
            "HTML files (*.html *.htm);;All files (*)",
        )
        if not file_paths:
            return

        if self.attach_tabs.count() == 1:
            placeholder = self._current_attachment_widget()
            if placeholder is not None and not placeholder.html_editor.toPlainText().strip():
                self.attach_tabs.removeTab(0)

        imported = 0
        self._show_subject_body_loader("Uploading attachment files.")
        self._workspace_loading = True
        try:
            for file_name in file_paths:
                if self.attach_tabs.count() >= MAX_ATTACHMENT_TABS:
                    self.notify("Maximum of 50 content tabs reached")
                    break
                path = Path(file_name)
                if path.suffix.lower() in {".html", ".htm"}:
                    imported += self._import_html_attachment_file(path)
                else:
                    self._log_action(f"Skipped unsupported content file: {path.name}")
        except Exception as exc:
            self._log_action(f"Failed to upload attachment files: {exc}")
            self.notify(f"Unable to upload content: {exc}")
        finally:
            self._workspace_loading = False
            self._hide_subject_body_loader()

        if self.attach_tabs.count() == 0:
            self._add_attachment_draft_tab(select=True)
        else:
            self._refresh_attachment_tab_labels()
            self._update_attachment_tab_controls()
            self._sync_active_attachment_widget_refs()

        if imported > 0:
            self._log_action(f"Uploaded {imported} attachment content(s)")
            self.notify(f"Uploaded {imported} content(s)")
        self._persist_attachment_state()

    def _load_attachment_tabs_from_local(self) -> None:
        state_payload = _load_attachment_state()
        _delete_local_drafts("attachment")
        tab_rows = state_payload.get("tabs") if isinstance(state_payload, dict) else []
        if not isinstance(tab_rows, list):
            tab_rows = []
        selected_index = int(state_payload.get("selected_index") or 0) if isinstance(state_payload, dict) else 0
        self.attach_tabs.blockSignals(True)
        try:
            for index in reversed(range(self.attach_tabs.count())):
                self.attach_tabs.removeTab(index)
            for row in tab_rows:
                if isinstance(row, dict):
                    self._add_attachment_draft_tab(row, select=False)
            if self.attach_tabs.count() == 0:
                self._add_attachment_draft_tab(select=True)
            else:
                self.attach_tabs.setCurrentIndex(min(max(selected_index, 0), self.attach_tabs.count() - 1))
            convert_enabled = bool(state_payload.get("convert_enabled", True)) if isinstance(state_payload, dict) else True
            format_value = self._normalize_attachment_format_value(str(state_payload.get("format_value") or self.attach_format_value or "PDF document")) if isinstance(state_payload, dict) else self._normalize_attachment_format_value(self.attach_format_value)
            file_name_mode = str(state_payload.get("file_name_mode") or self.attach_file_name_mode or "auto") if isinstance(state_payload, dict) else self.attach_file_name_mode
            file_name_value = str(state_payload.get("file_name_value") or self.attach_file_name_value or "") if isinstance(state_payload, dict) else self.attach_file_name_value
            self.attach_convert_checkbox.blockSignals(True)
            try:
                self.attach_convert_checkbox.setChecked(convert_enabled)
            finally:
                self.attach_convert_checkbox.blockSignals(False)
            self.attach_format_value = format_value or "PDF document"
            self.attach_file_name_mode = file_name_mode if file_name_mode in {"auto", "custom"} else "auto"
            self.attach_file_name_value = file_name_value
            self.attach_format_label.setText(self._attachment_format_summary(self.attach_format_value))
        finally:
            self.attach_tabs.blockSignals(False)
        self._refresh_attachment_tab_labels()
        self._update_attachment_tab_controls()
        self._sync_active_attachment_widget_refs()
        self._refresh_campaign_action_state()

    def _schedule_attachment_save(self) -> None:
        if self._workspace_loading:
            return
        self._refresh_campaign_action_state()
        self._attachment_save_timer.start()

    def _persist_attachment_state(self) -> None:
        if self._workspace_loading:
            return
        try:
            if self.attach_tabs.count() == 0:
                self._add_attachment_draft_tab(select=True)
            tabs_payload: list[dict[str, object]] = []
            for index in range(self.attach_tabs.count()):
                widget = self.attach_tabs.widget(index)
                if not isinstance(widget, AttachmentDraftEditor):
                    continue
                payload = widget.payload()
                title = payload["title"] or f"Content {index + 1}"
                html_text = payload["html_text"]
                tabs_payload.append(
                    {
                        "title": title[:64] or "Content",
                        "html_text": html_text,
                        "kind": "attachment",
                    }
                )
                widget.local_only = True
                self.attach_tabs.setTabText(index, self._attachment_tab_label(widget, index))
            active_widget = self._current_attachment_widget()
            if active_widget is not None:
                self.state.html_template_text = active_widget.html_editor.toPlainText()
            _delete_local_drafts("attachment")
            _upsert_attachment_state(
                {
                    "tabs": tabs_payload,
                    "selected_index": self.attach_tabs.currentIndex(),
                    "convert_enabled": self.attach_convert_checkbox.isChecked(),
                    "format_value": self.attach_format_value,
                    "file_name_mode": self.attach_file_name_mode,
                    "file_name_value": self.attach_file_name_value,
                }
            )
        except Exception as exc:
            self._log_action(f"Failed to save attachment content: {exc}")
        finally:
            self._refresh_campaign_action_state()

    def _preview_attachment_editor_html(self, widget: AttachmentDraftEditor | None) -> None:
        if widget is None:
            self.notify("No attachment content available")
            return
        html_content = widget.html_editor.toPlainText().strip()
        if not html_content:
            self.notify("Add HTML content first")
            return
        title = widget.title_text() or "HTML Content Preview"
        dialog = self._build_preview_dialog(title, html_content, "Previewing the selected HTML attachment content.")
        dialog.show()
        dialog.raise_()
        dialog.activateWindow()
        self._log_action(f"Opened attachment preview: {title}")
        self.notify("Content preview opened")

    def _complete_launch(self, target: int) -> None:
        self.window().hide_launch_loader()
        if target <= 0:
            return
        launched: list[BrowserSessionHandle] = []
        try:
            for index in range(1, target + 1):
                launched.append(self._launch_browser_process(index))
        except Exception as exc:
            self._terminate_browser_sessions()
            self._log_action(f"Browser launch failed: {exc}")
            self.notify("Unable to launch browser")
            return

        self._browser_sessions = launched
        self._sync_session_state_from_handles()
        self._refresh_sessions()
        self._browser_watch_timer.start()
        for session in self._browser_sessions:
            try:
                _upsert_local_browser_session(
                    self.state.username,
                    session.session_id,
                    session.title,
                    session.browser_name,
                    session.mode,
                    session.status,
                    session.process.pid if session.process is not None else None,
                    self.state.launch_preset or "Default",
                    {
                        "browser_mode": self.state.browser_mode,
                        "profile_dir": str(session.profile_dir) if session.profile_dir else "",
                    },
                )
            except Exception:
                pass
        self._log_action(f"Started {len(launched)} browser window(s)")
        self.notify(f"Launch started for {len(launched)} browser window(s)")
        if self._pending_campaign_payload:
            payload = self._pending_campaign_payload
            QTimer.singleShot(2500, lambda p=payload: self._execute_campaign_send(p))

    def _clear_subject_body(self) -> None:
        self._subject_body_save_timer.stop()
        self.subject_input.clear()
        self.state.subject_text = ""
        current = self._current_body_widget()
        if current is None:
            current = self._add_body_draft_tab(select=True)
        current.set_content("Body 1", "Normal Message", "", "")
        self.state.body_mode = "Normal Message"
        self.state.plain_body_text = ""
        self.state.html_message_text = ""
        self.state.html_template_text = ""
        self._log_action("Cleared subject and body")
        self._schedule_subject_body_save()

    def _save_current_body_draft(self) -> None:
        if self._workspace_loading or not self.state.logged_in or not self.state.username:
            return
        self._show_subject_body_loader("Saving body.")
        try:
            self._persist_subject_body_state()
            current = self._current_body_widget()
            title = current.title_text() if current is not None else "Body"
            self._log_action(f"Saved body: {title or 'Body'}")
            self.notify("Body saved")
        except Exception as exc:
            self._log_action(f"Failed to save body: {exc}")
            self.notify("Unable to save body")
        finally:
            self._hide_subject_body_loader()

    def _show_subject_body_loader(self, subtitle: str) -> None:
        self.window().show_launch_loader("Please wait", subtitle)
        QApplication.processEvents()

    def _hide_subject_body_loader(self) -> None:
        self.window().hide_launch_loader()

    def _schedule_subject_body_save(self) -> None:
        if self._workspace_loading or not self.state.logged_in or not self.state.username:
            return
        self._refresh_campaign_action_state()
        self._subject_body_save_timer.start()

    def _read_text_template_file(self, path: Path) -> str:
        return path.read_text(encoding="utf-8-sig").strip()

    def _load_subject_from_file(self) -> None:
        file_path, _ = QFileDialog.getOpenFileName(
            self,
            "Load subject template",
            "",
            "CSV files (*.csv);;All files (*)",
        )
        if not file_path:
            return

        path = Path(file_path)
        self._show_subject_body_loader(f"Loading subject template from {path.name}.")
        try:
            if path.suffix.lower() != ".csv":
                raise ValueError("Please choose a CSV file.")
            raw_text = self._read_text_template_file(path)
            subjects = []
            for row in csv.reader(raw_text.splitlines()):
                for value in row:
                    subject = str(value or "").strip()
                    if subject:
                        subjects.append(subject)
            if len(subjects) > 100:
                raise ValueError("CSV can contain at most 100 subjects.")
            if self.subject_drafts_list.count() + len(subjects) > MAX_SUBJECTS:
                raise ValueError("Total subjects cannot exceed 100.")
            for subject in subjects:
                if len(subject) > 300:
                    raise ValueError("Each subject must be 300 characters or less.")
            if not subjects:
                subjects = [path.stem.replace("_", " ").strip() or "Subject"]

            self._workspace_loading = True
            self.subject_drafts_list.blockSignals(True)
            try:
                for subject in subjects:
                    item = QListWidgetItem()
                    self._set_subject_item_data(
                        item,
                        None,
                        subject[:64] or "Subject",
                        subject,
                        local_only=True,
                        local_draft_id=None,
                    )
                    self.subject_drafts_list.addItem(item)
                self.subject_drafts_list.setCurrentRow(self.subject_drafts_list.count() - 1)
            finally:
                self.subject_drafts_list.blockSignals(False)
            current = self.subject_drafts_list.currentItem()
            if current is not None:
                self.subject_input.setText(self._subject_item_subject(current))
                self.state.subject_text = self.subject_input.text().strip()
            self._update_subject_count_label()
            self._update_subject_toggle_visibility()
            self._persist_subject_body_state()
            self._log_action(f"Loaded {len(subjects)} subject(s) from {path.name}")
            self.notify(f"Loaded {len(subjects)} subject(s)")
        except Exception as exc:
            self.notify(f"Unable to load subject file: {exc}")
            self._log_action(f"Failed to load subject template from {path.name}: {exc}")
        finally:
            self._workspace_loading = False
            self._hide_subject_body_loader()

    def _load_body_from_file(self) -> None:
        file_path, _ = QFileDialog.getOpenFileName(
            self,
            "Load body template",
            "",
            "Text/HTML files (*.txt *.html *.htm);;All files (*)",
        )
        if not file_path:
            return

        path = Path(file_path)
        self._show_subject_body_loader(f"Loading body template from {path.name}.")
        try:
            body = self._read_text_template_file(path)
            is_html = path.suffix.lower() in {".html", ".htm"} or "<html" in body.lower()
            title = self._next_body_title()
            self._workspace_loading = True
            body_widget = self._create_and_store_body_tab(
                title=title,
                mode="HTML Message" if is_html else "Normal Message",
                plain_text="" if is_html else body,
                html_text=body if is_html else "",
                source_name=path.name,
            )
            if body_widget is None:
                return
            if is_html:
                body_widget.set_content(title, "HTML Message", "", body)
            else:
                body_widget.set_content(title, "Normal Message", body, "")
            self._sync_active_body_widget_refs()
            self._log_action(f"Loaded body from {path.name}")
            self.notify(f"Loaded body from {path.name}")
        except Exception as exc:
            self.notify(f"Unable to load body file: {exc}")
            self._log_action(f"Failed to load body template from {path.name}: {exc}")
        finally:
            self._workspace_loading = False
            self._hide_subject_body_loader()

    def _upload_body_files(self) -> None:
        if not self.state.logged_in or not self.state.auth_token:
            self.notify("Sign in first to upload bodies")
            return

        file_paths, _ = QFileDialog.getOpenFileNames(
            self,
            "Upload bodies",
            "",
            "CSV or HTML files (*.csv *.html *.htm);;CSV files (*.csv);;HTML files (*.html *.htm);;All files (*)",
        )
        if not file_paths:
            return

        if self.body_tabs.count() == 1:
            placeholder = self._current_body_widget()
            if placeholder is not None:
                payload = placeholder.payload()
                plain = str(payload.get("plain_text") or "").strip()
                html_text = str(payload.get("html_text") or "").strip()
                if plain in {"", "Hello {{first_name}},\n\nThis is a body message."} and not html_text:
                    self.body_tabs.removeTab(0)

        imported = 0
        self._show_subject_body_loader("Uploading body files.")
        self._workspace_loading = True
        try:
            for file_name in file_paths:
                if self.body_tabs.count() >= MAX_BODY_TABS:
                    self.notify("Maximum of 50 bodies reached")
                    break

                path = Path(file_name)
                suffix = path.suffix.lower()
                if suffix == ".csv":
                    imported += self._import_csv_body_file(path)
                elif suffix in {".html", ".htm"}:
                    imported += self._import_html_body_file(path)
                else:
                    self._log_action(f"Skipped unsupported body file: {path.name}")
        except Exception as exc:
            self._log_action(f"Failed to upload body files: {exc}")
            self.notify(f"Unable to upload bodies: {exc}")
        finally:
            self._workspace_loading = False
            self._hide_subject_body_loader()

        if self.body_tabs.count() == 0:
            self._add_body_draft_tab(select=True)
        else:
            self._refresh_body_tab_labels()
            self._update_body_tab_controls()
            self._sync_active_body_widget_refs()

        if imported > 0:
            self._log_action(f"Uploaded {imported} body(s)")
            self.notify(f"Uploaded {imported} body(s)")
        self._schedule_subject_body_save()

    def _next_body_title(self) -> str:
        return f"Body {self.body_tabs.count() + 1}"

    def _import_csv_body_file(self, path: Path) -> int:
        imported = 0
        with path.open("r", encoding="utf-8-sig", newline="") as handle:
            reader = csv.reader(handle)
            rows = [row for row in reader if any(str(cell).strip() for cell in row)]

        for index, row in enumerate(rows, start=1):
            if self.body_tabs.count() >= MAX_BODY_TABS:
                break
            cells = [str(cell).strip() for cell in row if str(cell).strip()]
            if not cells:
                continue
            body_text = cells[0]
            title = self._next_body_title()
            self._create_and_store_body_tab(
                title=title,
                mode="Normal Message",
                plain_text=body_text,
                html_text="",
                source_name=path.name,
            )
            imported += 1
        return imported

    def _import_html_body_file(self, path: Path) -> int:
        body = self._read_text_template_file(path)
        if not body.strip():
            return 0
        title = self._next_body_title()
        self._create_and_store_body_tab(
            title=title,
            mode="HTML Message",
            plain_text="",
            html_text=body,
            source_name=path.name,
        )
        return 1

    def _create_and_store_body_tab(
        self,
        *,
        title: str,
        mode: str,
        plain_text: str,
        html_text: str,
        source_name: str = "",
    ) -> BodyDraftEditor | None:
        if self.body_tabs.count() >= MAX_BODY_TABS:
            return None
        record = {
            "title": title[:64] or "Body",
            "mode": mode,
            "plain_text": plain_text,
            "html_text": html_text if mode == "HTML Message" else "",
        }
        widget = self._add_body_draft_tab(record, select=False)
        if widget is None:
            return None
        self._refresh_body_tab_labels()
        self._update_body_tab_controls()
        self._persist_subject_body_state()
        return widget

    def _persist_subject_body_state(self) -> None:
        if self._workspace_loading or not self.state.logged_in or not self.state.username:
            return

        subject = self.subject_input.text().strip()
        current_subject = self._selected_subject_item()
        current_body = self._current_body_widget()
        current_attachment = self._current_attachment_widget()
        if current_body is None:
            current_body = self._add_body_draft_tab(select=True)
        body_payload = current_body.payload() if current_body is not None else {"mode": "Normal Message", "plain_text": "", "html_text": "", "title": "Body"}
        body_mode = str(body_payload["mode"])
        plain_body = str(body_payload["plain_text"])
        html_body = str(body_payload["html_text"])
        body_title = str(body_payload["title"] or "Body")
        attachment_html = current_attachment.html_editor.toPlainText() if current_attachment is not None else ""

        self.state.subject_text = subject
        self.state.plain_body_text = plain_body
        self.state.html_message_text = html_body
        self.state.body_mode = body_mode
        self.state.html_template_text = attachment_html

        try:
            subject_rows = []
            for index in range(self.subject_drafts_list.count()):
                item = self.subject_drafts_list.item(index)
                if item is None:
                    continue
                row_subject = self._subject_item_subject(item).strip() or item.text().strip()
                if row_subject:
                    subject_rows.append({"title": row_subject[:64] or "Subject", "subject": row_subject})
            if subject:
                if current_subject is not None:
                    current_index = self.subject_drafts_list.row(current_subject)
                    if current_index >= 0 and current_index < len(subject_rows):
                        subject_rows[current_index] = {"title": subject[:64] or "Subject", "subject": subject}
                    else:
                        subject_rows.append({"title": subject[:64] or "Subject", "subject": subject})
                else:
                    subject_rows.append({"title": subject[:64] or "Subject", "subject": subject})
            _upsert_ui_state(
                LOCAL_SUBJECT_STATE_KEY,
                {
                    "subjects": subject_rows,
                    "selected_index": max(self.subject_drafts_list.currentRow(), 0) if self.subject_drafts_list.count() > 0 else -1,
                },
            )

            body_rows = []
            for index in range(self.body_tabs.count()):
                widget = self.body_tabs.widget(index)
                if not isinstance(widget, BodyDraftEditor):
                    continue
                payload = widget.payload()
                mode = payload["mode"]
                title = payload["title"] or f"Body {index + 1}"
                plain = payload["plain_text"]
                html_text = payload["html_text"]
                body_rows.append(
                    {
                        "title": title[:64] or "Body",
                        "mode": mode,
                        "plain_text": plain,
                        "html_text": html_text if mode == "HTML Message" else "",
                    }
                )
            _upsert_ui_state(
                LOCAL_BODY_STATE_KEY,
                {
                    "tabs": body_rows,
                    "selected_index": self.body_tabs.currentIndex(),
                },
            )

            self._rename_body_tab(current_body, body_title)
        except Exception as exc:
            self._log_action(f"Failed to save subject/body state: {exc}")
        finally:
            self._refresh_campaign_action_state()

    def _save_subject_body_draft(self) -> None:
        if not self.state.logged_in or not self.state.username:
            self.notify("Sign in first to save changes")
            return

        self._show_subject_body_loader("Saving subject and body.")
        try:
            self._persist_subject_body_state()
            title = self.subject_input.text().strip() or "Untitled Subject"
            self._log_action(f"Saved subject and body: {title}")
            self.notify("Subject and body saved")
        except Exception as exc:
            self._log_action(f"Failed to save subject and body: {exc}")
            self.notify("Unable to save subject and body")
        finally:
            self._hide_subject_body_loader()

    def _decode_setting_value(self, raw_value) -> object:
        if raw_value is None:
            return None
        if isinstance(raw_value, str):
            try:
                return json.loads(raw_value)
            except Exception:
                return raw_value
        return raw_value

    def _update_subject_count_label(self) -> None:
        self.subject_count_label.setText(_subject_count_text(self.subject_drafts_list.count()))

    def _load_subjects(self) -> None:
        self._workspace_loading = True
        try:
            payload = _load_ui_state(LOCAL_SUBJECT_STATE_KEY)
            subject_rows = payload.get("subjects") or []
            if not isinstance(subject_rows, list):
                subject_rows = []

            self.subject_drafts_list.blockSignals(True)
            try:
                self.subject_drafts_list.clear()
                for row in reversed(subject_rows):
                    if not isinstance(row, dict):
                        continue
                    item = QListWidgetItem()
                    title = str(row.get("title") or row.get("subject") or "Subject")
                    subject = str(row.get("subject") or row.get("title") or "")
                    self._set_subject_item_data(
                        item,
                        None,
                        title,
                        subject,
                        local_only=True,
                        local_draft_id=None,
                    )
                    self.subject_drafts_list.addItem(item)
                if self.subject_drafts_list.count() > 0:
                    self.subject_drafts_list.setCurrentRow(self.subject_drafts_list.count() - 1)
            finally:
                self.subject_drafts_list.blockSignals(False)

            current_subject = self.subject_drafts_list.currentItem()
            if current_subject is not None:
                subject = self._subject_item_subject(current_subject)
                self.subject_input.blockSignals(True)
                try:
                    self.subject_input.setText(subject)
                finally:
                    self.subject_input.blockSignals(False)
                self.state.subject_text = subject

            self._update_subject_count_label()
            self._update_subject_toggle_visibility()
        finally:
            self._workspace_loading = False
        self._refresh_campaign_action_state()

    def _schedule_pending_emails_save(self) -> None:
        if self._workspace_loading:
            return
        self.state.pending_emails_validated = False
        self._refresh_campaign_action_state()
        self._refresh_pending_email_summary()
        self._pending_emails_save_timer.start()

    def _refresh_pending_email_summary(self) -> None:
        if self._workspace_loading:
            return

        source_text = self.pending_emails_editor.toPlainText()
        candidates = self._extract_email_candidates(source_text)
        total_count = len(candidates)
        self.data_summary_labels["total"].setText(str(total_count))

        if self.state.pending_emails_validated:
            accepted_count = len(self.state.pending_recipients)
            self.data_summary_labels["valid"].setText(str(accepted_count))
            self.data_summary_labels["filter_count"].setText(str(accepted_count))
            self.data_summary_labels["duplicates"].setText("0")
        else:
            self.data_summary_labels["valid"].setText("0")
            self.data_summary_labels["filter_count"].setText("0")
            self.data_summary_labels["duplicates"].setText("0")

    def _persist_pending_emails_state(self) -> None:
        if self._workspace_loading:
            return
        try:
            source_text = self.pending_emails_editor.toPlainText()
            emails = self._extract_email_candidates(source_text)
            gmail_only = self.standard_email_radio.isChecked()
            if source_text.strip() or emails:
                _upsert_ui_state(
                    LOCAL_PENDING_EMAILS_STATE_KEY,
                    {
                        "raw_text": source_text,
                        "emails": emails,
                        "gmail_only": gmail_only,
                        "validated": bool(self.state.pending_emails_validated),
                    },
                )
            else:
                _delete_ui_state(LOCAL_PENDING_EMAILS_STATE_KEY)
        except Exception as exc:
            self._log_action(f"Failed to save pending emails: {exc}")
        finally:
            self._refresh_campaign_action_state()

    def _load_pending_emails_from_local(self) -> None:
        payload = _load_ui_state(LOCAL_PENDING_EMAILS_STATE_KEY)
        raw_text = str(payload.get("raw_text") or "")
        emails_value = payload.get("emails") or []
        emails: list[str] = []
        if isinstance(emails_value, list):
            for value in emails_value:
                email = str(value).strip().lower()
                if email:
                    emails.append(email)
        if not raw_text and emails:
            raw_text = "\n".join(emails)

        self.pending_emails_editor.blockSignals(True)
        try:
            self.pending_emails_editor.setPlainText(raw_text)
        finally:
            self.pending_emails_editor.blockSignals(False)

        self.state.pending_recipients = emails[:]
        self.state.pending_emails_validated = bool(payload.get("validated"))
        self._refresh_pending_email_summary()
        self._refresh_campaign_action_state()

    def _sync_subject_body_widgets(self) -> None:
        self._workspace_loading = True
        try:
            if self.body_tabs.count() == 0:
                self._load_body_tabs_from_state()
            self._sync_active_body_widget_refs()
            self._update_body_tab_controls()
            self._sync_active_attachment_widget_refs()
        finally:
            self._workspace_loading = False

    def _load_body_tabs_from_state(self) -> None:
        payload = _load_ui_state(LOCAL_BODY_STATE_KEY)
        tab_rows = payload.get("tabs") or []
        if not isinstance(tab_rows, list):
            tab_rows = []
        selected_index = int(payload.get("selected_index") or -1)

        self.body_tabs.blockSignals(True)
        try:
            for index in reversed(range(self.body_tabs.count())):
                self.body_tabs.removeTab(index)
            for row in tab_rows:
                if not isinstance(row, dict):
                    continue
                self._add_body_draft_tab(
                    {
                        "title": str(row.get("title") or row.get("label") or "Body"),
                        "mode": str(row.get("mode") or "Normal Message"),
                        "plain_text": str(row.get("plain_text") or ""),
                        "html_text": str(row.get("html_text") or ""),
                    },
                    select=False,
                )
            if self.body_tabs.count() == 0:
                self._add_body_draft_tab(select=True)
            elif 0 <= selected_index < self.body_tabs.count():
                self.body_tabs.setCurrentIndex(selected_index)
            else:
                self.body_tabs.setCurrentIndex(self.body_tabs.count() - 1)
        finally:
            self.body_tabs.blockSignals(False)
        self._refresh_body_tab_labels()
        self._update_body_tab_controls()
        self._sync_active_body_widget_refs()
        self._refresh_campaign_action_state()

    def load_user_workspace(self) -> None:
        if not self.state.logged_in or not self.state.auth_token:
            return

        self._workspace_loading = True
        try:
            settings_payload = api_get_settings(self.state.auth_token)
        except Exception:
            settings_payload = {}

        try:
            settings_rows = settings_payload.get("settings") or []
            settings_map: dict[str, object] = {}
            for row in settings_rows:
                key = str(row.get("setting_key") or "")
                if not key:
                    continue
                settings_map[key] = self._decode_setting_value(row.get("setting_value_json"))

            self.state.body_mode = str(settings_map.get("body_mode") or self.state.body_mode or "Normal Message")
            self.state.subject_text = str(settings_map.get("subject_text") or self.state.subject_text or "")
            self.state.plain_body_text = str(settings_map.get("plain_body_text") or self.state.plain_body_text or "")
            self.state.html_message_text = str(settings_map.get("html_message_text") or self.state.html_message_text or "")
            self.state.html_template_text = str(settings_map.get("html_template_text") or self.state.html_template_text or "")
            self._load_sending_settings_state(settings_map)
            self._load_browser_state()
            self._load_subjects()
            self._load_body_tabs_from_state()
            self._load_attachment_tabs_from_local()
            self._load_pending_emails_from_local()
            self._load_tags_state()
        finally:
            self._workspace_loading = False
            self._sync_subject_body_widgets()

    def _clear_pending_emails(self) -> None:
        self.pending_emails_editor.clear()
        self.state.pending_recipients = []
        self.state.pending_emails_validated = False
        _delete_ui_state(LOCAL_PENDING_EMAILS_STATE_KEY)
        self._refresh_pending_email_summary()
        self._log_action("Cleared pending email list")

    def _show_email_loader(self, subtitle: str) -> None:
        self.window().show_launch_loader("Please wait", subtitle)
        QApplication.processEvents()

    def _hide_email_loader(self) -> None:
        self.window().hide_launch_loader()

    def _extract_email_candidates(self, text: str) -> list[str]:
        pattern = re.compile(r"[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}")
        candidates: list[str] = []
        for match in pattern.finditer(text or ""):
            email = match.group(0).strip().strip("<>[]{}()\"'.,;:")
            if email:
                candidates.append(email.lower())
        return candidates

    def _load_pending_emails_from_file(self) -> None:
        file_path, _ = QFileDialog.getOpenFileName(
            self,
            "Load recipient file",
            "",
            "Data files (*.csv *.xlsx *.xls *.txt);;CSV files (*.csv);;Excel files (*.xlsx *.xls);;Text files (*.txt);;All files (*)",
        )
        if not file_path:
            return

        path = Path(file_path)
        self._show_email_loader(f"Loading recipient data from {path.name}.")
        try:
            emails = self._read_email_file(path)
        except Exception as exc:
            self.notify(f"Unable to load file: {exc}")
            self._log_action(f"Failed to load pending emails from {path.name}: {exc}")
            return
        finally:
            self._hide_email_loader()

        self.pending_emails_editor.blockSignals(True)
        try:
            self.pending_emails_editor.setPlainText("\n".join(emails))
        finally:
            self.pending_emails_editor.blockSignals(False)
        self.state.pending_recipients = emails[:]
        self.state.pending_emails_validated = False
        self._persist_pending_emails_state()
        self._log_action(f"Loaded pending emails from {path.name}")
        self.notify(f"Loaded {len(emails)} email(s)")
        self._prompt_validate_pending_emails()

    def _read_email_file(self, path: Path) -> list[str]:
        suffix = path.suffix.lower()
        rows: list[str] = []
        if suffix in {".xlsx", ".xls"}:
            try:
                import pandas as pd
            except Exception as exc:
                raise RuntimeError("Excel support requires pandas and openpyxl.") from exc
            frame = pd.read_excel(path, header=None, dtype=str)
            for value in frame.fillna("").astype(str).to_numpy().ravel().tolist():
                rows.extend(self._extract_email_candidates(value))
        else:
            with path.open("r", encoding="utf-8-sig", newline="") as handle:
                if suffix == ".csv":
                    reader = csv.reader(handle)
                    for row in reader:
                        for cell in row:
                            rows.extend(self._extract_email_candidates(cell))
                else:
                    rows.extend(self._extract_email_candidates(handle.read()))

        deduped: list[str] = []
        seen: set[str] = set()
        for email in rows:
            if email not in seen:
                seen.add(email)
                deduped.append(email)
        return deduped

    def _prompt_validate_pending_emails(self) -> None:
        reply = QMessageBox.question(
            self,
            "Validate Emails",
            "Do you want to validate the customer email list now?",
            QMessageBox.Yes | QMessageBox.No,
            QMessageBox.No,
        )
        if reply == QMessageBox.Yes:
            self._validate_pending_emails(confirm=False)
        else:
            self._log_action("Skipped email validation by user choice")

    def _validate_pending_emails(self, confirm: bool = True) -> None:
        if confirm:
            reply = QMessageBox.question(
                self,
                "Validate Emails",
                "Do you want to validate the customer email list now?",
                QMessageBox.Yes | QMessageBox.No,
                QMessageBox.No,
            )
            if reply != QMessageBox.Yes:
                self._log_action("Skipped email validation by user choice")
                return

        self._show_email_loader("Validating customer email addresses.")
        try:
            source_text = self.pending_emails_editor.toPlainText()
            candidates = self._extract_email_candidates(source_text)
            gmail_only = self.standard_email_radio.isChecked()

            accepted: list[str] = []
            rejected: list[str] = []
            seen: set[str] = set()
            duplicates = 0

            for email in candidates:
                domain = email.rsplit("@", 1)[-1].lower() if "@" in email else ""
                allowed = (domain == "gmail.com") if gmail_only else bool(domain)
                if not allowed:
                    rejected.append(email)
                    continue
                if email in seen:
                    duplicates += 1
                    continue
                seen.add(email)
                accepted.append(email)

            self.state.pending_recipients = accepted[:]
            self.state.pending_emails_validated = True
            self.pending_emails_editor.blockSignals(True)
            try:
                self.pending_emails_editor.setPlainText("\n".join(accepted))
            finally:
                self.pending_emails_editor.blockSignals(False)
            self.data_summary_labels["total"].setText(str(len(candidates)))
            self.data_summary_labels["valid"].setText(str(len(accepted)))
            self.data_summary_labels["filter_count"].setText(str(len(accepted)))
            self.data_summary_labels["duplicates"].setText(str(duplicates))
            self._persist_pending_emails_state()

            mode_label = "gmail.com only" if gmail_only else "mixed domains"
            self._log_action(
                f"Validated {len(candidates)} email candidate(s) in {mode_label}: "
                f"{len(accepted)} valid, {len(rejected)} filtered out, {duplicates} duplicates"
            )
            self.notify(f"{len(accepted)} valid email(s) ready")
        finally:
            self._hide_email_loader()

    def eventFilter(self, obj, event):
        if obj is self.pending_emails_editor and event.type() == QEvent.KeyPress:
            if event.matches(QKeySequence.Paste):
                result = super().eventFilter(obj, event)
                QTimer.singleShot(0, self._prompt_validate_pending_emails)
                return result
        return super().eventFilter(obj, event)

    def _refresh_sessions(self) -> None:
        self.session_list.clear()
        self._sync_session_state_from_handles()
        for index, item in enumerate(self._browser_sessions, start=1):
            row_widget = self._session_row(item, index)
            list_item = QListWidgetItem()
            list_item.setSizeHint(row_widget.sizeHint())
            self.session_list.addItem(list_item)
            self.session_list.setItemWidget(list_item, row_widget)

    def _refresh_activity(self) -> None:
        if not self.state.activity_log:
            self.activity_log_view.setPlainText("[--:--:--] No activity yet.")
        else:
            lines = "\n".join(self.state.activity_log[-30:][::-1])
            self.activity_log_view.setPlainText(lines)

        if not self._campaign_send_log_entries:
            self.send_log_view.setPlainText("[--:--:--] No send events yet.")
        else:
            self.send_log_view.setPlainText("\n".join(self._campaign_send_log_entries[-80:]))

    def _refresh_controls(self) -> None:
        self.incognito_button.setChecked(self.state.browser_mode == "Incognito")
        self.normal_button.setChecked(self.state.browser_mode == "Normal")
        self.normal_message_button.setChecked(self.state.body_mode == "Normal Message")
        self.html_message_button.setChecked(self.state.body_mode == "HTML Message")
        self.sender_limit.blockSignals(True)
        self.sender_limit.setValue(int(getattr(self.state, "sender_limit", 300)))
        self.sender_limit.blockSignals(False)
        self.delay_from.blockSignals(True)
        self.delay_from.setValue(float(getattr(self.state, "delay_from", 0.5)))
        self.delay_from.blockSignals(False)
        self.delay_to.blockSignals(True)
        self.delay_to.setValue(float(getattr(self.state, "delay_to", 1.0)))
        self.delay_to.blockSignals(False)
        self.retry_count.blockSignals(True)
        self.retry_count.setValue(int(getattr(self.state, "retry_count", 3)))
        self.retry_count.blockSignals(False)
        self.retry_enable_checkbox.blockSignals(True)
        self.retry_enable_checkbox.setChecked(bool(getattr(self.state, "retry_enabled", True)))
        self.retry_enable_checkbox.blockSignals(False)
        self.delay_fixed_radio.blockSignals(True)
        self.delay_random_radio.blockSignals(True)
        self.delay_human_radio.blockSignals(True)
        delay_type = getattr(self.state, "delay_type", "Random range")
        self.delay_fixed_radio.setChecked(delay_type == "Fixed")
        self.delay_random_radio.setChecked(delay_type != "Fixed" and delay_type != "Human-like pattern")
        self.delay_human_radio.setChecked(delay_type == "Human-like pattern")
        self.delay_fixed_radio.blockSignals(False)
        self.delay_random_radio.blockSignals(False)
        self.delay_human_radio.blockSignals(False)
        self.send_seq_radio.blockSignals(True)
        self.send_rand_radio.blockSignals(True)
        self.send_seq_radio.setChecked(getattr(self.state, "email_send_order", "Sequential") != "Random shuffle")
        self.send_rand_radio.setChecked(getattr(self.state, "email_send_order", "Sequential") == "Random shuffle")
        self.send_seq_radio.blockSignals(False)
        self.send_rand_radio.blockSignals(False)
        self.window_parallel_radio.blockSignals(True)
        self.window_sequential_radio.blockSignals(True)
        self.window_parallel_radio.setChecked(getattr(self.state, "window_send_mode", "Parallel") != "Sequential")
        self.window_sequential_radio.setChecked(getattr(self.state, "window_send_mode", "Parallel") == "Sequential")
        self.window_parallel_radio.blockSignals(False)
        self.window_sequential_radio.blockSignals(False)
        current_body = self._current_body_widget()
        if current_body is not None:
            current_body.set_mode(self.state.body_mode)
        self.active_windows_value.setText(str(len(self._browser_sessions)))
        self.launch_preset_label.setText(self.state.launch_preset or "None")
        self.progress_bar.setValue(0)
        self.ai_provider_combo.blockSignals(True)
        self.ai_provider_combo.setCurrentText(self.state.ai_provider)
        self.ai_provider_combo.blockSignals(False)
        self.ai_api_key_input.blockSignals(True)
        self.ai_api_key_input.setText(getattr(self.state, "ai_api_key", ""))
        self.ai_api_key_input.blockSignals(False)
        self._refresh_ai_models()
        self._sync_ai_connection_ui()

    def _log_action(self, message: str) -> None:
        timestamp = QDateTime.currentDateTime().toString("hh:mm:ss")
        self.state.activity_log.append(f"[{timestamp}] {message}")
        if self.state.logged_in and self.state.username:
            try:
                record_activity(
                    self.state.username,
                    message,
                    category="ui",
                    user_id=None,
                )
            except Exception:
                pass
        self._refresh_activity()
        if callable(self.notify):
            self.notify(message)

    def _session_row(self, session: BrowserSessionHandle, index: int) -> QWidget:
        row = QFrame()
        row.setObjectName("sessionRow")
        row_layout = QVBoxLayout(row)
        row_layout.setContentsMargins(_scaled_int(8, self._scale), _scaled_int(7, self._scale), _scaled_int(8, self._scale), _scaled_int(7, self._scale))
        row_layout.setSpacing(_scaled_int(4, self._scale))
        row.setMinimumHeight(_scaled_int(52, self._scale))

        dot = QLabel("●")
        dot.setObjectName("sessionDot")
        label = QLabel(session.title)
        label.setObjectName("sessionTitleSmall")
        state = QLabel(f"{session.mode} - {session.status}")
        state.setObjectName("sessionState")
        count = QLabel(f"({session.send_completed}/{session.send_total})")
        count.setObjectName("sessionState")
        label.setMinimumWidth(_scaled_int(72, self._scale))
        state.setMinimumWidth(_scaled_int(80, self._scale))
        state.setAlignment(Qt.AlignVCenter | Qt.AlignLeft)
        close_button = QPushButton("✕")
        close_button.setObjectName("dangerButton")
        close_button.setFixedWidth(_scaled_int(28, self._scale))
        close_button.clicked.connect(lambda _, sid=session.session_id: self._close_session(sid))
        close_button.setToolTip("Close this browser window")

        top_row = QHBoxLayout()
        top_row.setSpacing(_scaled_int(6, self._scale))
        top_row.addWidget(dot)
        top_row.addWidget(label)
        top_row.addStretch()

        bottom_row = QHBoxLayout()
        bottom_row.setSpacing(_scaled_int(6, self._scale))
        bottom_row.addWidget(state)
        bottom_row.addWidget(count)
        bottom_row.addStretch()
        bottom_row.addWidget(close_button)

        row_layout.addLayout(top_row)
        row_layout.addLayout(bottom_row)
        return row

    def refresh(self) -> None:
        self.window_spin.setValue(self.state.window_count)
        self._refresh_controls()
        self._sync_subject_body_widgets()
        self._refresh_sessions()
        self._refresh_activity()


class MainWindow(QMainWindow):
    def _system_is_dark(self) -> bool:
        palette = QApplication.palette()
        return palette.window().color().lightness() < 150

    def __init__(self):
        super().__init__()
        self.state = AppState()
        self._toasts = []
        self._pending_launch_target = 0
        self._last_login_username = ""
        self._browser_bootstrap_thread: QThread | None = None
        self._browser_bootstrap_worker: BrowserBootstrapWorker | None = None
        self._browser_bootstrap_timer: QTimer | None = None
        self._browser_bootstrap_running = False
        self._session_check_timer = QTimer(self)
        self._session_check_timer.setInterval(30000)
        self._session_check_timer.timeout.connect(self._check_active_login_session)
        self._scale = _compute_layout_scale(QApplication.primaryScreen())
        self._text_scale = _compute_text_scale(QApplication.primaryScreen())
        self._centered_once = False
        self.setWindowTitle(APP_TITLE)
        self.setWindowFlags(Qt.Window | Qt.FramelessWindowHint)
        self.resize(_scaled_int(1120, self._scale), _scaled_int(760, self._scale))
        self._build_ui()

    def _build_ui(self) -> None:
        container = QWidget()
        root = QVBoxLayout(container)
        root.setContentsMargins(0, 0, 0, 0)
        root.setSpacing(0)

        self.title_bar = TitleBar(self, self.close, scale=self._scale)
        self.title_bar.set_logout_handler(self.handle_logout)
        self.title_bar.sync_window_state()
        self.launch_loader = LaunchLoaderDialog(self, scale=self._scale)

        self.stack = QStackedWidget()
        self.login_page = LoginPage(self.handle_login, scale=self._scale)
        self.dashboard_page = DashboardPage(self.state, self.handle_logout, self.show_toast, scale=self._scale)
        self.stack.addWidget(self.login_page)
        self.stack.addWidget(self.dashboard_page)
        root.addWidget(self.title_bar)
        root.addWidget(self.stack, 1)
        self.setCentralWidget(container)
        self._apply_styles()
        self.show_login()

    def show_toast(self, message: str, kind: str = "info") -> None:
        toast = Toast(self, message, kind, scale=self._scale)
        self._toasts.append(toast)
        toast.destroyed.connect(lambda: self._remove_toast(toast))
        toast.adjustSize()
        toast.move(
            max(12, self.width() - toast.width() - 16),
            max(12, self.title_bar.height() + 12),
        )
        toast.show()
        toast.raise_()

    def _remove_toast(self, toast: QWidget) -> None:
        if toast in self._toasts:
            self._toasts.remove(toast)

    def show_launch_loader(self, title: str, subtitle: str) -> None:
        self.launch_loader.set_message(title, subtitle, status_prefix="Preparing")
        self.launch_loader.show()
        self.launch_loader.raise_()
        self.launch_loader.activateWindow()

    def show_bootstrap_loader(self, title: str, subtitle: str) -> None:
        # This is a visual progress overlay, not a modal dialog. Application
        # modality can leave the parent disabled after the overlay closes.
        self.launch_loader.setWindowModality(Qt.NonModal)
        self.launch_loader.setModal(False)
        self.launch_loader.set_busy(title, subtitle, status_prefix="Configuring")
        self.launch_loader.show()
        self.launch_loader.raise_()
        self.launch_loader.activateWindow()

    def hide_launch_loader(self) -> None:
        self.launch_loader.setWindowModality(Qt.NonModal)
        self.launch_loader.setModal(False)
        self.launch_loader.releaseKeyboard()
        self.launch_loader.releaseMouse()
        self.launch_loader.hide()
        self.launch_loader.close()
        self.setEnabled(True)
        self.activateWindow()

    def _start_windows_browser_bootstrap(self) -> None:
        if self._browser_bootstrap_running:
            return

        browser_cache_dir = _browser_cache_dir()
        self._browser_bootstrap_running = True
        self.show_bootstrap_loader(
            "Preparing browser and app dependencies",
            "Checking the local runtime before the workspace opens.",
        )
        self.launch_loader.set_message(
            "Downloading Chromium",
            "Starting the one-time Chromium CDN download • ETA calculating",
            status_prefix="Downloading",
        )
        self.launch_loader.set_progress(2, 100, status_prefix="Downloading")

        # Do not call _log_action here: it records to the remote activity API
        # synchronously and can block the GUI before the downloader starts.
        self.state.activity_log.append("Preparing browser dependencies after login")
        self.dashboard_page._refresh_activity()

        worker = BrowserBootstrapWorker(browser_cache_dir)
        thread = threading.Thread(
            target=worker.run,
            name="ezymailer-browser-bootstrap",
            daemon=True,
        )
        self._browser_bootstrap_thread = thread
        self._browser_bootstrap_worker = worker
        self._browser_bootstrap_timer = QTimer(self)
        self._browser_bootstrap_timer.setInterval(250)
        self._browser_bootstrap_timer.timeout.connect(self._poll_browser_bootstrap)
        self._browser_bootstrap_timer.start()
        # Give the loader one event-loop tick to paint before network I/O.
        QTimer.singleShot(500, thread.start)

    def _poll_browser_bootstrap(self) -> None:
        worker = self._browser_bootstrap_worker
        if worker is None:
            return
        with worker._state_lock:
            progress = worker.latest_progress
            result = worker.result
        if progress is not None:
            self._update_browser_bootstrap_message(*progress)
        if result is not None:
            if self._browser_bootstrap_timer is not None:
                self._browser_bootstrap_timer.stop()
                self._browser_bootstrap_timer.deleteLater()
                self._browser_bootstrap_timer = None
            self._finish_windows_browser_bootstrap(*result)

    def _update_browser_bootstrap_message(self, title: str, subtitle: str, value: int, total: int) -> None:
        if self.launch_loader.isVisible():
            self.launch_loader.set_message(title, subtitle, status_prefix="Configuring")
            self.launch_loader.set_progress(value, total, status_prefix="Configuring")
            if value >= total and title in {"Browser ready", "Dependencies already ready"}:
                # Close synchronously from the GUI timer that received the
                # ready state. A deferred callback can be overwritten by the
                # loader's animation timer and leave the ready screen open.
                self.hide_launch_loader()

    def _finish_windows_browser_bootstrap(self, success: bool, message: str) -> None:
        if not self._browser_bootstrap_running and self._browser_bootstrap_worker is None:
            return
        self._browser_bootstrap_running = False
        self._browser_bootstrap_worker = None
        self._browser_bootstrap_thread = None

        browser_binary = self._browser_binary()
        if success and browser_binary is not None:
            if self.launch_loader.isVisible():
                self.launch_loader.set_message(
                    "Additional files configured",
                    "The app finished preparing its browser runtime.",
                    status_prefix="Done",
                )
                self.launch_loader.set_progress(100, 100, status_prefix="Done")
            self.show_toast("Browser runtime ready", "success")
            if hasattr(self, "dashboard_page"):
                try:
                    if getattr(self.dashboard_page, "_pending_launch_target", 0):
                        QTimer.singleShot(450, self.dashboard_page._resume_pending_launch)
                    QTimer.singleShot(100, self.hide_launch_loader)
                except Exception:
                    QTimer.singleShot(100, self.hide_launch_loader)
            else:
                QTimer.singleShot(100, self.hide_launch_loader)
            return

        self.hide_launch_loader()
        warning_message = message or "Unable to auto-configure the browser runtime."
        print(f"Browser runtime bootstrap failed: {warning_message}")
        self.show_toast("Browser runtime setup needs attention", "warning")
        QMessageBox.warning(
            self,
            "Browser runtime setup",
            "The app could not auto-configure the browser runtime on this device.\n\n"
            "The login screen is still available, but campaign browser automation may require the bundled Chromium runtime to finish downloading.",
        )

    def changeEvent(self, event) -> None:
        super().changeEvent(event)
        if event.type() == QEvent.WindowStateChange:
            self.title_bar.sync_window_state()

    def showEvent(self, event) -> None:
        super().showEvent(event)
        if not self._centered_once:
            self._centered_once = True
            QTimer.singleShot(0, self._center_window_on_screen)

    def _center_window_on_screen(self) -> None:
        screen = self.screen() or QApplication.primaryScreen()
        if screen is None:
            return

        geometry = screen.availableGeometry()
        x = geometry.x() + (geometry.width() - self.width()) // 2
        y = geometry.y() + (geometry.height() - self.height()) // 2
        self.move(x, y)

    def _apply_styles(self) -> None:
        self.setFont(QFont("Segoe UI", _scaled_int(9 if IS_WINDOWS else 10, self._text_scale)))
        style = """
            QMainWindow {
                background: #1e1e1e;
            }
            QWidget {
                color: #d4d4d4;
                font-family: "Segoe UI Variable Text", "Segoe UI", sans-serif;
                font-size: 7.2pt;
            }
            QWidget#tabPage {
                background: #1e1e1e;
            }
            QScrollArea,
            QAbstractScrollArea,
            QScrollArea > QWidget,
            QScrollArea > QWidget > QWidget {
                background: transparent;
                border: none;
            }
            QScrollBar:vertical {
                background: #1e1e1e;
                width: 10px;
                margin: 0px;
            }
            QScrollBar::handle:vertical {
                background: #4b4b4b;
                min-height: 24px;
                border-radius: 4px;
            }
            QScrollBar::handle:vertical:hover {
                background: #5d5d5d;
            }
            QScrollBar::add-line:vertical,
            QScrollBar::sub-line:vertical,
            QScrollBar::add-page:vertical,
            QScrollBar::sub-page:vertical {
                background: transparent;
                border: none;
            }
            QScrollBar:horizontal {
                background: #1e1e1e;
                height: 10px;
                margin: 0px;
            }
            QScrollBar::handle:horizontal {
                background: #4b4b4b;
                min-width: 24px;
                border-radius: 4px;
            }
            QScrollBar::handle:horizontal:hover {
                background: #5d5d5d;
            }
            QScrollBar::add-line:horizontal,
            QScrollBar::sub-line:horizontal,
            QScrollBar::add-page:horizontal,
            QScrollBar::sub-page:horizontal {
                background: transparent;
                border: none;
            }
            QFrame#heroPanel,
            QFrame#loginCard,
            QFrame#topBar,
            QFrame#sidebar,
            QFrame#panelCard,
            QFrame#contentArea {
                background: #1e1e1e;
                border: 1px solid #333333;
                border-radius: 8px;
            }
            QFrame#sidebar {
                border-radius: 0px;
                border-left: none;
                border-top: none;
                border-bottom: none;
            }
            QFrame#contentArea {
                border-radius: 0px;
                border-right: none;
                border-top: none;
                border-bottom: none;
            }
            QFrame#panelCard {
                border-radius: 6px;
                background: #252526;
            }
            QFrame#loginShell {
                background: #252526;
                border: 1px solid #333333;
                border-radius: 10px;
            }
            QFrame#topBar {
                background: #1f1f1f;
                border-radius: 0px;
                border-left: none;
                border-right: none;
                border-top: none;
                border-bottom: 1px solid #333333;
            }
            QFrame#heroPanel {
                background: #252526;
            }
            QFrame#loginCard {
                background: #252526;
            }
            QFrame#toast {
                background: #202020;
                border: 1px solid #3a3d41;
                border-radius: 6px;
            }
            QFrame#loaderCard {
                background: #252526;
                border: 1px solid #333333;
                border-radius: 12px;
            }
            QDialog#launchLoader {
                background: rgba(0, 0, 0, 0.32);
            }
            QFrame#toast[kind="info"] {
                border-left: 3px solid #007acc;
            }
            QFrame#toast[kind="success"] {
                border-left: 3px solid #6a9955;
            }
            QFrame#toast[kind="warning"] {
                border-left: 3px solid #d7ba7d;
            }
            QFrame#toast[kind="error"] {
                border-left: 3px solid #f48771;
            }
            QFrame#toast[kind="info"] QLabel#toastIcon {
                color: #9cdcfe;
            }
            QFrame#toast[kind="success"] QLabel#toastIcon {
                color: #89d185;
            }
            QFrame#toast[kind="warning"] QLabel#toastIcon {
                color: #d7ba7d;
            }
            QFrame#toast[kind="error"] QLabel#toastIcon {
                color: #f48771;
            }
            QFrame#toast[kind="info"] QLabel#toastText {
                color: #d4d4d4;
            }
            QFrame#toast[kind="success"] QLabel#toastText {
                color: #89d185;
            }
            QFrame#toast[kind="warning"] QLabel#toastText {
                color: #d7ba7d;
            }
            QFrame#toast[kind="error"] QLabel#toastText {
                color: #f48771;
            }
            QLabel#heroBadge {
                color: #569cd6;
                font-size: 8pt;
                font-weight: 700;
                letter-spacing: 0.6px;
            }
            QLabel#loginAppName {
                color: #ffffff;
                font-size: 14pt;
                font-weight: 800;
            }
            QLabel#loginKicker {
                color: #569cd6;
                font-size: 8pt;
                font-weight: 700;
            }
            QLabel#heroTitle {
                color: #ffffff;
                font-size: 18pt;
                font-weight: 800;
            }
            QLabel#heroSubtitle,
            QLabel#heroPoints,
            QLabel#heroFooter,
            QLabel#loginSubtitle,
            QLabel#loginHint,
            QLabel#sectionSubtitle,
            QLabel#sectionHint,
            QLabel#placeholderText,
            QLabel#brandSubtitle {
                color: #94a3b8;
            }
            QLabel#loginTitle {
                color: #ffffff;
                font-size: 16pt;
                font-weight: 800;
            }
            QLabel#loaderTitle {
                color: #ffffff;
                font-size: 12pt;
                font-weight: 800;
            }
            QLabel#loaderSubtitle,
            QLabel#loaderStatus {
                color: #9e9e9e;
            }
            QLabel#loginError {
                color: #fda4af;
                min-height: 18px;
            }
            QLabel#brandTitle {
                color: #ffffff;
                font-size: 10pt;
                font-weight: 800;
            }
            QLabel#versionBadge {
                background: #007acc;
                color: white;
                border-radius: 4px;
                padding: 1px 6px;
                font-size: 7pt;
                font-weight: 700;
            }
            QLabel#statusBadge {
                background: #1e4620;
                color: #89d185;
                border: 1px solid #2d5e34;
                border-radius: 4px;
                padding: 3px 7px;
                font-weight: 700;
            }
            QLabel#sectionTitle {
                color: #ffffff;
                font-size: 10pt;
                font-weight: 800;
            }
            QLabel#fieldLabel {
                color: #c8c8c8;
                min-width: 100px;
            }
            QLabel#windowPill {
                background: #3c3c3c;
                color: #ffffff;
                border: 1px solid #4b4b4b;
                border-radius: 4px;
                padding: 3px 8px;
                font-weight: 700;
                max-width: 110px;
            }
            QLabel#toastIcon {
                color: #9cdcfe;
                font-weight: 900;
            }
            QLabel#toastText {
                color: #d4d4d4;
            }
            QLineEdit,
            QTextEdit,
            QSpinBox,
            QComboBox,
            QTableWidget {
                background: #1e1e1e;
                color: #d4d4d4;
                border: 1px solid #3c3c3c;
                border-radius: 4px;
                padding: 5px 8px;
            }
            QLineEdit:focus,
            QTextEdit:focus,
            QSpinBox:focus,
            QComboBox:focus {
                border: 1px solid #007acc;
            }
            QSpinBox {
                min-height: 18px;
            }
            QTableWidget {
                gridline-color: #333333;
                selection-background-color: #094771;
            }
            QHeaderView::section {
                background: #2d2d30;
                color: #d4d4d4;
                border: none;
                padding: 6px 8px;
                font-weight: 700;
            }
            QListWidget {
                background: #1e1e1e;
                border: 1px solid #333333;
                border-radius: 4px;
                padding: 3px;
                color: #d4d4d4;
            }
            QListWidget::item {
                padding: 5px 7px;
                margin: 1px;
                border-radius: 4px;
            }
            QListWidget::item:selected {
                background: #094771;
            }
            QTabWidget::pane {
                border: none;
                background: transparent;
            }
            QTabBar::tab {
                background: #2d2d30;
                color: #d4d4d4;
                padding: 7px 10px;
                margin-right: 6px;
                border-top-left-radius: 4px;
                border-top-right-radius: 4px;
                min-width: 96px;
            }
            QTabBar::tab:selected {
                background: #007acc;
                color: white;
            }
            QPushButton {
                border: none;
                border-radius: 4px;
                padding: 6px 10px;
                font-weight: 700;
            }
            QPushButton#macTrafficLightButton {
                background: transparent;
                border: none;
                padding: 0px;
                margin: 0px;
            }
            QPushButton#macTrafficLightButton:pressed {
                transform: none;
            }
            QPushButton#primaryButton,
            QPushButton#blastButton {
                background: #0e639c;
                color: white;
            }
            QPushButton#secondaryButton {
                background: #3c3c3c;
                color: #d4d4d4;
            }
            QPushButton#warningButton {
                background: #4d3d1f;
                color: #d7ba7d;
            }
            QPushButton#dangerButton {
                background: #5a1d1d;
                color: #f48771;
            }
            QPushButton#closeButton {
                background: #5a1d1d;
                color: #ffffff;
                border: 1px solid #733838;
                border-radius: 4px;
                padding: 0px;
                font-weight: 900;
            }
            QPushButton#closeButton:hover {
                background: #c43c3c;
                color: white;
            }
            QPushButton#windowControlButton {
                background: #3c3c3c;
                color: #d4d4d4;
                border: 1px solid #4b4b4b;
                border-radius: 4px;
                padding: 0px;
                font-weight: 800;
            }
            QPushButton#windowControlButton:hover {
                background: #505050;
                color: white;
            }
            QPushButton:checked {
                background: #007acc;
                color: white;
            }
            QPushButton:pressed {
                transform: translateY(1px);
            }
            QProgressBar {
                background: #1e1e1e;
                border: 1px solid #3c3c3c;
                border-radius: 4px;
                text-align: center;
                color: #d4d4d4;
                height: 20px;
            }
            QProgressBar::chunk {
                background: #007acc;
                border-radius: 4px;
            }
            QLabel#miniStatLabel {
                color: #9e9e9e;
                font-size: 8pt;
                text-transform: uppercase;
            }
            QLabel#countValue {
                color: #ffffff;
                font-size: 13pt;
                font-weight: 800;
            }
            QFrame#miniStat {
                background: #252526;
                border: 1px solid #333333;
                border-radius: 4px;
            }
            QFrame#sessionRow {
                background: #252526;
                border: 1px solid #333333;
                border-radius: 4px;
            }
            QLabel#sessionDot {
                color: #569cd6;
                font-size: 8pt;
            }
            QLabel#sessionTitleSmall {
                color: #ffffff;
                font-weight: 700;
                font-size: 8.5pt;
            }
            QLabel#sessionState {
                color: #9e9e9e;
                font-size: 7.5pt;
            }
            QPushButton#sessionAction {
                background: #0e639c;
                color: white;
                padding: 2px 6px;
                border-radius: 4px;
                min-height: 18px;
                font-size: 7.5pt;
            }
            QTextEdit#activityList,
            QTextEdit#bodyEditor,
            QTextEdit#htmlEditor {
                background: #1e1e1e;
                border: 1px solid #3c3c3c;
                border-radius: 4px;
                color: #d4d4d4;
                font-family: Consolas, "Cascadia Code", "Courier New", monospace;
            }
            QFrame#panelCard QWidget QLabel#sectionHint {
                color: #9e9e9e;
            }
            QFrame#dialogCard {
                background: #252526;
                border: 1px solid #333333;
                border-radius: 4px;
            }
            QFrame#confirmCard {
                background: #252526;
                border: 1px solid #333333;
                border-radius: 4px;
            }
            QDialog#outputDialog {
                background: #1e1e1e;
            }
            QDialog#previewDialog {
                background: rgba(18, 18, 18, 0.96);
            }
            QDialog#confirmDialog {
                background: rgba(0, 0, 0, 0.28);
            }
            QLabel#confirmTitle {
                color: #ffffff;
                font-size: 11pt;
                font-weight: 800;
            }
            QLabel#confirmText {
                color: #d4d4d4;
            }
            QDialog#outputDialog QLabel,
            QDialog#outputDialog QRadioButton {
                color: #d4d4d4;
            }
            QDialog#outputDialog QDialogButtonBox {
                spacing: 8px;
            }
            QFrame#previewCard {
                background: #1e1e1e;
                border: 1px solid #333333;
                border-radius: 8px;
            }
            QDialog#previewDialog QPushButton#secondaryButton {
                background: #2d2d30;
                color: #d4d4d4;
                padding: 5px 10px;
                min-width: 54px;
            }
            QDialog#previewDialog QPushButton#secondaryButton:hover {
                background: #3e3e42;
                color: white;
            }
            QLabel#previewMeta {
                color: #9e9e9e;
            }
            QTextBrowser#previewBrowser {
                background: #111827;
                color: #d4d4d4;
                border: 1px solid #3c3c3c;
                border-radius: 6px;
                padding: 12px;
                font-family: "Segoe UI Variable Text", "Segoe UI", sans-serif;
                font-size: 9pt;
            }
            QTextEdit#sourceEditor {
                background: #0f172a;
                color: #cbd5e1;
                border: 1px solid #334155;
                border-radius: 6px;
                padding: 10px;
                font-family: Consolas, "Cascadia Code", "Courier New", monospace;
                font-size: 8.5pt;
            }
            QTextBrowser#previewBrowser QScrollBar:vertical,
            QTextBrowser#previewBrowser QScrollBar:horizontal {
                background: #111827;
            }
            """
        import re

        def scale_font(match) -> str:
            value = float(match.group(1))
            return f"font-size: {value * self._text_scale:.1f}pt;"

        style = re.sub(r"font-size:\s*([0-9]+(?:\.[0-9]+)?)pt;", scale_font, style)
        style = style.replace("min-width: 96px;", f"min-width: {_scaled_int(96, self._scale)}px;")
        style = style.replace("width: 10px;", f"width: {_scaled_int(10, self._scale)}px;")
        style = style.replace("height: 10px;", f"height: {_scaled_int(10, self._scale)}px;")
        self.setStyleSheet(style)

    def show_login(self) -> None:
        self.hide_launch_loader()
        self.title_bar.set_state("", False)
        self.stack.setCurrentWidget(self.login_page)

    def show_dashboard(self) -> None:
        self.hide_launch_loader()
        self.title_bar.set_state(self.state.username, self.state.logged_in)
        self.dashboard_page.refresh()
        self.stack.setCurrentWidget(self.dashboard_page)

    def handle_login(self, username: str, auth_token: str = "", role: str = "", reset_workspace: bool = False) -> None:
        previous_username = self._last_login_username.strip()
        should_reset_workspace = bool(reset_workspace or (previous_username and previous_username != username))
        self.state.username = username
        self.state.role = role
        self.state.logged_in = True
        self.state.auth_token = auth_token
        self._last_login_username = username
        self.title_bar.set_state(self.state.username, self.state.logged_in)
        self.show_dashboard()
        try:
            if should_reset_workspace:
                self.dashboard_page._reset_campaign_form_state(confirm=False)
            self.dashboard_page.load_user_workspace()
        except Exception as exc:
            self.dashboard_page._log_action(f"Workspace load failed after login: {exc}")
            self.show_toast("Signed in, but workspace reload had an issue", "warning")
        self.dashboard_page._log_action("User authenticated")
        self.show_toast("Signed in successfully", "success")
        self._session_check_timer.start()
        QTimer.singleShot(150, self._start_windows_browser_bootstrap)

    def _check_active_login_session(self) -> None:
        if not self.state.logged_in or not self.state.auth_token:
            return
        try:
            api_get_settings(self.state.auth_token)
        except urllib.error.HTTPError as exc:
            if exc.code in {401, 403}:
                self._perform_logout(confirm=False, reason="This account was signed in on another device.")
        except Exception:
            pass

    def _perform_logout(self, *, confirm: bool, reason: str = "") -> None:
        if confirm and self.state.logged_in and self.state.username:
            reply = QMessageBox.question(
                self,
                "Confirm Logout",
                f"Log out user {self.state.username}?",
                QMessageBox.Yes | QMessageBox.No,
                QMessageBox.No,
            )
            if reply != QMessageBox.Yes:
                return
        self._session_check_timer.stop()
        self._last_login_username = self.state.username
        if self.state.logged_in and self.state.username:
            self.dashboard_page._log_action("User signed out")
        self.hide_launch_loader()
        try:
            self.dashboard_page._handle_campaign_cancel()
            self.dashboard_page._persist_sending_settings_state()
        except Exception:
            pass
        self.dashboard_page._terminate_browser_sessions()
        self.state = AppState()
        self.dashboard_page.state = self.state
        self.dashboard_page.refresh()
        self.login_page.username_input.clear()
        self.login_page.password_input.clear()
        self.login_page.error_label.setText("")
        self.show_login()
        self.show_toast(reason or "Logged out", "warning")

    def handle_logout(self) -> None:
        self._perform_logout(confirm=True)

    def closeEvent(self, event) -> None:
        try:
            if hasattr(self, "dashboard_page"):
                self.dashboard_page._handle_campaign_cancel()
                self.dashboard_page._persist_subject_body_state()
                self.dashboard_page._persist_attachment_state()
                self.dashboard_page._persist_sending_settings_state()
        except Exception:
            pass
        super().closeEvent(event)


def main() -> int:
    try:
        ensure_api_server()
    except Exception as exc:
        print(f"Local login API failed to start: {exc}")

    app = QApplication(sys.argv)
    window = MainWindow()
    window.show()
    return app.exec()


if __name__ == "__main__":
    raise SystemExit(main())
